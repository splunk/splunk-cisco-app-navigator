from pptx import Presentation
prs = Presentation('os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'docs', 'Cisco_Splunk_TA_Strategy_Presentation.pptx')')
print(f"Slides: {len(prs.slides)}")
print(f"Size: {prs.slide_width}, {prs.slide_height}")
for i, slide in enumerate(prs.slides):
    layout = slide.slide_layout.name
    shapes = len(slide.shapes)
    title = ""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            title = ph.text[:60]
            break
    print(f"  [{i+1:2d}] Layout: {layout:35s} | Shapes: {shapes:3d} | Title: {title}")
