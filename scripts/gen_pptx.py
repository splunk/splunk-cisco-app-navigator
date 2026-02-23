#!/usr/bin/env python3
"""Generate Cisco-Splunk TA Strategy Presentation as PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.shapes import MSO_SHAPE

# ─── Colors ───
CISCO_BLUE   = RGBColor(0x04, 0x9F, 0xD9)
SPLUNK_GREEN = RGBColor(0x65, 0xA6, 0x37)
DEBT_RED     = RGBColor(0xD3, 0x2F, 0x2F)
WARN_AMBER   = RGBColor(0xF9, 0xA8, 0x25)
DARK_BG      = RGBColor(0x1A, 0x1A, 0x2E)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY     = RGBColor(0x9E, 0x9E, 0x9E)
BLACK        = RGBColor(0x00, 0x00, 0x00)
DARK_TEXT     = RGBColor(0x21, 0x21, 0x21)
PHASE_GREEN  = RGBColor(0x2E, 0x7D, 0x32)
PHASE_BLUE   = RGBColor(0x15, 0x65, 0xC0)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Helpers ───

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bold_items=None, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = spacing
        if bold_items and i in bold_items:
            p.font.bold = True
    return txBox

def add_table(slide, left, top, width, rows_data, col_widths=None,
              header_color=CISCO_BLUE, font_size=11):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(0.3 * rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Arial"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    return table_shape

def add_stat_callout(slide, left, top, number, label, color=DEBT_RED):
    # Number
    add_text_box(slide, left, top, 2.5, 0.8, number, font_size=42, bold=True,
                 color=color, alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left, top + 0.7, 2.5, 0.5, label, font_size=14,
                 color=MED_GRAY, alignment=PP_ALIGN.CENTER)

def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "Arial"
        p.font.bold = True
    return shape


# ==========================================
# SLIDE 1: Title
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

add_text_box(slide, 1.5, 1.0, 10, 1.2,
             "Cisco-Splunk TA Strategy",
             font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 2.2, 10, 0.8,
             "Findings, Challenges & The Path Forward",
             font_size=24, color=CISCO_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 3.5, 10, 0.5,
             "Engage with Cisco Stakeholders to Re-evaluate Success Factors for TA Authoring & Maintenance",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 5.0, 5, 0.8,
             "Alan Ivarson\nGroup Vice President, Field Solutions",
             font_size=14, color=WHITE, alignment=PP_ALIGN.RIGHT)
add_text_box(slide, 7.0, 5.0, 5, 0.8,
             "Amir (AK) Khamis\nPrincipal Architect",
             font_size=14, color=WHITE, alignment=PP_ALIGN.LEFT)

add_text_box(slide, 1.5, 6.2, 10, 0.5,
             "February 2026  •  45-Minute Presentation",
             font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ==========================================
# SLIDE 2: Agenda
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.4, 11, 0.7, "Agenda", font_size=36, bold=True, color=DARK_BG)

agenda_items = [
    "1.  The Problem — The Current State of Cisco Apps & TAs on Splunkbase",
    "2.  Stakeholder Feedback — Top 10 Challenges & The Magic 6 Best Practices",
    "3.  Why We Can't Hire Our Way Out — Option 1: Traditional TA Development",
    "4.  Why Bundling Alone Isn't Enough — Option 2: The Super App",
    "5.  Two Deliverables — Quick Win (CCC) + Strategic Play (Super App + DCE)",
    "6.  The AI-Powered Path Forward — Lower Cost, Increase Speed & Quality",
    "7.  Next Steps & Roadmap",
]
add_bullet_list(slide, 1.5, 1.5, 10, 5, agenda_items, font_size=20, color=DARK_TEXT,
                bold_items=set(range(7)), spacing=Pt(12))


# ==========================================
# SLIDE 3: The Problem — 100 Apps
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.3, 11, 0.7, "1. The Problem — 100 Apps, No Clear Answer",
             font_size=32, bold=True, color=DEBT_RED)

# CDF Last Mile context
add_text_box(slide, 0.8, 0.95, 7.5, 0.7,
             'Strategic Context: Per Kunal Mukerjee (SVP Engineering), high-fidelity Splunk integrations\n'
             'are the "Last Mile" of the Cisco Data Fabric (CDF). Remediating TAs to Gold Standard enables\n'
             '1,900+ out-of-the-box detections in Splunk Enterprise Security.',
             font_size=11, bold=True, color=CISCO_BLUE)

add_text_box(slide, 8.5, 0.95, 4.3, 0.7,
             '"If we don\'t have good integration and data,\nwe can\'t tell a good story or have full coverage\nspanning all Cisco products in the portfolio."',
             font_size=10, color=MED_GRAY, alignment=PP_ALIGN.LEFT)

# 3 big stat callouts
add_stat_callout(slide, 1.0, 1.8, "41%", "Unsupported", DEBT_RED)
add_stat_callout(slide, 4.5, 1.8, "52%", "Archived / Removed", WARN_AMBER)
add_stat_callout(slide, 8.0, 1.8, "46%", "Community-Built (No QA)", MED_GRAY)

# Pie chart — App Health
chart_data = ChartData()
chart_data.categories = ['Live (42)', 'Archived (39)', 'Manually Archived (13)', 'Flagged (4)', 'Deprecated (13)']
chart_data.add_series('Status', (42, 39, 13, 4, 13))
chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(0.5), Inches(3.3), Inches(4), Inches(3.5), chart_data
)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.font.size = Pt(9)
chart.legend.include_in_layout = False
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.font.size = Pt(9)
data_labels.font.color.rgb = WHITE
# Color the slices
from pptx.oxml.ns import qn
colors_pie = [SPLUNK_GREEN, DEBT_RED, WARN_AMBER, MED_GRAY, RGBColor(0x78, 0x78, 0x78)]
for i, color in enumerate(colors_pie):
    point = plot.series[0].points[i]
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = color

# Pie chart — Support
chart_data2 = ChartData()
chart_data2.categories = ['Not Supported (41)', 'Developer (38)', 'Cisco (13)', 'Splunk (6)']
chart_data2.add_series('Support', (41, 38, 13, 6))
chart_frame2 = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(4.8), Inches(3.3), Inches(4), Inches(3.5), chart_data2
)
chart2 = chart_frame2.chart
chart2.has_legend = True
chart2.legend.position = XL_LEGEND_POSITION.BOTTOM
chart2.legend.font.size = Pt(9)
chart2.legend.include_in_layout = False
plot2 = chart2.plots[0]
plot2.has_data_labels = True
plot2.data_labels.font.size = Pt(9)
plot2.data_labels.font.color.rgb = WHITE
colors_pie2 = [DEBT_RED, WARN_AMBER, CISCO_BLUE, SPLUNK_GREEN]
for i, color in enumerate(colors_pie2):
    point = plot2.series[0].points[i]
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = color

# Pie chart — Developer
chart_data3 = ChartData()
chart_data3.categories = ['Community (46)', 'Cisco (41)', 'Splunk (13)']
chart_data3.add_series('Developer', (46, 41, 13))
chart_frame3 = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(9.1), Inches(3.3), Inches(3.8), Inches(3.5), chart_data3
)
chart3 = chart_frame3.chart
chart3.has_legend = True
chart3.legend.position = XL_LEGEND_POSITION.BOTTOM
chart3.legend.font.size = Pt(9)
chart3.legend.include_in_layout = False
plot3 = chart3.plots[0]
plot3.has_data_labels = True
plot3.data_labels.font.size = Pt(9)
plot3.data_labels.font.color.rgb = WHITE
colors_pie3 = [MED_GRAY, CISCO_BLUE, SPLUNK_GREEN]
for i, color in enumerate(colors_pie3):
    point = plot3.series[0].points[i]
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = color

# Financial impact callout
add_text_box(slide, 0.8, 6.9, 12, 0.4,
             "Financial Impact: Fragmented TAs → thousands of wasted TAC hours, delayed Time-to-Value for multi-million dollar EAs.",
             font_size=11, bold=True, color=DEBT_RED)


# ==========================================
# SLIDE 4: The Confusion Problem — Firewall
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.3, 11, 0.7,
             "2. The Confusion Problem — 13 Apps, 1 Correct Answer",
             font_size=32, bold=True, color=DEBT_RED)

add_text_box(slide, 0.8, 1.0, 11, 0.5,
             'A customer searching Splunkbase for "Cisco Firewall" encounters 13 different results:',
             font_size=14, color=DARK_TEXT)

firewall_data = [
    ["#", "App Name", "Status", "Built By"],
    ["1", "Splunk Add-on for Cisco FireSIGHT", "Archived", "Splunk"],
    ["2", "Cisco Secure eStreamer Client Add-On", "Legacy", "Cisco"],
    ["3", "Cisco Firepower eNcore App", "Legacy", "Community"],
    ["4", "Cisco Secure Firewall App for Splunk", "Superseded", "Cisco"],
    ["5", "Splunk Add-on for Cisco ASA", "Legacy", "Splunk"],
    ["6", "Enosys eStreamer Add-on", "Legacy", "Community"],
    ["7", "Cisco FTD sourcetype TA", "Legacy", "Community"],
    ["8", "Cisco FTD Dashboards", "Legacy", "Community"],
    ["9", "Cisco Firepower pcap Add-on", "Legacy", "Community"],
    ["10", "Cisco eStreamer Client for Splunk", "Legacy", "Cisco"],
    ["11", "Firegen Log Analyzer for Cisco ASA", "Legacy", "Community"],
    ["12", "Cisco Security Suite", "Legacy", "Community"],
    ["13", "Cisco Suite for Splunk", "Legacy", "Community"],
]
add_table(slide, 0.5, 1.5, 8, firewall_data, col_widths=[0.4, 4.5, 1.3, 1.5],
          header_color=DEBT_RED, font_size=10)

# Correct answer callout
add_rounded_rect(slide, 9.0, 1.8, 3.8, 1.5, SPLUNK_GREEN,
                 "CORRECT ANSWER:\n\nInstall Cisco\nSecurity Cloud\n(1 add-on)", font_size=16)

# Top 10 worst offenders (right side)
add_text_box(slide, 9.0, 3.8, 3.8, 0.4, "Top 10 Worst Offenders:",
             font_size=14, bold=True, color=DARK_TEXT)

offenders = [
    "Firewall: 13 legacy apps",
    "Webex: 8",
    "Meraki: 5",
    "Catalyst Center: 5",
    "Nexus: 5",
    "Umbrella: 4",
    "NVM: 3  •  SNA: 3",
    "ThousandEyes: 3  •  ISE: 3",
]
add_bullet_list(slide, 9.0, 4.2, 3.8, 3, offenders, font_size=11, color=DARK_TEXT,
                bold_items={0}, spacing=Pt(4))

add_text_box(slide, 0.5, 6.8, 12, 0.5,
             "93 total legacy entries across 29 of 57 products (51%).  "
             "Wrong app = blank dashboards, broken lookups, duplicate sourcetypes, unsupported TAC tickets.",
             font_size=11, bold=True, color=DEBT_RED)


# ==========================================
# SLIDE 5: Top 10 Challenges + Magic 6
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.3, 11, 0.7,
             "3. Stakeholder Feedback — 20+ SMEs Consulted",
             font_size=32, bold=True, color=DARK_BG)

# Named SME quotes — left column
add_text_box(slide, 0.3, 1.0, 6.3, 0.35,
             '"The current state is a disaster resulting from a lack of holistic strategy."',
             font_size=10, bold=True, color=DEBT_RED)
add_text_box(slide, 0.3, 1.3, 6.3, 0.2,
             '— James Young, Principal Global Security Product Specialist',
             font_size=9, color=MED_GRAY)

add_text_box(slide, 0.3, 1.55, 6.3, 0.35,
             '"One Cisco marketing is outpacing our technical execution — customers are calling us out."',
             font_size=10, bold=True, color=DEBT_RED)
add_text_box(slide, 0.3, 1.85, 6.3, 0.2,
             '— Bryan Pluta, Principal Solutions Engineer',
             font_size=9, color=MED_GRAY)

# SME quotes — right column
add_text_box(slide, 7.0, 1.0, 5.8, 0.35,
             '"CyberCX built their own TAs for 7 years — we risk losing mindshare."',
             font_size=10, bold=True, color=DEBT_RED)
add_text_box(slide, 7.0, 1.3, 5.8, 0.2,
             '— Daniel Peluso, ANZ Splunk Partner SE Manager',
             font_size=9, color=MED_GRAY)

add_text_box(slide, 7.0, 1.55, 5.8, 0.35,
             '"TA sprawl is a hurdle for the Cisco sales force — we need a turnkey experience."',
             font_size=10, bold=True, color=DEBT_RED)
add_text_box(slide, 7.0, 1.85, 5.8, 0.2,
             '— Matt Poland, Sr. Director, SE',
             font_size=9, color=MED_GRAY)

challenges_data = [
    ["#", "Challenge", "Evidence"],
    ["1", "Too many apps", "13 for Firewall; 110 total"],
    ["2", "47% unsupported", "52 apps with zero support"],
    ["3", "51% archived but still referenced", "Old blog posts \u2192 dead apps"],
    ["4", "50% community-built, no QA", "48 developers; CIM conflicts"],
    ["5", "Blank dashboards", "Install \u2192 empty panels \u2192 TAC"],
    ["6", "Version fragmentation", "No in-product update alerts"],
    ["7", "No migration path", "Legacy coexists with modern"],
    ["8", "Duplicate sourcetypes", "Two TAs claim cisco:asa"],
    ["9", "Platform incompatibility", "Cloud vs Enterprise"],
    ["10", "Knowledge gap", "Splunk + Python + Cisco rare"],
    ["11", "Support Loop", "Cisco TAC requires Cisco contract; Splunk admins get bounced"],
    ["12", "Competitive risk", "Palo Alto has AI parsing; partners use Microsoft"],
]
add_table(slide, 0.3, 2.15, 6.5, challenges_data, col_widths=[0.4, 2.8, 3.0],
          header_color=DEBT_RED, font_size=9)

# Brand risk callout
add_rounded_rect(slide, 0.3, 6.0, 6.5, 0.65, RGBColor(0x4A, 0x14, 0x8C),
                 '\u26a0 BRAND RISK: 50% community-built (48 developers) = Cisco doesn\'t "own" its telemetry story.', 10, WHITE)

# End-to-End mandate callout
add_rounded_rect(slide, 0.3, 6.7, 12.5, 0.55, DARK_BG,
                 'END-TO-END MANDATE (James Young & Dimitri McKay): Data \u2192 CIM \u2192 Dashboards \u2192 ESCU Detections \u2192 SOAR Playbooks', 10, WHITE)

# Magic 6 — right side
add_text_box(slide, 7.2, 2.1, 5.5, 0.5, 'The "Magic 6" Best Practices',
             font_size=20, bold=True, color=SPLUNK_GREEN)

magic6_data = [
    ["#", "Best Practice", "Impact"],
    ["1", "Unified add-on family", "1 install, no conflicts"],
    ["2", "Dynamic Config Engine (DCE)", "Zero-footprint, no bloat"],
    ["3", "Sourcetype validation (MTTI)", "Prove data flows → reduce TAC"],
    ["4", "Legacy debt detection", "Auto-flag what to remove"],
    ["5", "Platform awareness", "Right config every time"],
    ["6", "Centralized catalog", "One pane of glass"],
]
add_table(slide, 7.2, 2.7, 5.5, magic6_data, col_widths=[0.4, 2.5, 2.3],
          header_color=SPLUNK_GREEN, font_size=10)

add_text_box(slide, 7.2, 5.5, 5.5, 0.6,
             "If every TA followed these 6 principles,\nthe top 12 challenges would disappear.",
             font_size=14, bold=True, color=SPLUNK_GREEN)


# ==========================================
# SLIDE 6: Option 1 — Can't Hire Our Way Out
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.3, 11, 0.7,
             "4. Option 1: Hire & Build — Doesn't Scale",
             font_size=32, bold=True, color=DARK_BG)

# Scale math
scale_data = [
    ["Scenario", "Calculation", "Result"],
    ["Build 1 TA", "3–4 engineer-months", "Feasible"],
    ["23 TAs (current)", "69–92 eng-months", "6–8 engineer-years"],
    ["57 products", "171–228 eng-months", "14–19 engineer-years"],
    ["400 data sources", "1,200–1,600 eng-months", "100–133 engineer-years"],
]
add_table(slide, 0.5, 1.2, 6, scale_data, col_widths=[2.0, 2.2, 2.2],
          header_color=DARK_BG, font_size=11)

# Bar chart — engineer years
chart_data_bar = CategoryChartData()
chart_data_bar.categories = ['1 TA', '23 TAs', '57 Products', '400 Sources']
chart_data_bar.add_series('Engineer-Years', (0.3, 7, 17, 133))
chart_frame_bar = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(7), Inches(1.0), Inches(5.5), Inches(3.5), chart_data_bar
)
chart_bar = chart_frame_bar.chart
chart_bar.has_legend = False
chart_bar.value_axis.maximum_scale = 140
chart_bar.value_axis.has_title = True
chart_bar.value_axis.axis_title.text_frame.paragraphs[0].text = "Engineer-Years"
chart_bar.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
series = chart_bar.series[0]
# Color the bars
bar_colors = [SPLUNK_GREEN, WARN_AMBER, DEBT_RED, RGBColor(0x8B, 0x00, 0x00)]
for i, color in enumerate(bar_colors):
    point = series.points[i]
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = color

# Knowledge Venn — 4 domain boxes
add_text_box(slide, 0.5, 4.2, 6, 0.5, "The Knowledge Problem — 4 Domains, Very Few Unicorns",
             font_size=16, bold=True, color=DARK_TEXT)

add_rounded_rect(slide, 0.5, 4.8, 2.8, 0.7, CISCO_BLUE, "Splunk UCC Framework", 11, WHITE)
add_rounded_rect(slide, 3.5, 4.8, 2.8, 0.7, PHASE_GREEN, "Python & Scripting", 11, WHITE)
add_rounded_rect(slide, 0.5, 5.6, 2.8, 0.7, RGBColor(0x6A, 0x1B, 0x9A), "Splunk Platform\nprops / CIM / SPL", 10, WHITE)
add_rounded_rect(slide, 3.5, 5.6, 2.8, 0.7, DEBT_RED, "Cisco Domain\nAPIs / Syslog / MIBs", 10, WHITE)
add_rounded_rect(slide, 7.0, 5.1, 3.0, 0.9, WARN_AMBER, "⭐ THE UNICORN\nKnows All Four", 14, BLACK)

# Arrow shapes pointing to unicorn
for (x, y) in [(3.2, 5.15), (6.0, 5.15), (3.2, 5.95), (6.0, 5.95)]:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(0.6), Inches(0.25)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MED_GRAY
    shape.line.fill.background()

add_text_box(slide, 0.5, 6.6, 12, 0.5,
             "Conclusion: We cannot hire our way out of this problem.",
             font_size=18, bold=True, color=DEBT_RED, alignment=PP_ALIGN.CENTER)


# ==========================================
# SLIDE 7: Option 2 — Super App Bundling
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.3, 11, 0.7,
             "5. Option 2: Super App Bundling — Right Direction, Still Expensive",
             font_size=28, bold=True, color=DARK_BG)

# Consolidation ratios
consol_data = [
    ["Super Add-on", "Products", "Ratio"],
    ["Cisco Security Cloud", "19 security", "19 : 1"],
    ["Cisco Catalyst Add-on", "13 networking", "13 : 1"],
    ["Cisco DC Networking", "3 data center", "3 : 1"],
    ["Total", "35 → 3 add-ons", "~12 : 1"],
]
add_table(slide, 0.5, 1.2, 5.5, consol_data, col_widths=[2.5, 1.8, 1.0],
          header_color=CISCO_BLUE, font_size=11)

# Before / After
ba_data = [
    ["Metric", "Before", "After"],
    ["Customer installs", "23 separate TAs", "3 Super Add-ons"],
    ["Upgrade cycles", "23 paths", "3"],
    ["Splunkbase confusion", "110 apps", "3 clear choices"],
    ["Version conflicts", "Constant", "Eliminated"],
    ["Support", "Fragmented", "Unified"],
]
add_table(slide, 6.5, 1.2, 6.3, ba_data, col_widths=[2.0, 2.0, 2.0],
          header_color=SPLUNK_GREEN, font_size=11)

# Remaining challenges
add_text_box(slide, 0.5, 4.0, 12, 0.5, "But It Still Has Challenges:",
             font_size=18, bold=True, color=WARN_AMBER)

remaining = [
    "Maintenance cost — many mini-TAs inside; still need per-product parsing updates",
    "Regression risk — updating Firewall parsing could break Duo parsing",
    "Release velocity — can't ship a Meraki fix without regression-testing all 13 Catalyst products",
    "Blank dashboard problem — app installed ≠ data flowing; static configs validate nothing",
    "Knowledge bottleneck — still need engineers who understand each product's data format",
]
add_bullet_list(slide, 1.0, 4.5, 11, 2.5, remaining, font_size=13, color=DARK_TEXT, spacing=Pt(6))

add_text_box(slide, 0.5, 6.8, 12, 0.5,
             "Super App is the right architecture, but doesn't solve cost, speed, and quality by itself.",
             font_size=16, bold=True, color=WARN_AMBER, alignment=PP_ALIGN.CENTER)


# ==========================================
# SLIDE 8: Two Deliverables
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.3, 11, 0.7,
             "6. Two Deliverables — Quick Win + Strategic Play",
             font_size=32, bold=True, color=CISCO_BLUE)

# CCC box (left)
add_rounded_rect(slide, 0.5, 1.3, 5.8, 1.0, WARN_AMBER,
                 'DELIVERABLE 1: Cisco Control Center (CCC)\n"The Front Door"  •  Ships in Weeks', 14, BLACK)

ccc_items = [
    "✅  Catalog 57 products with correct add-on for each",
    "✅  Monitor 321 sourcetypes — Proof of Innocence (MTTI)",
    "✅  Auto-detect 93 legacy apps and warn what to remove",
    "✅  Map SOAR (10), ITSI (4), alert actions (3)",
    "✅  One-click dashboard launch for 26 products",
    "✅  Platform-aware guidance (Cloud vs Enterprise)",
    "",
    '\"The network is innocent until proven guilty —',
    'CCC provides the proof in seconds.\"',
    "Lightweight: no conf generation, no data collection",
    "Splunkbase in 2–3 weeks (AppInspect + packaging)",
]
add_bullet_list(slide, 0.7, 2.5, 5.3, 3.5, ccc_items, font_size=11, color=DARK_TEXT, spacing=Pt(3))

# Super App box (right)
add_rounded_rect(slide, 6.8, 1.3, 6.0, 1.0, PHASE_GREEN,
                 'DELIVERABLE 2: Enterprise Networking Super App\n"The Engineering Solution"  •  3–6 Months', 14, WHITE)

sa_items = [
    "Dynamic Configuration Engine (DCE) — just-in-time activation",
    "Master Library — gold-standard mini-TAs per product",
    "Logical Isolation — Meraki error can't crash ASA or ISE",
    "Tier-Aware — Full Stack / Heavy Forwarder / Search Head Only",
    "Zero-footprint — nothing active until admin enables",
    "6 products onboarded (ASA, DC Net, Intersight, Meraki, UCS, Webex)",
    "🔒 Security Posture: removes 93 legacy apps → reduced attack surface",
    "",
    "CDF Alignment — mini-TAs = Data Products for Data Fabric",
    "MCP Integration — AI-ready telemetry for autonomous agents",
    "NetFlow/IPFIX — first-class foundational signal",
]
add_bullet_list(slide, 7.0, 2.5, 5.5, 3.5, sa_items, font_size=11, color=DARK_TEXT, spacing=Pt(3))

# Comparison table
compare_data = [
    ["", "CCC (Quick Win)", "Super App (Strategic)"],
    ["Purpose", "Guide to the right TA", "Replace legacy TAs entirely"],
    ["Ships", "2–3 weeks", "3–6 months (v1)"],
    ["Customer value", "What do I install?", "One app, activate what I need"],
    ["Effort", "Low — packaging", "Medium — expand & harden"],
    ["Risk", "Very low", "Moderate"],
    ["Solves confusion", "✅ Immediately", "✅ Permanently"],
    ["Solves blank dashboards", "✅ Detects missing data", "✅ Active products only"],
    ["Solves legacy debt", "✅ Warns about legacy", "✅ Replaces entirely"],
]
add_table(slide, 1.5, 6.0, 10, compare_data, col_widths=[2.5, 3.5, 3.5],
          header_color=CISCO_BLUE, font_size=9)

# Aspirin / Cure tagline — draw across the bottom
add_text_box(slide, 0.5, 5.4, 12, 0.5,
             'CCC is the aspirin.  The Super App is the cure.',
             font_size=20, bold=True, color=CISCO_BLUE, alignment=PP_ALIGN.CENTER)


# ==========================================
# SLIDE 9: AI-Powered Path Forward
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.3, 11, 0.7,
             "7. The AI-Powered Path Forward",
             font_size=32, bold=True, color=DARK_BG)

add_text_box(slide, 0.8, 1.0, 11, 0.5,
             'Leveraging FDSE (Rafal Piekarz) for AI-driven maintenance + HCL for validation.\n'
             '"Palo Alto has already implemented AI-driven parsing, proving viability." — Dimitri McKay',
             font_size=12, color=MED_GRAY)

# Key enablers callout boxes
add_rounded_rect(slide, 6.5, 5.2, 3.0, 0.7, CISCO_BLUE,
                 'TACO Framework\nCertification in <24hrs\n(vs. multi-week manual)', 9, WHITE)
add_rounded_rect(slide, 9.7, 5.2, 3.2, 0.7, SPLUNK_GREEN,
                 'Activation Telemetry\nper Kunal: every DCE toggle\nemits ROI proof data', 9, WHITE)

# AI comparison table
ai_data = [
    ["Step", "Traditional", "AI-Assisted"],
    ["Data source analysis", "1–2 weeks", "Hours"],
    ["Input configuration", "1–2 weeks", "Hours"],
    ["Parsing (props/transforms)", "2–4 weeks", "Hours"],
    ["CIM mapping", "1–2 weeks", "Hours"],
    ["Dashboards", "2–4 weeks", "Hours"],
    ["Human validation (HCL)", "—", "1–2 weeks"],
    ["Certification", "1–2 weeks", "1 week"],
    ["TOTAL", "3–4 months", "2–3 weeks"],
]
add_table(slide, 0.5, 1.6, 5.5, ai_data, col_widths=[2.5, 1.5, 1.5],
          header_color=DARK_BG, font_size=11)

# Bar chart — engineer years by approach
chart_data_approaches = CategoryChartData()
chart_data_approaches.categories = ['Option 1:\nHire (57)', 'Option 1:\nHire (400)', 'Option 2:\nStatic Super', 'Our Approach:\nAI + DCE']
chart_data_approaches.add_series('Engineer-Years', (17, 133, 7, 1))
chart_frame_ap = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(6.5), Inches(1.5), Inches(6.3), Inches(3.5), chart_data_approaches
)
chart_ap = chart_frame_ap.chart
chart_ap.has_legend = False
chart_ap.value_axis.maximum_scale = 140
chart_ap.value_axis.has_title = True
chart_ap.value_axis.axis_title.text_frame.paragraphs[0].text = "Engineer-Years"
chart_ap.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
series_ap = chart_ap.series[0]
ap_colors = [WARN_AMBER, DEBT_RED, CISCO_BLUE, SPLUNK_GREEN]
for i, color in enumerate(ap_colors):
    point = series_ap.points[i]
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = color

# Options comparison table
options_data = [
    ["", "Option 1: Hire", "Option 2: Static", "Our Approach"],
    ["Architecture", "1 TA per product", "Bundled, static", "Dynamic Super App + DCE"],
    ["Intelligence", "None", "None", "CCC + MTTI"],
    ["AI authoring", "No", "No", "AI drafts, humans validate"],
    ["Cost", "$$$$$ (100+ eng-yrs)", "$$$ (6–8 eng-yrs)", "$ (<1 eng-yr)"],
    ["Speed", "3–4 months/TA", "2–3 months", "2–3 weeks"],
    ["Scale to 400", "Impossible", "Very expensive", "Feasible"],
]
add_table(slide, 0.5, 5.3, 12, options_data, col_widths=[1.8, 3.0, 3.0, 3.5],
          header_color=CISCO_BLUE, font_size=10)


# ==========================================
# SLIDE 10: Next Steps & Roadmap
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.3, 11, 0.7,
             "8. Next Steps & Roadmap",
             font_size=32, bold=True, color=DARK_BG)

# Immediate actions
actions_data = [
    ["#", "Action", "Owner", "Timeline"],
    ["1", "Ship CCC to Splunkbase — package, AppInspect, certify", "AK", "2–3 weeks"],
    ["2", "Expand Super App Master Library (next 5 products)", "AK + HCL", "0–3 months"],
    ["3", "Pilot AI-assisted TA authoring on 3 products", "AK + HCL", "0–3 months"],
    ["4", "⚠️ DECISION: Approve FDSE partnership (Rafal Piekarz) — dedicated AI-driven TA maintenance cycles", "Alan / Exec", "This quarter"],
    ["5", "⚠️ DECISION: Allocate HCL validation resources (2–3 engineers) for Super App QA", "Alan / Exec", "This quarter"],
    ["6", "Coordinate exec summary with Kunal Mukerjee — align CDF/MDL positioning", "Alan / AK", "Next 2 weeks"],
]
add_table(slide, 0.5, 1.1, 12, actions_data, col_widths=[0.4, 6.5, 2.0, 2.0],
          header_color=CISCO_BLUE, font_size=11)

# Phase boxes (roadmap)
add_text_box(slide, 0.5, 3.8, 12, 0.5, "Value Realization — Three-Phase Journey",
             font_size=18, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Phase 1
add_rounded_rect(slide, 0.3, 4.4, 3.8, 2.5, WARN_AMBER, "", 1, BLACK)
add_text_box(slide, 0.5, 4.5, 3.4, 0.4, 'Phase 1: "Stop the Bleeding"',
             font_size=14, bold=True, color=BLACK)
add_text_box(slide, 0.5, 4.9, 3.4, 0.3, '0–3 months  •  Immediate Visibility',
             font_size=10, color=RGBColor(0x55, 0x55, 0x55))
phase1_items = [
    "Ship CCC to Splunkbase",
    "Expand Master Library to 12+",
    "Pilot AI authoring w/ FDSE",
    "Migrate top 10 legacy offenders",
]
add_bullet_list(slide, 0.5, 5.3, 3.4, 1.5, phase1_items, font_size=10, color=BLACK, spacing=Pt(3))

# Phase 2
add_rounded_rect(slide, 4.5, 4.4, 3.8, 2.5, PHASE_GREEN, "", 1, WHITE)
add_text_box(slide, 4.7, 4.5, 3.4, 0.4, 'Phase 2: "Scale the Portfolio"',
             font_size=14, bold=True, color=WHITE)
add_text_box(slide, 4.7, 4.9, 3.4, 0.3, '3–6 months  •  Operational Efficiency',
             font_size=10, color=RGBColor(0xC8, 0xE6, 0xC9))
phase2_items = [
    "Ship Super App v1 to Splunkbase",
    "AI pipeline: 5–10 mini-TAs/month",
    "NetFlow/IPFIX hardened",
    "HCL validation cadence",
]
add_bullet_list(slide, 4.7, 5.3, 3.4, 1.5, phase2_items, font_size=10, color=WHITE, spacing=Pt(3))

# Phase 3
add_rounded_rect(slide, 8.7, 4.4, 4.2, 2.5, PHASE_BLUE, "", 1, WHITE)
add_text_box(slide, 8.9, 4.5, 3.8, 0.4, 'Phase 3: "Own the Data Fabric"',
             font_size=14, bold=True, color=WHITE)
add_text_box(slide, 8.9, 4.9, 3.8, 0.3, '6–12 months  •  Strategic Differentiation',
             font_size=10, color=RGBColor(0xBB, 0xDE, 0xFB))
phase3_items = [
    "Full CDF/MDL integration",
    "MCP-driven AI features",
    "200+ data sources covered",
    'CCC + Super App = "One Cisco"',
]
add_bullet_list(slide, 8.9, 5.3, 3.8, 1.5, phase3_items, font_size=10, color=WHITE, spacing=Pt(3))

# Arrows between phases
for x in [4.1, 8.3]:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(5.5), Inches(0.4), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MED_GRAY
    shape.line.fill.background()


# ==========================================
# SLIDE 11: Summary
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, 0.8, 0.3, 11, 0.7, "Summary",
             font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

summary_data = [
    ["The Problem", "Our Answer"],
    ["110 apps, customers confused", "CCC: ships in weeks, shows the right answer"],
    ["Blank dashboards, no validation", "CCC monitors 321 sourcetypes; Super App activates what has data"],
    ["Legacy apps cause conflicts", "CCC detects 93 legacy; Super App replaces entirely"],
    ["TAs cost 3–4 months each", "AI authoring: 2–3 weeks per product"],
    ["Can't hire (400 = 133 eng-yrs)", "AI drafts, humans validate, Super App scales"],
    ["Static Super Apps have regression risk", "DCE: each product is an isolated mini-TA"],
]
add_table(slide, 1.0, 1.3, 11, summary_data, col_widths=[4.5, 6.5],
          header_color=CISCO_BLUE, font_size=12)

# The Funnel visual
add_text_box(slide, 0.5, 4.8, 12, 0.5, "The Consolidation Journey",
             font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, 0.8, 5.4, 3.5, 1.2, DEBT_RED,
                 "TODAY\n110 apps  •  60% dead\nNo guidance", 13, WHITE)

shape = slide.shapes.add_shape(
    MSO_SHAPE.RIGHT_ARROW, Inches(4.4), Inches(5.7), Inches(0.6), Inches(0.4)
)
shape.fill.solid()
shape.fill.fore_color.rgb = MED_GRAY
shape.line.fill.background()

add_rounded_rect(slide, 5.1, 5.4, 3.5, 1.2, WARN_AMBER,
                 'QUICK WIN: CCC\n"The Front Door"\nShips in weeks', 13, BLACK)

shape2 = slide.shapes.add_shape(
    MSO_SHAPE.RIGHT_ARROW, Inches(8.7), Inches(5.7), Inches(0.6), Inches(0.4)
)
shape2.fill.solid()
shape2.fill.fore_color.rgb = MED_GRAY
shape2.line.fill.background()

add_rounded_rect(slide, 9.4, 5.4, 3.5, 1.2, PHASE_GREEN,
                 'STRATEGIC: Super App\n"Activate What You Need"\n3–6 months', 13, WHITE)

add_text_box(slide, 0.5, 6.8, 12, 0.5,
             "Together, they move Cisco from fragmented products to a unified Security & Networking platform.",
             font_size=14, bold=True, color=CISCO_BLUE, alignment=PP_ALIGN.CENTER)


# ==========================================
# SLIDE 12: Stakeholder Endorsements
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.3, 11, 0.7,
             "Stakeholder Endorsements — 20+ SMEs Consulted",
             font_size=32, bold=True, color=DARK_BG)

endorse_data = [
    ["Stakeholder", "Role", "Key Insight"],
    ["James Young", "Principal Global Security Specialist", '"A disaster... practitioner-led dev is faster than siloed 5-yr projects"'],
    ["Bryan Pluta", "Principal Solutions Engineer", '"One Cisco marketing is outpacing execution — customers calling us out"'],
    ["Daniel Peluso", "ANZ Splunk Partner SE Manager", '"CyberCX built own TAs for 7 yrs — we risk losing mindshare"'],
    ["Matt Poland", "Sr. Director, SE", '"TA sprawl hurts the sales force — need turnkey, outcome-focused experience"'],
    ["Colin Gibbens", "Principal Engineer (CSC App)", 'Built 16 integrations in 1 yr via SoftServe; supports expanding'],
    ["Dimitri McKay", "Principal Security Architect", 'Validated DCE for resource bloat; Palo Alto already has AI parsing'],
    ["Sarav Radhakrishnan", "Distinguished Engineer", 'Validated DCE POC; supports Two Super App strategy (Security vs Net)'],
    ["Yaron Caspy", "Engineering PM", 'Built SSE App because one-size-fits-all failed (millions of missed events)'],
    ["Steven Moore", "Cisco Solutions Architect", '"Network is a black box — infra is first point of blame" (coined MTTI)'],
    ["Kunal Mukerjee", "SVP Engineering", '"TA Flywheel" vision; Master Library + DCE = CDF Data Products'],
]
add_table(slide, 0.3, 1.1, 12.5, endorse_data, col_widths=[2.0, 3.5, 7.0],
          header_color=CISCO_BLUE, font_size=10)

add_text_box(slide, 0.5, 5.8, 12, 0.5,
             "Feedback unanimously supports the CCC + Dynamic Super App strategy.",
             font_size=16, bold=True, color=SPLUNK_GREEN, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 0.5, 6.3, 12, 0.5,
             'Strategy: CCC is the aspirin. The Super App is the cure.',
             font_size=22, bold=True, color=CISCO_BLUE, alignment=PP_ALIGN.CENTER)


# ==========================================
# SLIDE 13: Thank You
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, 1.5, 2.0, 10, 1.0, "Thank You",
             font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 3.2, 10, 0.5, "Questions & Discussion",
             font_size=24, color=CISCO_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 4.5, 10, 1.0,
             "Alan Ivarson  •  Group Vice President, Field Solutions\n"
             "Amir (AK) Khamis  •  Principal Architect",
             font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ─── Save ───
output_path = "/Users/akhamis/repo/cisco_control_center_app/docs/Cisco_Splunk_TA_Strategy_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
