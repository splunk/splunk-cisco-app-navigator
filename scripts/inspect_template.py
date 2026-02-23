#!/usr/bin/env python3
"""Inspect the Cisco PowerPoint template to discover available slide layouts and their placeholders."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

TEMPLATE = '/Users/akhamis/repo/cisco_control_center_app/docs/Cisco_PowerPoint_Template_LIGHT.potx'
prs = Presentation(TEMPLATE)

print(f"Slide width: {prs.slide_width} EMU = {prs.slide_width / 914400:.2f} inches")
print(f"Slide height: {prs.slide_height} EMU = {prs.slide_height / 914400:.2f} inches")
print(f"Number of slide layouts: {len(prs.slide_layouts)}")
print()

for i, layout in enumerate(prs.slide_layouts):
    print(f"{'='*60}")
    print(f"Layout [{i}]: '{layout.name}'")
    print(f"  Placeholders: {len(layout.placeholders)}")
    for ph in layout.placeholders:
        print(f"    idx={ph.placeholder_format.idx}, type={ph.placeholder_format.type}, "
              f"name='{ph.name}', "
              f"pos=({ph.left/914400:.2f}\", {ph.top/914400:.2f}\"), "
              f"size=({ph.width/914400:.2f}\" x {ph.height/914400:.2f}\")")
    # Also check for shapes that aren't placeholders
    shape_count = len(layout.shapes) - len(layout.placeholders)
    if shape_count > 0:
        print(f"  Additional shapes (non-placeholder): {shape_count}")
        for shape in layout.shapes:
            if not shape.is_placeholder:
                print(f"    shape: '{shape.name}', type={shape.shape_type}, "
                      f"pos=({shape.left/914400:.2f}\", {shape.top/914400:.2f}\"), "
                      f"size=({shape.width/914400:.2f}\" x {shape.height/914400:.2f}\")")
    print()

# Also check slide masters
print(f"\n{'='*60}")
print(f"Slide Masters: {len(prs.slide_masters)}")
for j, master in enumerate(prs.slide_masters):
    print(f"  Master [{j}]: '{master.name if hasattr(master, 'name') else 'N/A'}'")
    print(f"    Placeholders: {len(master.placeholders)}")
    for ph in master.placeholders:
        print(f"      idx={ph.placeholder_format.idx}, name='{ph.name}'")
