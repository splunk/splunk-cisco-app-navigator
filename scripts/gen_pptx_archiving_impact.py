#!/usr/bin/env python3
"""
Generate single-slide PPTX: Splunkbase Mass Archiving — Impact on Cisco Ecosystem.
Uses the Cisco PowerPoint Template if available, otherwise a blank presentation.
"""

import shutil
import zipfile
import os
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Paths ───
_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
_REPO_ROOT  = os.path.join(_SCRIPT_DIR, '..')
TEMPLATE_POTX = os.path.join(_REPO_ROOT, 'docs', 'Cisco_PowerPoint_Template_LIGHT.potx')
TEMPLATE_PPTX = os.path.join(_REPO_ROOT, 'docs', 'cisco_template_converted.pptx')
OUTPUT_PATH   = os.path.join(_REPO_ROOT, 'docs', 'Splunkbase_Archiving_Impact.pptx')

# ─── Convert .potx → .pptx ───
def convert_potx_to_pptx(potx_path, pptx_path):
    shutil.copy2(potx_path, pptx_path)
    tmpdir = tempfile.mkdtemp()
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        zin.extractall(tmpdir)
    ct_path = os.path.join(tmpdir, '[Content_Types].xml')
    with open(ct_path, 'r', encoding='utf-8') as f:
        content = f.read()
    content = content.replace(
        'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
    )
    with open(ct_path, 'w', encoding='utf-8') as f:
        f.write(content)
    with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for root, dirs, files in os.walk(tmpdir):
            for fname in files:
                full_path = os.path.join(root, fname)
                arcname = os.path.relpath(full_path, tmpdir)
                zout.write(full_path, arcname)
    shutil.rmtree(tmpdir)

# ─── Load template or create blank widescreen presentation ───
use_template = False
if os.path.exists(TEMPLATE_POTX):
    convert_potx_to_pptx(TEMPLATE_POTX, TEMPLATE_PPTX)
    prs = Presentation(TEMPLATE_PPTX)
    use_template = True
    # Remove all pre-existing template slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId is None:
            for k, v in prs.slides._sldIdLst[0].attrib.items():
                if k.endswith('}id') or k == 'r:id':
                    rId = v
                    break
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    print("Using Cisco PowerPoint Template (LIGHT)")
else:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    print("Cisco template not found — using blank widescreen presentation")

# ─── Colors ───
CISCO_BLUE   = RGBColor(0x04, 0x9F, 0xD9)
SPLUNK_GREEN = RGBColor(0x65, 0xA6, 0x37)
DEBT_RED     = RGBColor(0xD3, 0x2F, 0x2F)
WARN_AMBER   = RGBColor(0xF9, 0xA8, 0x25)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY     = RGBColor(0x9E, 0x9E, 0x9E)
DARK_TEXT     = RGBColor(0x21, 0x21, 0x21)
PURPLE       = RGBColor(0x4A, 0x14, 0x8C)
PHASE_GREEN  = RGBColor(0x2E, 0x7D, 0x32)
PHASE_BLUE   = RGBColor(0x15, 0x65, 0xC0)
ARCHIVE_ORANGE = RGBColor(0xE6, 0x51, 0x00)

# ─── Helpers ───

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    if font_name:
        p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text_box(slide, left, top, width, height, runs_list, alignment=PP_ALIGN.LEFT):
    """Add a text box with mixed formatting. runs_list = list of paragraphs,
    each paragraph is a list of (text, font_size, bold, color) tuples."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for pi, para_runs in enumerate(runs_list):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = alignment
        for ri, (text, font_size, bold, color) in enumerate(para_runs):
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
    return shape


def add_stat_callout(slide, left, top, number, label, color=DEBT_RED):
    add_text_box(slide, left, top, 2.2, 0.65, number, font_size=36, bold=True,
                 color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + 0.55, 2.2, 0.4, label, font_size=11,
                 color=MED_GRAY, alignment=PP_ALIGN.CENTER)


def add_section_box(slide, left, top, width, height, header, header_color,
                    body_lines, body_size=10):
    """Rounded rect with header + body text."""
    # Background rect
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(0.75)

    # Header bar
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.32)
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = header_color
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = header
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True

    # Body text
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.12), Inches(top + 0.38),
        Inches(width - 0.24), Inches(height - 0.44)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        # Parse inline **bold** markers anywhere in the line
        import re
        parts = re.split(r'(\*\*.*?\*\*)', line)
        if len(parts) > 1:
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    run = p.add_run()
                    run.text = part[2:-2]
                    run.font.size = Pt(body_size)
                    run.font.bold = True
                    run.font.color.rgb = DARK_TEXT
                elif part:
                    run = p.add_run()
                    run.text = part
                    run.font.size = Pt(body_size)
                    run.font.color.rgb = DARK_TEXT
        else:
            p.text = line
            p.font.size = Pt(body_size)
            p.font.color.rgb = DARK_TEXT

    return shape


def set_title(slide, text, color=DARK_TEXT):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = text
            for paragraph in ph.text_frame.paragraphs:
                paragraph.font.color.rgb = color
            return ph
    return None


# ==========================================
# SLIDE 1: Splunkbase Mass Archiving Impact
# ==========================================
if use_template:
    slide = prs.slides.add_slide(prs.slide_layouts[9])  # LAYOUT_TITLE_ONLY
    set_title(slide, "The Splunkbase Mass Archiving — Impact on Cisco Ecosystem", DEBT_RED)
else:
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # Blank layout
    add_text_box(slide, 0.3, 0.2, 12.7, 0.65,
                 "The Splunkbase Mass Archiving — Impact on Cisco Ecosystem",
                 font_size=26, bold=True, color=DEBT_RED)

# ── Big stat callouts across the top ──
add_stat_callout(slide, 0.2,  1.15, "121",  "Apps & Add-ons\n(Before)", CISCO_BLUE)
add_stat_callout(slide, 2.7,  1.15, "62",   "Archived\n(51%)", DEBT_RED)
add_stat_callout(slide, 5.2,  1.15, "59",   "Remaining\n(After)", WARN_AMBER)
add_stat_callout(slide, 7.7,  1.15, "54",   "Active Today", PHASE_GREEN)
add_stat_callout(slide, 10.2, 1.15, "18",   "Enterprise-Grade\n(Cisco+Splunk)", PHASE_BLUE)

# Arrow connectors (simple text arrows between stats)
add_text_box(slide, 2.25, 1.25, 0.5, 0.5, "→", font_size=28, bold=True,
             color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 4.75, 1.25, 0.5, 0.5, "→", font_size=28, bold=True,
             color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 7.25, 1.25, 0.5, 0.5, "→", font_size=28, bold=True,
             color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 9.75, 1.25, 0.5, 0.5, "→", font_size=28, bold=True,
             color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ── Four section boxes ──
col_w = 3.0
gap = 0.15
x1 = 0.2
x2 = x1 + col_w + gap
x3 = x2 + col_w + gap
x4 = x3 + col_w + gap
box_top = 2.35
box_h = 2.45

# BEFORE
add_section_box(slide, x1, box_top, col_w, box_h,
    "BEFORE the Archive", CISCO_BLUE,
    [
        "**121** Apps & Add-ons on Splunkbase",
        "  • 49 apps + 72 add-ons",
        "**13** SOAR Connectors (tracked separately)",
        "**56** unique Cisco products covered",
        "",
        "**Support breakdown:**",
        "  Cisco 13 | Splunk 8",
        "  Developer 53 | Unsupported 47",
    ], body_size=9)

# ARCHIVED
add_section_box(slide, x2, box_top, col_w, box_h,
    "62 Archived (51%)", DEBT_RED,
    [
        "**22** apps + **40** add-ons removed",
        "",
        "**10** Cisco products lost all coverage",
        "  (all 10 already EOL/EOS by Cisco)",
        "",
        "**4** active products impacted:",
        "  AppDynamics, Webex,",
        "  Meeting Server, CUCM",
        "  (now developer-supported only)",
    ], body_size=9)

# AFTER
add_section_box(slide, x3, box_top, col_w, box_h,
    "AFTER the Archive", WARN_AMBER,
    [
        "**59** Apps & Add-ons remaining",
        "  • 27 apps + 32 add-ons",
        "**46** unique Cisco products",
        "",
        "**Support breakdown:**",
        "  Cisco 11 | Splunk 8  (32%)",
        "  Developer 30 | Unsupported 10  (68%)",
        "",
        "⚠ 68% lack enterprise-grade support",
    ], body_size=9)

# AFTER DEPRECATION — ENTERPRISE REALITY
add_section_box(slide, x4, box_top, col_w, box_h,
    "Enterprise Reality", PHASE_BLUE,
    [
        "**54** Apps & Add-ons in active use",
        "  • 24 apps + 30 add-ons",
        "",
        "**Only 18** Cisco + Splunk supported  (33%)",
        "**36** Developer + Unsupported  (67%)",
        "",
        "Most enterprise customers will only",
        "deploy officially supported apps",
        "→ effective coverage gap is real",
    ], body_size=9)

# ── Key Takeaway box ──
add_rounded_rect(slide, 0.2, 5.0, 12.9, 0.35, PURPLE,
                 "KEY TAKEAWAY", font_size=11, font_color=WHITE)

takeaway_lines = [
    [("51% ", 11, True, DEBT_RED),
     ("of the Cisco Splunkbase ecosystem was archived — ", 11, False, DARK_TEXT),
     ("10 retired products lost all coverage ", 11, True, WARN_AMBER),
     ("(all EOL/EOS)", 11, False, DARK_TEXT)],
    [("4 active products ", 11, True, ARCHIVE_ORANGE),
     ("(AppDynamics, Webex, CMS, CUCM) now rely solely on developer-supported add-ons — ", 11, False, DARK_TEXT),
     ("no Cisco/Splunk SLA", 11, True, DEBT_RED)],
    [("Only 18 of 54 remaining apps (33%) ", 11, True, PHASE_BLUE),
     ("carry Cisco or Splunk support — enterprise customers often treat the other ", 11, False, DARK_TEXT),
     ("67% as 'no coverage'", 11, True, DEBT_RED)],
    [("SOAR connectors (13) ", 11, True, CISCO_BLUE),
     ("were unaffected  •  ", 11, False, DARK_TEXT),
     ("Technical coverage preserved, but enterprise-grade coverage has real gaps", 11, True, PURPLE)],
]
add_rich_text_box(slide, 0.5, 5.4, 12.4, 1.2, takeaway_lines, alignment=PP_ALIGN.LEFT)

# ── Beyond the Archive: GTM Roadmap footnote ──
GTM_GRAY = RGBColor(0x45, 0x45, 0x45)
add_rounded_rect(slide, 0.2, 6.7, 12.9, 0.55, RGBColor(0x37, 0x47, 0x4F),
                 "", font_size=10, font_color=WHITE)
gtm_lines = [
    [("Beyond the Archive: ", 11, True, WARN_AMBER),
     ("16 additional Cisco products on the GTM roadmap have ", 11, False, WHITE),
     ("zero Splunkbase coverage today", 11, True, DEBT_RED)],
    [("These include Hypershield, Secure AI Factory, Security Cloud Control, pxGrid, Industrial Networking, and 11 more.", 9, False, RGBColor(0xB0, 0xBE, 0xC5))],
    [("The archiving didn't cause these gaps — they pre-date it — but the full picture shows ", 9, False, RGBColor(0xB0, 0xBE, 0xC5)),
     ("63 Cisco products ", 9, True, WHITE),
     ("in the portfolio with only ", 9, False, RGBColor(0xB0, 0xBE, 0xC5)),
     ("18 enterprise-grade integrations.", 9, True, WARN_AMBER)],
]
add_rich_text_box(slide, 0.5, 6.72, 12.4, 0.5, gtm_lines, alignment=PP_ALIGN.LEFT)


# ─── Save ───
prs.save(OUTPUT_PATH)
size_kb = os.path.getsize(OUTPUT_PATH) // 1024
print(f"\nSaved: {OUTPUT_PATH}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {size_kb:,} KB")
