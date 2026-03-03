#!/usr/bin/env python3
"""
Generate a Cisco-branded PowerPoint slide showing all Splunk coverage gap products.
These are Cisco products with no existing Splunk TA or app integration.

Usage:
    python3 scripts/gen_coverage_gaps_slide.py

Output:
    docs/Cisco_Splunk_Coverage_Gaps.pptx
"""

import configparser
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PRODUCTS_CONF = os.path.join(
    os.path.dirname(__file__), '..', 'packages', 'splunk-cisco-app-navigator',
    'src', 'main', 'resources', 'splunk', 'default', 'products.conf'
)
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), '..', 'docs', 'Cisco_Splunk_Coverage_Gaps.pptx')

# Cisco brand colors
CISCO_BLUE = RGBColor(0x04, 0x9F, 0xD9)
CISCO_DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_TEXT = RGBColor(0x2A, 0x2A, 0x2A)
RED_ACCENT = RGBColor(0xE5, 0x39, 0x35)
GREEN_ACCENT = RGBColor(0x2E, 0x7D, 0x32)
ORANGE_ACCENT = RGBColor(0xF5, 0x7C, 0x00)

# Category colors & emojis
CAT_COLORS = {
    'security': RGBColor(0xE5, 0x39, 0x35),
    'networking': RGBColor(0x00, 0x89, 0x7B),
    'collaboration': RGBColor(0x7B, 0x1F, 0xA2),
    'observability': RGBColor(0x15, 0x65, 0xC0),
}
CAT_LABELS = {
    'security': 'Security',
    'networking': 'Networking',
    'collaboration': 'Collaboration',
    'observability': 'Observability',
}


def load_coverage_gaps():
    """Read products.conf and return coverage gap products grouped by category."""
    with open(PRODUCTS_CONF) as f:
        raw = '[DEFAULT]\n' + f.read()
    cp = configparser.ConfigParser(interpolation=None)
    cp.read_string(raw)

    gaps = []
    for section in cp.sections():
        if section == 'DEFAULT':
            continue
        if cp.get(section, 'coverage_gap', fallback='') == 'true':
            gaps.append({
                'id': section,
                'name': cp.get(section, 'display_name', fallback=section),
                'category': cp.get(section, 'category', fallback=''),
                'subcategory': cp.get(section, 'subcategory', fallback=''),
                'description': cp.get(section, 'description', fallback=''),
            })

    # Group by category
    grouped = {}
    for g in sorted(gaps, key=lambda x: x['name']):
        cat = g['category']
        if cat not in grouped:
            grouped[cat] = []
        grouped[cat].append(g)

    return grouped, len(gaps)


def add_rect(slide, left, top, width, height, fill_color):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 font_color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Segoe UI'):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def build_presentation():
    grouped, total_count = load_coverage_gaps()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # Blank layout

    # ─── SLIDE 1: Title Slide ───
    slide = prs.slides.add_slide(slide_layout)

    # Dark background
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, CISCO_DARK)

    # Cisco blue accent bar at top
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), CISCO_BLUE)

    # Title
    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                 'Cisco–Splunk Coverage Gaps',
                 font_size=44, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
                 f'{total_count} Cisco Products with Zero Splunk Integration Today',
                 font_size=22, font_color=CISCO_BLUE, bold=False, alignment=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, Inches(2), Inches(4.3), Inches(9), Inches(1.0),
                 'These are established Cisco products that currently have no Splunk Technology Add-on (TA), '
                 'no visualization app, and no sourcetype coverage on Splunkbase.',
                 font_size=14, font_color=RGBColor(0xAA, 0xAA, 0xAA), bold=False,
                 alignment=PP_ALIGN.CENTER)

    # Footer
    add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
                 'Source: Splunk Cisco App Navigator (SCAN) · March 2026',
                 font_size=11, font_color=MEDIUM_GRAY, bold=False, alignment=PP_ALIGN.CENTER)

    # Blue accent bar at bottom
    add_rect(slide, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), CISCO_BLUE)

    # ─── SLIDE 2: Coverage Gaps Table ───
    slide2 = prs.slides.add_slide(slide_layout)

    # Light background
    add_rect(slide2, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)

    # Header bar
    add_rect(slide2, Inches(0), Inches(0), prs.slide_width, Inches(0.9), CISCO_DARK)
    add_text_box(slide2, Inches(0.5), Inches(0.15), Inches(8), Inches(0.6),
                 f'Coverage Gap Products — {total_count} Cisco Products Without Splunk Integration',
                 font_size=22, font_color=WHITE, bold=True)

    # Category columns
    cat_order = ['security', 'networking', 'collaboration', 'observability']
    cats_with_data = [c for c in cat_order if c in grouped]

    col_count = len(cats_with_data)
    margin = 0.4
    gap = 0.3
    total_width = prs.slide_width.inches - (2 * margin) - ((col_count - 1) * gap)
    col_width = total_width / col_count

    y_start = 1.2

    for i, cat in enumerate(cats_with_data):
        x = margin + i * (col_width + gap)
        products = grouped[cat]
        cat_color = CAT_COLORS.get(cat, CISCO_BLUE)
        cat_label = CAT_LABELS.get(cat, cat.title())

        # Category header pill
        header_shape = add_rect(slide2, Inches(x), Inches(y_start), Inches(col_width), Inches(0.45), cat_color)
        add_text_box(slide2, Inches(x), Inches(y_start + 0.05), Inches(col_width), Inches(0.35),
                     f'{cat_label} ({len(products)})',
                     font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Product list
        y = y_start + 0.6
        for product in products:
            # Product card
            card_height = 0.85
            card = add_rect(slide2, Inches(x), Inches(y), Inches(col_width), Inches(card_height), LIGHT_GRAY)
            card.shadow.inherit = False

            # Left accent bar
            add_rect(slide2, Inches(x), Inches(y), Inches(0.05), Inches(card_height), cat_color)

            # Product name
            name = product['name']
            # Trim long names
            if len(name) > 45:
                name = name[:42] + '…'
            add_text_box(slide2, Inches(x + 0.12), Inches(y + 0.05), Inches(col_width - 0.2), Inches(0.35),
                         name, font_size=11, font_color=DARK_TEXT, bold=True)

            # Description (truncated)
            desc = product['description']
            if len(desc) > 90:
                desc = desc[:87] + '…'
            if desc:
                add_text_box(slide2, Inches(x + 0.12), Inches(y + 0.35), Inches(col_width - 0.2), Inches(0.45),
                             desc, font_size=8, font_color=MEDIUM_GRAY, bold=False)

            y += card_height + 0.12

    # Footer
    add_rect(slide2, Inches(0), Inches(7.1), prs.slide_width, Inches(0.4), RGBColor(0xF0, 0xF0, 0xF0))
    add_text_box(slide2, Inches(0.5), Inches(7.12), Inches(12), Inches(0.35),
                 '⚠ No Splunk TA · No Splunk App · No Sourcetypes — These products represent GTM roadmap opportunities for new Cisco–Splunk integrations.',
                 font_size=10, font_color=MEDIUM_GRAY, bold=False)

    # Bottom accent
    add_rect(slide2, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), CISCO_BLUE)

    # Save
    prs.save(OUTPUT_PATH)
    print(f'✅ Saved: {OUTPUT_PATH}')
    print(f'   {total_count} coverage gap products across {len(cats_with_data)} categories')


if __name__ == '__main__':
    build_presentation()
