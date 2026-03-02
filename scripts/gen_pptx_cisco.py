#!/usr/bin/env python3
"""
Generate Cisco-Splunk TA Strategy Presentation as PPTX.
Uses the official Cisco PowerPoint Template (LIGHT).

Template conversion: .potx → .pptx (change content type in [Content_Types].xml)
"""

import shutil
import zipfile
import os
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.shapes import MSO_SHAPE

# ─── Paths (relative to this script's location: scripts/) ───
_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
_REPO_ROOT = os.path.join(_SCRIPT_DIR, '..')
TEMPLATE_POTX = os.path.join(_REPO_ROOT, 'docs', 'Cisco_PowerPoint_Template_LIGHT.potx')
TEMPLATE_PPTX = os.path.join(_REPO_ROOT, 'docs', 'cisco_template_converted.pptx')
OUTPUT_PATH = os.path.join(_REPO_ROOT, 'docs', 'Cisco_Splunk_TA_Strategy_Presentation.pptx')

# ─── Convert .potx → .pptx ───
def convert_potx_to_pptx(potx_path, pptx_path):
    shutil.copy2(potx_path, pptx_path)
    tmpdir = tempfile.mkdtemp()
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        zin.extractall(tmpdir)
    ct_path = os.path.join(tmpdir, '[Content_Types].xml')
    with open(ct_path, 'r', encoding='utf-8') as f:
        content = f.read()
    content = content.replace(
        'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
    )
    with open(ct_path, 'w', encoding='utf-8') as f:
        f.write(content)
    with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for root, dirs, files in os.walk(tmpdir):
            for fname in files:
                full_path = os.path.join(root, fname)
                arcname = os.path.relpath(full_path, tmpdir)
                zout.write(full_path, arcname)
    shutil.rmtree(tmpdir)
    print(f"Converted template: {pptx_path}")

convert_potx_to_pptx(TEMPLATE_POTX, TEMPLATE_PPTX)

# ─── Load template ───
prs = Presentation(TEMPLATE_PPTX)

# Remove all pre-existing slides from the template (keep only layouts/masters)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId is None:
        # Try attribute name 'r:id'
        rId_attr = prs.slides._sldIdLst[0].attrib
        for k, v in rId_attr.items():
            if k.endswith('}id') or k == 'r:id':
                rId = v
                break
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
print(f"Template slides removed. Starting with {len(prs.slides)} slides.")

# ─── Cisco Template Layout Map ───
# [0]  Title Slide 1, Two Speakers — title + subtitle + 2 speaker positions
# [8]  Agenda 1 — agenda with section markers
# [9]  Title Only 1 — clean: just title + footer/slide number (our workhorse)
# [10] Title, Subtitle Only 1 — title + subtitle
# [19] Title, Subtitle, 2 Columns — title + sub + left/right body
# [37] Quote 1 — large quote with attribution
# [41] Blank — fully blank (footer + page number only)
# [47] Segue 1 — section dividers
# [51] Thank you 1 — official thank you

LAYOUT_TITLE = 0        # Title Slide 1, Two Speakers
LAYOUT_AGENDA = 9        # We'll use Title Only — Agenda 1 has odd placeholder arrangement
LAYOUT_TITLE_ONLY = 9    # Title Only 1 (workhorse — clean canvas with title)
LAYOUT_TITLE_SUB = 10    # Title, Subtitle Only 1
LAYOUT_TWO_COL = 19      # Title, Subtitle, 2 Columns
LAYOUT_QUOTE = 37        # Quote 1
LAYOUT_BLANK = 41        # Blank
LAYOUT_SEGUE = 47        # Segue 1
LAYOUT_THANKYOU = 51     # Thank you 1

# ─── Colors ───
CISCO_BLUE   = RGBColor(0x04, 0x9F, 0xD9)
SPLUNK_GREEN = RGBColor(0x65, 0xA6, 0x37)
DEBT_RED     = RGBColor(0xD3, 0x2F, 0x2F)
WARN_AMBER   = RGBColor(0xF9, 0xA8, 0x25)
DARK_BG      = RGBColor(0x1A, 0x1A, 0x2E)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY     = RGBColor(0x9E, 0x9E, 0x9E)
BLACK        = RGBColor(0x00, 0x00, 0x00)
DARK_TEXT     = RGBColor(0x21, 0x21, 0x21)
PHASE_GREEN  = RGBColor(0x2E, 0x7D, 0x32)
PHASE_BLUE   = RGBColor(0x15, 0x65, 0xC0)
PURPLE       = RGBColor(0x4A, 0x14, 0x8C)

# ─── Helpers ───

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name=None):
    """Add a text box. font_name=None means inherit from template (Cisco Sans)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    if font_name:
        p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, bold_items=None, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = spacing
        if bold_items and i in bold_items:
            p.font.bold = True
    return txBox

def add_table(slide, left, top, width, rows_data, col_widths=None,
              header_color=CISCO_BLUE, font_size=11):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(0.3 * rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    return table_shape

def add_stat_callout(slide, left, top, number, label, color=DEBT_RED):
    add_text_box(slide, left, top, 2.5, 0.8, number, font_size=42, bold=True,
                 color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + 0.7, 2.5, 0.5, label, font_size=14,
                 color=MED_GRAY, alignment=PP_ALIGN.CENTER)

def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
    return shape

def set_title(slide, text):
    """Set the title placeholder text if available (idx=0)."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = text
            for paragraph in ph.text_frame.paragraphs:
                paragraph.font.color.rgb = DARK_TEXT
            return ph
    return None


# ==========================================
# SLIDE 1: Title
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
# Layout [0]: Title Slide 1, Two Speakers
# idx=0: Title, idx=13: subtitle line, idx=16: Speaker 1, idx=14: Speaker 2, idx=12: date/footer
for ph in slide.placeholders:
    idx = ph.placeholder_format.idx
    if idx == 0:
        ph.text = "Cisco-Splunk TA Strategy"
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(36)
    elif idx == 13:
        ph.text = "Findings, Challenges & The Path Forward"
    elif idx == 16:
        ph.text = "Alan Ivarson\nGroup Vice President, Field Solutions"
    elif idx == 14:
        ph.text = "Amir (AK) Khamis\nPrincipal Architect"
    elif idx == 12:
        ph.text = "February 2026  •  45-Minute Presentation"

# ==========================================
# SLIDE 2: Agenda
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
set_title(slide, "Agenda")

agenda_items = [
    "1.  The Problem — 110 Cisco Apps, No Clear Answer",
    "2.  The Confusion Problem — 13 Firewall Apps, 1 Correct Answer",
    "3.  Stakeholder Feedback — Top 12 Challenges & The Magic 6",
    "4.  Option 1: Hire & Build — Doesn't Scale",
    "5.  Option 2: Super App Bundling — Right Direction, Still Expensive",
    "6.  Two Deliverables — Quick Win (SCAN) + Strategic Play (Super App)",
    "7.  The AI-Powered Path Forward",
    "8.  Next Steps & Roadmap",
]
add_bullet_list(slide, 0.8, 1.3, 10, 5.5, agenda_items, font_size=20, color=DARK_TEXT,
                bold_items=set(range(8)), spacing=Pt(14))


# ==========================================
# SLIDE 3: The Problem — 110 Apps + Key Insights
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
set_title(slide, "1. The Problem — 110 Apps, No Clear Answer")
# Color the title red
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        for p in ph.text_frame.paragraphs:
            p.font.color.rgb = DEBT_RED

# CDF Last Mile context
add_text_box(slide, 0.35, 0.95, 12.6, 0.45,
             'Per Kunal Mukerjee (SVP Engineering): Splunk integrations are the "Last Mile" of the Cisco Data Fabric (CDF). '
             '125 listings audited; 13 SOAR connectors + 2 prerequisites excluded = 110 apps & add-ons.',
             font_size=10, bold=True, color=CISCO_BLUE)

# 4 big stat callouts
add_stat_callout(slide, 0.3, 1.45, "47%", "Unsupported", DEBT_RED)
add_stat_callout(slide, 3.3, 1.45, "60%", "Dead / Dying", WARN_AMBER)
add_stat_callout(slide, 6.3, 1.45, "50%", "Community-Built", MED_GRAY)
add_stat_callout(slide, 9.3, 1.45, "74%", "Non-Functional", PURPLE)

# KEY INSIGHTS TABLE
add_text_box(slide, 0.35, 2.85, 12.6, 0.3, "Key Insights — The Numbers That Should Surprise Everyone",
             font_size=14, bold=True, color=CISCO_BLUE)

insights_data = [
    ["#", "Insight", "Stat", "Why It Matters"],
    ["1", "Ecosystem is mostly dead", "66 of 110 (60%) archived/deprecated/flagged", "Customers navigate a graveyard"],
    ["2", "Almost no one has your back", "Only 19 of 110 (17%) Cisco+Splunk supported", "83% run on hope, not SLAs"],
    ["3", "Half is community-built", "55 of 110 (50%), 48 different developers", "No QA, no security review"],
    ["4", "Massive consolidation", "52 apps \u2192 absorbed by just 15", "App 7404 alone replaces 18"],
    ["5", "16 ticking time bombs", "16 of 49 live apps (33%) = zero support", "1/3 of live apps have no owner"],
    ["6", "5 zombie apps", "Deprecated but still downloadable", "Customers install dead apps daily"],
    ["7", "Dev > Cisco+Splunk", "39 dev-supported vs 19 official (2\u00d7)", "Volunteers carry our brand"],
    ["8", "Only 17% backed", "Cisco: 13, Splunk: 6, everyone else: 91", "$28B acquisition, <1 in 5 backed"],
]
add_table(slide, 0.2, 3.15, 12.9, insights_data, col_widths=[0.4, 2.8, 4.5, 4.8],
          header_color=DEBT_RED, font_size=9)

# Punchline
add_rounded_rect(slide, 0.2, 6.15, 12.9, 0.5, PURPLE,
                 '74% of Cisco\u2019s Splunkbase ecosystem is non-functional, unsupported, or flagged.  '
                 'SCAN is the aspirin. The Super App is the cure.', 11, WHITE)

# Financial impact
add_text_box(slide, 0.35, 6.75, 12.6, 0.3,
             "Financial Impact: Fragmented TAs \u2192 thousands of wasted TAC hours, delayed Time-to-Value for multi-million dollar EAs.",
             font_size=10, bold=True, color=DEBT_RED)


# ==========================================
# SLIDE 4: The Confusion Problem — Firewall
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
set_title(slide, "2. The Confusion Problem — 13 Apps, 1 Correct Answer")
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        for p in ph.text_frame.paragraphs:
            p.font.color.rgb = DEBT_RED

add_text_box(slide, 0.35, 1.0, 11, 0.4,
             'A customer searching Splunkbase for "Cisco Firewall" encounters 13 different results:',
             font_size=14, color=DARK_TEXT)

firewall_data = [
    ["#", "App Name", "Status", "Built By"],
    ["1", "Splunk Add-on for Cisco FireSIGHT", "Archived", "Splunk"],
    ["2", "Cisco Secure eStreamer Client Add-On", "Legacy", "Cisco"],
    ["3", "Cisco Firepower eNcore App", "Legacy", "Community"],
    ["4", "Cisco Secure Firewall App for Splunk", "Superseded", "Cisco"],
    ["5", "Splunk Add-on for Cisco ASA", "Legacy", "Splunk"],
    ["6", "Enosys eStreamer Add-on", "Legacy", "Community"],
    ["7", "Cisco FTD sourcetype TA", "Legacy", "Community"],
    ["8", "Cisco FTD Dashboards", "Legacy", "Community"],
    ["9", "Cisco Firepower pcap Add-on", "Legacy", "Community"],
    ["10", "Cisco eStreamer Client for Splunk", "Legacy", "Cisco"],
    ["11", "Firegen Log Analyzer for Cisco ASA", "Legacy", "Community"],
    ["12", "Cisco Security Suite", "Legacy", "Community"],
    ["13", "Cisco Suite for Splunk", "Legacy", "Community"],
]
add_table(slide, 0.35, 1.5, 8, firewall_data, col_widths=[0.4, 4.5, 1.3, 1.5],
          header_color=DEBT_RED, font_size=10)

add_rounded_rect(slide, 8.8, 1.8, 4.0, 1.5, SPLUNK_GREEN,
                 "CORRECT ANSWER:\n\nInstall Cisco\nSecurity Cloud\n(1 add-on)", font_size=16)

add_text_box(slide, 8.8, 3.8, 4.0, 0.4, "Top 10 Worst Offenders:",
             font_size=14, bold=True, color=DARK_TEXT)
offenders = [
    "Firewall: 13 legacy apps",
    "Webex: 8",
    "Meraki: 5  •  Catalyst Center: 5",
    "Nexus: 5  •  Umbrella: 4",
    "NVM: 3  •  SNA: 3  •  ISE: 3",
]
add_bullet_list(slide, 8.8, 4.2, 4.0, 2.5, offenders, font_size=11, color=DARK_TEXT,
                bold_items={0}, spacing=Pt(4))

add_text_box(slide, 0.35, 6.7, 12.6, 0.4,
             "93 total legacy entries across 29 of 57 products (51%).  Wrong app = blank dashboards, broken lookups, duplicate sourcetypes.",
             font_size=11, bold=True, color=DEBT_RED)


# ==========================================
# SLIDE 5: Stakeholder Feedback + Challenges + Magic 6
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
set_title(slide, "3. Stakeholder Feedback — 20+ SMEs Consulted")

# Named SME quotes — left
add_text_box(slide, 0.35, 1.0, 6.3, 0.35,
             '\u201cThe current state is a disaster resulting from a lack of holistic strategy.\u201d',
             font_size=10, bold=True, color=DEBT_RED)
add_text_box(slide, 0.35, 1.3, 6.3, 0.2,
             '\u2014 James Young, Principal Global Security Product Specialist',
             font_size=9, color=MED_GRAY)
add_text_box(slide, 0.35, 1.55, 6.3, 0.35,
             '\u201cOne Cisco marketing is outpacing our technical execution \u2014 customers are calling us out.\u201d',
             font_size=10, bold=True, color=DEBT_RED)
add_text_box(slide, 0.35, 1.85, 6.3, 0.2,
             '\u2014 Bryan Pluta, Principal Solutions Engineer',
             font_size=9, color=MED_GRAY)

# SME quotes — right
add_text_box(slide, 7.0, 1.0, 5.8, 0.35,
             '\u201cCyberCX built their own TAs for 7 years \u2014 we risk losing mindshare.\u201d',
             font_size=10, bold=True, color=DEBT_RED)
add_text_box(slide, 7.0, 1.3, 5.8, 0.2,
             '\u2014 Daniel Peluso, ANZ Splunk Partner SE Manager',
             font_size=9, color=MED_GRAY)
add_text_box(slide, 7.0, 1.55, 5.8, 0.35,
             '\u201cTA sprawl is a hurdle for the Cisco sales force \u2014 we need a turnkey experience.\u201d',
             font_size=10, bold=True, color=DEBT_RED)
add_text_box(slide, 7.0, 1.85, 5.8, 0.2,
             '\u2014 Matt Poland, Sr. Director, SE',
             font_size=9, color=MED_GRAY)

challenges_data = [
    ["#", "Challenge", "Evidence"],
    ["1", "Too many apps", "13 for Firewall; 110 total"],
    ["2", "47% unsupported", "52 apps with zero support"],
    ["3", "51% archived but still referenced", "Old blog posts \u2192 dead apps"],
    ["4", "50% community-built, no QA", "48 developers; CIM conflicts"],
    ["5", "Blank dashboards", "Install \u2192 empty panels \u2192 TAC"],
    ["6", "Version fragmentation", "No in-product update alerts"],
    ["7", "No migration path", "Legacy coexists with modern"],
    ["8", "Duplicate sourcetypes", "Two TAs claim cisco:asa"],
    ["9", "Platform incompatibility", "Cloud vs Enterprise"],
    ["10", "Knowledge gap", "Splunk + Python + Cisco rare"],
    ["11", "Support Loop", "Cisco TAC requires contract; admins bounced"],
    ["12", "Competitive risk", "Palo Alto AI parsing; partners use Microsoft"],
]
add_table(slide, 0.35, 2.15, 6.3, challenges_data, col_widths=[0.4, 2.7, 3.0],
          header_color=DEBT_RED, font_size=9)

add_rounded_rect(slide, 0.35, 5.95, 6.3, 0.6, PURPLE,
                 '\u26a0 BRAND RISK: 50% community-built (48 developers) = Cisco doesn\'t own its telemetry story.', 10, WHITE)

add_rounded_rect(slide, 0.35, 6.6, 12.6, 0.5, DARK_BG,
                 'END-TO-END MANDATE: Data \u2192 CIM \u2192 Dashboards \u2192 ESCU Detections \u2192 SOAR Playbooks', 10, WHITE)

# Magic 6 — right side
add_text_box(slide, 7.2, 2.1, 5.5, 0.5, 'The "Magic 6" Best Practices',
             font_size=20, bold=True, color=SPLUNK_GREEN)
magic6_data = [
    ["#", "Best Practice", "Impact"],
    ["1", "Unified add-on family", "1 install, no conflicts"],
    ["2", "Dynamic Config Engine (DCE)", "Zero-footprint, no bloat"],
    ["3", "Sourcetype validation (MTTI)", "Prove data flows \u2192 reduce TAC"],
    ["4", "Legacy debt detection", "Auto-flag what to remove"],
    ["5", "Platform awareness", "Right config every time"],
    ["6", "Centralized catalog", "One pane of glass"],
]
add_table(slide, 7.2, 2.7, 5.5, magic6_data, col_widths=[0.4, 2.5, 2.3],
          header_color=SPLUNK_GREEN, font_size=10)

add_text_box(slide, 7.2, 5.5, 5.5, 0.6,
             "If every TA followed these 6 principles,\nthe top 12 challenges would disappear.",
             font_size=14, bold=True, color=SPLUNK_GREEN)


# ==========================================
# SLIDE 6: Option 1 — Can't Hire Our Way Out
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
set_title(slide, "4. Option 1: Hire & Build — Doesn't Scale")

scale_data = [
    ["Scenario", "Calculation", "Result"],
    ["Build 1 TA", "3\u20134 engineer-months", "Feasible"],
    ["23 TAs (current)", "69\u201392 eng-months", "6\u20138 engineer-years"],
    ["57 products", "171\u2013228 eng-months", "14\u201319 engineer-years"],
    ["400 data sources", "1,200\u20131,600 eng-months", "100\u2013133 engineer-years"],
]
add_table(slide, 0.35, 1.2, 6, scale_data, col_widths=[2.0, 2.2, 2.2],
          header_color=DARK_BG, font_size=11)

# Bar chart
chart_data_bar = CategoryChartData()
chart_data_bar.categories = ['1 TA', '23 TAs', '57 Products', '400 Sources']
chart_data_bar.add_series('Engineer-Years', (0.3, 7, 17, 133))
chart_frame_bar = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(7), Inches(1.0), Inches(5.5), Inches(3.5), chart_data_bar
)
chart_bar = chart_frame_bar.chart
chart_bar.has_legend = False
chart_bar.value_axis.maximum_scale = 140
chart_bar.value_axis.has_title = True
chart_bar.value_axis.axis_title.text_frame.paragraphs[0].text = "Engineer-Years"
chart_bar.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
series = chart_bar.series[0]
bar_colors = [SPLUNK_GREEN, WARN_AMBER, DEBT_RED, RGBColor(0x8B, 0x00, 0x00)]
for i, color in enumerate(bar_colors):
    point = series.points[i]
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = color

# Knowledge Venn
add_text_box(slide, 0.35, 4.2, 6, 0.5, "The Knowledge Problem — 4 Domains, Very Few Unicorns",
             font_size=16, bold=True, color=DARK_TEXT)
add_rounded_rect(slide, 0.5, 4.8, 2.8, 0.7, CISCO_BLUE, "Splunk UCC Framework", 11, WHITE)
add_rounded_rect(slide, 3.5, 4.8, 2.8, 0.7, PHASE_GREEN, "Python & Scripting", 11, WHITE)
add_rounded_rect(slide, 0.5, 5.6, 2.8, 0.7, RGBColor(0x6A, 0x1B, 0x9A), "Splunk Platform\nprops / CIM / SPL", 10, WHITE)
add_rounded_rect(slide, 3.5, 5.6, 2.8, 0.7, DEBT_RED, "Cisco Domain\nAPIs / Syslog / MIBs", 10, WHITE)
add_rounded_rect(slide, 7.0, 5.1, 3.0, 0.9, WARN_AMBER, "\u2b50 THE UNICORN\nKnows All Four", 14, BLACK)

for (x, y) in [(3.2, 5.15), (6.0, 5.15), (3.2, 5.95), (6.0, 5.95)]:
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(0.6), Inches(0.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MED_GRAY
    shape.line.fill.background()

add_text_box(slide, 0.35, 6.5, 12.6, 0.5,
             "Conclusion: We cannot hire our way out of this problem.",
             font_size=18, bold=True, color=DEBT_RED, alignment=PP_ALIGN.CENTER)


# ==========================================
# SLIDE 7: Option 2 — Super App Bundling
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
set_title(slide, "5. Option 2: Super App Bundling — Right Direction, Still Expensive")
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(22)

consol_data = [
    ["Super Add-on", "Products", "Ratio"],
    ["Cisco Security Cloud", "19 security", "19 : 1"],
    ["Cisco Catalyst Add-on", "13 networking", "13 : 1"],
    ["Cisco DC Networking", "3 data center", "3 : 1"],
    ["Total", "35 \u2192 3 add-ons", "~12 : 1"],
]
add_table(slide, 0.35, 1.2, 5.5, consol_data, col_widths=[2.5, 1.8, 1.0],
          header_color=CISCO_BLUE, font_size=11)

ba_data = [
    ["Metric", "Before", "After"],
    ["Customer installs", "23 separate TAs", "3 Super Add-ons"],
    ["Upgrade cycles", "23 paths", "3"],
    ["Splunkbase confusion", "110 apps", "3 clear choices"],
    ["Version conflicts", "Constant", "Eliminated"],
    ["Support", "Fragmented", "Unified"],
]
add_table(slide, 6.5, 1.2, 6.3, ba_data, col_widths=[2.0, 2.0, 2.0],
          header_color=SPLUNK_GREEN, font_size=11)

add_text_box(slide, 0.35, 4.0, 12.6, 0.5, "But It Still Has Challenges:",
             font_size=18, bold=True, color=WARN_AMBER)
remaining = [
    "Maintenance cost \u2014 many mini-TAs inside; still need per-product parsing updates",
    "Regression risk \u2014 updating Firewall parsing could break Duo parsing",
    "Release velocity \u2014 can't ship a fix without regression-testing all products",
    "Blank dashboard problem \u2014 app installed \u2260 data flowing; static configs validate nothing",
    "Knowledge bottleneck \u2014 still need engineers who understand each product's data format",
]
add_bullet_list(slide, 0.8, 4.5, 11.5, 2.5, remaining, font_size=13, color=DARK_TEXT, spacing=Pt(6))

add_text_box(slide, 0.35, 6.7, 12.6, 0.5,
             "Super App is the right architecture, but doesn't solve cost, speed, and quality by itself.",
             font_size=16, bold=True, color=WARN_AMBER, alignment=PP_ALIGN.CENTER)


# ==========================================
# SLIDE 8: Two Deliverables
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
set_title(slide, "6. Two Deliverables — Quick Win + Strategic Play")
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        for p in ph.text_frame.paragraphs:
            p.font.color.rgb = CISCO_BLUE

# SCAN (left)
add_rounded_rect(slide, 0.35, 1.2, 6.0, 0.9, WARN_AMBER,
                 'DELIVERABLE 1: Splunk Cisco App Navigator (SCAN)\n"The Front Door"  \u2022  Ships in Weeks', 14, BLACK)
ccc_items = [
    "\u2705  Catalog 57 products with correct add-on for each",
    "\u2705  Monitor 321 sourcetypes \u2014 Proof of Innocence (MTTI)",
    "\u2705  Auto-detect 93 legacy apps and warn what to remove",
    "\u2705  Map SOAR (10), ITSI (4), alert actions (3)",
    "\u2705  One-click dashboard launch for 26 products",
    "\u2705  Platform-aware guidance (Cloud vs Enterprise)",
    "",
    '\u201cThe network is innocent until proven guilty \u2014',
    'SCAN provides the proof in seconds.\u201d',
    "Lightweight: no conf generation, no data collection",
    "Splunkbase in 2\u20133 weeks (AppInspect + packaging)",
]
add_bullet_list(slide, 0.5, 2.25, 5.7, 3.5, ccc_items, font_size=11, color=DARK_TEXT, spacing=Pt(3))

# Super App (right)
add_rounded_rect(slide, 6.9, 1.2, 6.1, 0.9, PHASE_GREEN,
                 'DELIVERABLE 2: Enterprise Networking Super App\n"The Engineering Solution"  \u2022  3\u20136 Months', 14, WHITE)
sa_items = [
    "Dynamic Configuration Engine (DCE) \u2014 just-in-time activation",
    "Master Library \u2014 gold-standard mini-TAs per product",
    "Logical Isolation \u2014 Meraki error can't crash ASA or ISE",
    "Tier-Aware \u2014 Full Stack / Heavy Forwarder / Search Head Only",
    "Zero-footprint \u2014 nothing active until admin enables",
    "6 products onboarded (ASA, DC Net, Intersight, Meraki, UCS, Webex)",
    "\U0001f512 Security: removes 93 legacy apps \u2192 reduced attack surface",
    "",
    "CDF Alignment \u2014 mini-TAs = Data Products for Data Fabric",
    "MCP Integration \u2014 AI-ready telemetry for autonomous agents",
    "NetFlow/IPFIX \u2014 first-class foundational signal",
]
add_bullet_list(slide, 7.1, 2.25, 5.7, 3.5, sa_items, font_size=11, color=DARK_TEXT, spacing=Pt(3))

# Aspirin / Cure tagline
add_text_box(slide, 0.35, 5.3, 12.6, 0.5,
             'SCAN is the aspirin.  The Super App is the cure.',
             font_size=20, bold=True, color=CISCO_BLUE, alignment=PP_ALIGN.CENTER)

# Comparison table
compare_data = [
    ["", "SCAN (Quick Win)", "Super App (Strategic)"],
    ["Purpose", "Guide to the right TA", "Replace legacy TAs entirely"],
    ["Ships", "2\u20133 weeks", "3\u20136 months (v1)"],
    ["Customer value", "What do I install?", "One app, activate what I need"],
    ["Effort", "Low \u2014 packaging", "Medium \u2014 expand & harden"],
    ["Risk", "Very low", "Moderate"],
    ["Solves confusion", "\u2705 Immediately", "\u2705 Permanently"],
    ["Solves blank dashboards", "\u2705 Detects missing data", "\u2705 Active products only"],
    ["Solves legacy debt", "\u2705 Warns about legacy", "\u2705 Replaces entirely"],
]
add_table(slide, 1.2, 5.85, 11, compare_data, col_widths=[2.5, 3.8, 3.8],
          header_color=CISCO_BLUE, font_size=9)


# ==========================================
# SLIDE 9: AI-Powered Path Forward
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
set_title(slide, "7. The AI-Powered Path Forward")

add_text_box(slide, 0.35, 1.0, 12.6, 0.45,
             'Leveraging FDSE (Rafal Piekarz) for AI-driven maintenance + HCL for validation.\n'
             '\u201cPalo Alto has already implemented AI-driven parsing, proving viability.\u201d \u2014 Dimitri McKay',
             font_size=12, color=MED_GRAY)

# Key enabler callouts
add_rounded_rect(slide, 6.5, 5.2, 3.0, 0.7, CISCO_BLUE,
                 'TACO Framework\nCertification in <24hrs\n(vs. multi-week manual)', 9, WHITE)
add_rounded_rect(slide, 9.7, 5.2, 3.2, 0.7, SPLUNK_GREEN,
                 'Activation Telemetry\nper Kunal: every DCE toggle\nemits ROI proof data', 9, WHITE)

ai_data = [
    ["Step", "Traditional", "AI-Assisted"],
    ["Data source analysis", "1\u20132 weeks", "Hours"],
    ["Input configuration", "1\u20132 weeks", "Hours"],
    ["Parsing (props/transforms)", "2\u20134 weeks", "Hours"],
    ["CIM mapping", "1\u20132 weeks", "Hours"],
    ["Dashboards", "2\u20134 weeks", "Hours"],
    ["Human validation (HCL)", "\u2014", "1\u20132 weeks"],
    ["Certification", "1\u20132 weeks", "1 week"],
    ["TOTAL", "3\u20134 months", "2\u20133 weeks"],
]
add_table(slide, 0.35, 1.6, 5.5, ai_data, col_widths=[2.5, 1.5, 1.5],
          header_color=DARK_BG, font_size=11)

# Bar chart — approaches comparison
chart_data_ap = CategoryChartData()
chart_data_ap.categories = ['Option 1:\nHire (57)', 'Option 1:\nHire (400)', 'Option 2:\nStatic Super', 'Our Approach:\nAI + DCE']
chart_data_ap.add_series('Engineer-Years', (17, 133, 7, 1))
chart_frame_ap = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(6.5), Inches(1.5), Inches(6.3), Inches(3.5), chart_data_ap
)
chart_ap = chart_frame_ap.chart
chart_ap.has_legend = False
chart_ap.value_axis.maximum_scale = 140
chart_ap.value_axis.has_title = True
chart_ap.value_axis.axis_title.text_frame.paragraphs[0].text = "Engineer-Years"
chart_ap.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
series_ap = chart_ap.series[0]
ap_colors = [WARN_AMBER, DEBT_RED, CISCO_BLUE, SPLUNK_GREEN]
for i, color in enumerate(ap_colors):
    point = series_ap.points[i]
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = color

# Options comparison table
options_data = [
    ["", "Option 1: Hire", "Option 2: Static", "Our Approach"],
    ["Architecture", "1 TA per product", "Bundled, static", "Dynamic Super App + DCE"],
    ["Intelligence", "None", "None", "SCAN + MTTI"],
    ["AI authoring", "No", "No", "AI drafts, humans validate"],
    ["Cost", "$$$$$ (100+ eng-yrs)", "$$$ (6\u20138 eng-yrs)", "$ (<1 eng-yr)"],
    ["Speed", "3\u20134 months/TA", "2\u20133 months", "2\u20133 weeks"],
    ["Scale to 400", "Impossible", "Very expensive", "Feasible"],
]
add_table(slide, 0.35, 5.3, 12.6, options_data, col_widths=[1.8, 3.2, 3.2, 3.6],
          header_color=CISCO_BLUE, font_size=10)


# ==========================================
# SLIDE 10: Next Steps & Roadmap
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
set_title(slide, "8. Next Steps & Roadmap")

actions_data = [
    ["#", "Action", "Owner", "Timeline"],
    ["1", "Ship SCAN to Splunkbase \u2014 package, AppInspect, certify", "AK", "2\u20133 weeks"],
    ["2", "Expand Super App Master Library (next 5 products)", "AK + HCL", "0\u20133 months"],
    ["3", "Pilot AI-assisted TA authoring on 3 products", "AK + HCL", "0\u20133 months"],
    ["\u26a0", "DECISION: Approve FDSE partnership (Rafal Piekarz) for AI-driven TA maintenance", "Alan / Exec", "This quarter"],
    ["\u26a0", "DECISION: Allocate HCL validation resources (2\u20133 engineers) for QA", "Alan / Exec", "This quarter"],
    ["6", "Coordinate exec summary with Kunal Mukerjee \u2014 align CDF/MDL", "Alan / AK", "Next 2 weeks"],
]
add_table(slide, 0.35, 1.1, 12.6, actions_data, col_widths=[0.4, 6.5, 2.5, 2.0],
          header_color=CISCO_BLUE, font_size=11)

# Phase boxes
add_text_box(slide, 0.35, 3.8, 12.6, 0.5, "Value Realization \u2014 Three-Phase Journey",
             font_size=18, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Phase 1
add_rounded_rect(slide, 0.35, 4.4, 3.8, 2.5, WARN_AMBER, "", 1, BLACK)
add_text_box(slide, 0.5, 4.5, 3.5, 0.4, 'Phase 1: "Stop the Bleeding"',
             font_size=14, bold=True, color=BLACK)
add_text_box(slide, 0.5, 4.9, 3.5, 0.3, '0\u20133 months  \u2022  Immediate Visibility',
             font_size=10, color=RGBColor(0x55, 0x55, 0x55))
add_bullet_list(slide, 0.5, 5.3, 3.5, 1.5, [
    "Ship SCAN to Splunkbase", "Expand Master Library to 12+",
    "Pilot AI authoring w/ FDSE", "Migrate top 10 legacy offenders",
], font_size=10, color=BLACK, spacing=Pt(3))

# Phase 2
add_rounded_rect(slide, 4.55, 4.4, 3.8, 2.5, PHASE_GREEN, "", 1, WHITE)
add_text_box(slide, 4.7, 4.5, 3.5, 0.4, 'Phase 2: "Scale the Portfolio"',
             font_size=14, bold=True, color=WHITE)
add_text_box(slide, 4.7, 4.9, 3.5, 0.3, '3\u20136 months  \u2022  Operational Efficiency',
             font_size=10, color=RGBColor(0xC8, 0xE6, 0xC9))
add_bullet_list(slide, 4.7, 5.3, 3.5, 1.5, [
    "Ship Super App v1", "AI pipeline: 5\u201310 mini-TAs/month",
    "NetFlow/IPFIX hardened", "HCL validation cadence",
], font_size=10, color=WHITE, spacing=Pt(3))

# Phase 3
add_rounded_rect(slide, 8.75, 4.4, 4.2, 2.5, PHASE_BLUE, "", 1, WHITE)
add_text_box(slide, 8.9, 4.5, 3.9, 0.4, 'Phase 3: "Own the Data Fabric"',
             font_size=14, bold=True, color=WHITE)
add_text_box(slide, 8.9, 4.9, 3.9, 0.3, '6\u201312 months  \u2022  Strategic Differentiation',
             font_size=10, color=RGBColor(0xBB, 0xDE, 0xFB))
add_bullet_list(slide, 8.9, 5.3, 3.9, 1.5, [
    "Full CDF/MDL integration", "MCP-driven AI features",
    "200+ data sources covered", 'SCAN + Super App = "One Cisco"',
], font_size=10, color=WHITE, spacing=Pt(3))

# Arrows between phases
for x in [4.15, 8.35]:
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(5.5), Inches(0.4), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MED_GRAY
    shape.line.fill.background()


# ==========================================
# SLIDE 11: Summary
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
# Manual dark background for Summary
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BG

add_text_box(slide, 0.35, 0.3, 12.6, 0.7, "Summary",
             font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

summary_data = [
    ["The Problem", "Our Answer"],
    ["110 apps, customers confused", "CCC: ships in weeks, shows the right answer"],
    ["Blank dashboards, no validation", "SCAN monitors 321 sourcetypes; Super App activates what has data"],
    ["Legacy apps cause conflicts", "SCAN detects 93 legacy; Super App replaces entirely"],
    ["TAs cost 3\u20134 months each", "AI authoring: 2\u20133 weeks per product"],
    ["Can't hire (400 = 133 eng-yrs)", "AI drafts, humans validate, Super App scales"],
    ["Static Super Apps have regression risk", "DCE: each product is an isolated mini-TA"],
]
add_table(slide, 0.8, 1.3, 11.5, summary_data, col_widths=[4.5, 7.0],
          header_color=CISCO_BLUE, font_size=12)

add_text_box(slide, 0.35, 4.8, 12.6, 0.5, "The Consolidation Journey",
             font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, 0.8, 5.4, 3.5, 1.2, DEBT_RED,
                 "TODAY\n110 apps  \u2022  60% dead\nNo guidance", 13, WHITE)
shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.4), Inches(5.7), Inches(0.6), Inches(0.4))
shape.fill.solid()
shape.fill.fore_color.rgb = MED_GRAY
shape.line.fill.background()

add_rounded_rect(slide, 5.1, 5.4, 3.5, 1.2, WARN_AMBER,
                 'QUICK WIN: SCAN\n"The Front Door"\nShips in weeks', 13, BLACK)
shape2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.7), Inches(5.7), Inches(0.6), Inches(0.4))
shape2.fill.solid()
shape2.fill.fore_color.rgb = MED_GRAY
shape2.line.fill.background()

add_rounded_rect(slide, 9.4, 5.4, 3.5, 1.2, PHASE_GREEN,
                 'STRATEGIC: Super App\n"Activate What You Need"\n3\u20136 months', 13, WHITE)

add_text_box(slide, 0.35, 6.8, 12.6, 0.5,
             "Together, they move Cisco from fragmented products to a unified Security & Networking platform.",
             font_size=14, bold=True, color=CISCO_BLUE, alignment=PP_ALIGN.CENTER)


# ==========================================
# SLIDE 12: Stakeholder Endorsements
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
set_title(slide, "Stakeholder Endorsements — 20+ SMEs Consulted")

endorse_data = [
    ["Stakeholder", "Role", "Key Insight"],
    ["James Young", "Principal Global Security Specialist", '\u201cA disaster... practitioner-led dev is faster than siloed 5-yr projects\u201d'],
    ["Bryan Pluta", "Principal Solutions Engineer", '\u201cOne Cisco marketing is outpacing execution \u2014 customers calling us out\u201d'],
    ["Daniel Peluso", "ANZ Splunk Partner SE Manager", '\u201cCyberCX built own TAs for 7 yrs \u2014 we risk losing mindshare\u201d'],
    ["Matt Poland", "Sr. Director, SE", '\u201cTA sprawl hurts the sales force \u2014 need turnkey, outcome-focused experience\u201d'],
    ["Colin Gibbens", "Principal Engineer (CSC App)", 'Built 16 integrations in 1 yr via SoftServe; supports expanding'],
    ["Dimitri McKay", "Principal Security Architect", 'Validated DCE for resource bloat; Palo Alto already has AI parsing'],
    ["Sarav Radhakrishnan", "Distinguished Engineer", 'Validated DCE POC; supports Two Super App strategy (Security vs Net)'],
    ["Yaron Caspy", "Engineering PM", 'Built SSE App because one-size-fits-all failed (millions of missed events)'],
    ["Steven Moore", "Cisco Solutions Architect", '\u201cNetwork is a black box \u2014 infra is first point of blame\u201d (coined MTTI)'],
    ["Kunal Mukerjee", "SVP Engineering", '\u201cTA Flywheel\u201d vision; Master Library + DCE = CDF Data Products'],
]
add_table(slide, 0.2, 1.1, 12.9, endorse_data, col_widths=[2.0, 3.5, 7.0],
          header_color=CISCO_BLUE, font_size=10)

add_text_box(slide, 0.35, 5.7, 12.6, 0.5,
             "Feedback unanimously supports the SCAN + Dynamic Super App strategy.",
             font_size=16, bold=True, color=SPLUNK_GREEN, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 0.35, 6.2, 12.6, 0.5,
             'SCAN is the aspirin. The Super App is the cure.',
             font_size=22, bold=True, color=CISCO_BLUE, alignment=PP_ALIGN.CENTER)


# ==========================================
# SLIDE 13: Thank You
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_THANKYOU])
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "Thank You"

add_text_box(slide, 1.5, 4.2, 10, 0.5, "Questions & Discussion",
             font_size=24, color=CISCO_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 5.0, 10, 1.0,
             "Alan Ivarson  \u2022  Group Vice President, Field Solutions\n"
             "Amir (AK) Khamis  \u2022  Principal Architect",
             font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 6.2, 10, 0.4,
             "February 2026",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ─── Save ───
prs.save(OUTPUT_PATH)
import os
size_kb = os.path.getsize(OUTPUT_PATH) // 1024
print(f"\nSaved: {OUTPUT_PATH}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {size_kb:,} KB")
print("Using Cisco PowerPoint Template (LIGHT)")
