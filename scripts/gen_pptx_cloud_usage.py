#!/usr/bin/env python3
"""
Generate single-slide PPTX: Cisco Apps Still Running in Splunk Cloud.
Shows that 89 of 121 Cisco apps are still deployed across 2,973 Cloud stacks.
"""

import shutil
import zipfile
import os
import re
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Paths ───
_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
_REPO_ROOT  = os.path.join(_SCRIPT_DIR, '..')
TEMPLATE_POTX = os.path.join(_REPO_ROOT, 'docs', 'Cisco_PowerPoint_Template_LIGHT.potx')
TEMPLATE_PPTX = os.path.join(_REPO_ROOT, 'docs', 'cisco_template_converted_cloud.pptx')
OUTPUT_PATH   = os.path.join(_REPO_ROOT, 'docs', 'Splunkbase_Cloud_Usage.pptx')

# ─── Convert .potx → .pptx ───
def convert_potx_to_pptx(potx_path, pptx_path):
    shutil.copy2(potx_path, pptx_path)
    tmpdir = tempfile.mkdtemp()
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        zin.extractall(tmpdir)
    ct_path = os.path.join(tmpdir, '[Content_Types].xml')
    with open(ct_path, 'r', encoding='utf-8') as f:
        content = f.read()
    content = content.replace(
        'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
    )
    with open(ct_path, 'w', encoding='utf-8') as f:
        f.write(content)
    with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for root, dirs, files in os.walk(tmpdir):
            for fname in files:
                full_path = os.path.join(root, fname)
                arcname = os.path.relpath(full_path, tmpdir)
                zout.write(full_path, arcname)
    shutil.rmtree(tmpdir)

# ─── Load template or blank widescreen ───
use_template = False
if os.path.exists(TEMPLATE_POTX):
    convert_potx_to_pptx(TEMPLATE_POTX, TEMPLATE_PPTX)
    prs = Presentation(TEMPLATE_PPTX)
    use_template = True
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId is None:
            for k, v in prs.slides._sldIdLst[0].attrib.items():
                if k.endswith('}id') or k == 'r:id':
                    rId = v
                    break
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    print("Using Cisco PowerPoint Template (LIGHT)")
else:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    print("Cisco template not found — using blank widescreen presentation")

# ─── Colors ───
CISCO_BLUE     = RGBColor(0x04, 0x9F, 0xD9)
CISCO_BLUE_DK  = RGBColor(0x00, 0x7C, 0x92)
SPLUNK_GREEN   = RGBColor(0x65, 0xA6, 0x37)
DEBT_RED       = RGBColor(0xD3, 0x2F, 0x2F)
WARN_AMBER     = RGBColor(0xF9, 0xA8, 0x25)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY     = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY       = RGBColor(0x9E, 0x9E, 0x9E)
DARK_TEXT       = RGBColor(0x21, 0x21, 0x21)
PURPLE         = RGBColor(0x4A, 0x14, 0x8C)
PHASE_GREEN    = RGBColor(0x2E, 0x7D, 0x32)
PHASE_BLUE     = RGBColor(0x15, 0x65, 0xC0)
ARCHIVE_ORANGE = RGBColor(0xE6, 0x51, 0x00)
TEAL           = RGBColor(0x00, 0x79, 0x6B)
BAR_BLUE_1     = RGBColor(0x01, 0x57, 0x9B)  # darkest tier
BAR_BLUE_2     = RGBColor(0x02, 0x88, 0xD1)
BAR_BLUE_3     = RGBColor(0x03, 0xA9, 0xF4)
BAR_BLUE_4     = RGBColor(0x4F, 0xC3, 0xF7)  # lightest tier

# ─── Data: All 89 Cisco apps deployed in Splunk Cloud ───
# (app_package_name, versions, stacks)
CLOUD_DATA = [
    ("Splunk_TA_cisco-asa", 10, 1427),
    ("TA-cisco_ios", 15, 1295),
    ("Splunk_TA_cisco-ise", 6, 917),
    ("Splunk_TA_cisco_meraki", 13, 833),
    ("cisco_ios", 13, 741),
    ("CiscoSecurityCloud", 25, 622),
    ("Splunk_TA_cisco-esa", 7, 463),
    ("Splunk_CiscoISE", 7, 459),
    ("duo_splunkapp", 11, 428),
    ("TA-eStreamer", 20, 334),
    ("Splunk_TA_Talos_Intelligence", 2, 304),
    ("Splunk_TA_sourcefire", 3, 276),
    ("TA-cisco-cloud-security-addon", 9, 268),
    ("TA-meraki", 4, 256),
    ("TA-cisco-cloud-security-umbrella-addon", 7, 236),
    ("cisco-cloud-security", 13, 229),
    ("firepower_dashboard", 6, 207),
    ("TA-cisco_acs", 5, 203),
    ("Splunk_TA_cisco-wsa", 8, 186),
    ("TA_cisco_catalyst", 7, 161),
    ("TA-cisco_umbrella", 8, 154),
    ("ta_cisco_thousandeyes", 12, 149),
    ("Splunk_TA_cisco-ucs", 8, 107),
    ("cisco-catalyst-app", 3, 106),
    ("Splunk_TA_CCX_Unified_Cisco_Firepower_eStreamer", 3, 89),
    ("Splunk_TA_AppDynamics", 13, 78),
    ("Splunk_CiscoDNACenter", 2, 68),
    ("cisco_dc_networking_app_for_splunk", 3, 56),
    ("TA-cisco-amp4e", 5, 52),
    ("Splunk_CiscoSecuritySuite", 6, 50),
    ("cisco_cdr", 24, 50),
    ("TA_cisco-ACI", 6, 47),
    ("cisco-app-ACI", 6, 45),
    ("Alef_TA_MIMEDecoder", 6, 40),
    ("amp4e_events_input", 4, 40),
    ("Splunk_TA_CCX_DUO_Connector_Support", 3, 38),
    ("TA_cisco_cdr", 8, 36),
    ("cisco-sdwan-app", 4, 36),
    ("ta-cisco-sdwan", 3, 36),
    ("Splunk_TA_CCX_Cisco_ISE", 4, 35),
    ("Splunk_TA_Cisco_Intersight", 4, 35),
    ("Splunk_TA_CCX_Cisco_Meraki", 2, 31),
    ("Splunk-TA-cisco-dnacenter", 3, 30),
    ("Splunk-cisco-catalyst-center", 1, 30),
    ("splunk_app_stream_ipfix_cisco_hsl", 2, 29),
    ("Splunk-TA-cisco-catalyst-center", 1, 27),
    ("TA-cisco-stealthwatch", 1, 26),
    ("Splunk_TA_CCX_Cisco_Secure_Endpoint", 3, 24),
    ("TA-cisco-firepower-threat-defense-ftd-sourcetype", 3, 24),
    ("CiscoNVM", 6, 22),
    ("TA_cisco-Nexus-9k", 2, 22),
    ("cisco_cybervision_app_for_splunk", 4, 18),
    ("TA-Cisco-NVM", 3, 15),
    ("cisco_webex_meetings_app_for_splunk", 2, 15),
    ("Splunk_TA_CCX_Cisco_Secure_Network_Analytics", 1, 14),
    ("WebexTeams_AlertAction", 1, 14),
    ("ta_cisco_webex_add_on_for_splunk", 8, 14),
    ("TA-cisco_meraki_operations", 1, 13),
    ("cisco_sna_app", 2, 13),
    ("ta-cisco-webex-meetings-add-on-for-splunk", 2, 13),
    ("TA-stealthwatch_dataexporter", 1, 11),
    ("TA-cisco_cybervision", 4, 10),
    ("Splunk_TA_cisco-estreamer_Enosys", 1, 9),
    ("webex_alert", 1, 9),
    ("TA-cisco-threat-response", 3, 7),
    ("firegen_cisco_asa", 1, 7),
    ("Splunk_TA_cisco-ips", 2, 6),
    ("TA-add-on-for-cisco-prime-infrastructure", 3, 6),
    ("TA-cisco-firepower-threat-defense-ftd-dashboards", 1, 6),
    ("eStreamer-Dashboard", 3, 6),
    ("ta_cisco_spaces", 4, 6),
    ("TA_Meraki_SNMP_trap", 1, 5),
    ("opendns_investigate", 2, 5),
    ("tr-splunk-relay", 1, 5),
    ("SA_cisco_cdr_axl", 2, 4),
    ("TA-cisco-etd-connector", 1, 4),
    ("cisco-app-Nexus-9k", 1, 3),
    ("eStreamer", 1, 3),
    ("Cisco-suite-for-splunk", 1, 2),
    ("Cisco_Bug_Search_and_Analytics", 2, 2),
    ("Cisco_WSA_Insight", 1, 2),
    ("TA-cisco-firepower-pcap-add-on", 1, 2),
    ("TA-cisco-threat-grid", 1, 2),
    ("TA_cisco-NI", 1, 2),
    ("cisco_thousandeyes_alerting_app_for_splunk", 1, 2),
    ("TA-intersight-addon", 1, 1),
    ("TA_cisco-candid", 1, 1),
    ("cisco-app-NI", 1, 1),
    ("splunk_ta_cisco-ise", 1, 1),
]

# Friendly display names for the top 20
DISPLAY_NAMES = {
    "Splunk_TA_cisco-asa": "Cisco ASA (TA)",
    "TA-cisco_ios": "Cisco IOS (TA)",
    "Splunk_TA_cisco-ise": "Cisco ISE (TA)",
    "Splunk_TA_cisco_meraki": "Cisco Meraki (TA)",
    "cisco_ios": "Cisco IOS (App)",
    "CiscoSecurityCloud": "Security Cloud (SCAN)",
    "Splunk_TA_cisco-esa": "Cisco ESA (TA)",
    "Splunk_CiscoISE": "Cisco ISE (App)",
    "duo_splunkapp": "Cisco Duo (App)",
    "TA-eStreamer": "Firepower eStreamer (TA)",
    "Splunk_TA_Talos_Intelligence": "Talos Intelligence (TA)",
    "Splunk_TA_sourcefire": "Sourcefire (TA)",
    "TA-cisco-cloud-security-addon": "Umbrella Cloud Sec (TA)",
    "TA-meraki": "Meraki (Community TA)",
    "TA-cisco-cloud-security-umbrella-addon": "Umbrella DNS (TA)",
    "cisco-cloud-security": "Umbrella Cloud Sec (App)",
    "firepower_dashboard": "Firepower (Dashboard)",
    "TA-cisco_acs": "Cisco ACS (TA)",
    "Splunk_TA_cisco-wsa": "Cisco WSA (TA)",
    "TA_cisco_catalyst": "Cisco Catalyst (TA)",
}

# ─── Compute stats ───
TOTAL_APPS = 89
TOTAL_STACKS = 2973
TOTAL_SPLUNKBASE = 121
ARCHIVED_COUNT = 62
REMAINING_COUNT = 59

# Tier breakdown
tier1 = [(n, v, s) for n, v, s in CLOUD_DATA if s >= 500]
tier2 = [(n, v, s) for n, v, s in CLOUD_DATA if 100 <= s < 500]
tier3 = [(n, v, s) for n, v, s in CLOUD_DATA if 10 <= s < 100]
tier4 = [(n, v, s) for n, v, s in CLOUD_DATA if s < 10]

# Archived apps still running = total deployed - apps still on Splunkbase
# 89 deployed, 59 still on Splunkbase → at least 30 are from the archived set
ARCHIVED_STILL_RUNNING = TOTAL_APPS - REMAINING_COUNT

print(f"\n── Cloud Usage Stats ──")
print(f"Total apps deployed: {TOTAL_APPS}")
print(f"Total Cloud stacks: {TOTAL_STACKS:,}")
print(f"Pct of 121: {TOTAL_APPS/TOTAL_SPLUNKBASE*100:.0f}%")
print(f"Archived still running: ~{ARCHIVED_STILL_RUNNING}")
print(f"Tier 1 (500+): {len(tier1)} apps")
print(f"Tier 2 (100-499): {len(tier2)} apps")
print(f"Tier 3 (10-99): {len(tier3)} apps")
print(f"Tier 4 (1-9): {len(tier4)} apps")

# ─── Helpers ───

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    if font_name:
        p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text_box(slide, left, top, width, height, runs_list, alignment=PP_ALIGN.LEFT):
    """runs_list = list of paragraphs, each paragraph = list of (text, size, bold, color)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for pi, para_runs in enumerate(runs_list):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = alignment
        for text, font_size, bold, color in para_runs:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
    return shape


def add_stat_callout(slide, left, top, number, label, color=DEBT_RED, num_size=36, label_size=11):
    add_text_box(slide, left, top, 2.2, 0.65, number, font_size=num_size, bold=True,
                 color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + 0.55, 2.2, 0.4, label, font_size=label_size,
                 color=MED_GRAY, alignment=PP_ALIGN.CENTER)


def set_title(slide, text, color=DARK_TEXT):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = text
            for paragraph in ph.text_frame.paragraphs:
                paragraph.font.color.rgb = color
            return ph
    return None


# ================================================================
# SLIDE: Cisco Apps Still Running in Splunk Cloud
# ================================================================
if use_template:
    slide = prs.slides.add_slide(prs.slide_layouts[9])
    set_title(slide, "Still Running — Cisco Apps in Splunk Cloud", CISCO_BLUE)
else:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, 0.3, 0.15, 12.7, 0.6,
                 "Still Running — Cisco Apps in Splunk Cloud",
                 font_size=26, bold=True, color=CISCO_BLUE)

# Subtitle
add_text_box(slide, 0.3, 0.7, 12.7, 0.3,
             "Despite the Splunkbase mass archiving, 73% of Cisco apps remain actively deployed across Splunk Cloud stacks",
             font_size=12, color=MED_GRAY)

# ── Hero stat callouts ──
stat_y = 1.05
add_stat_callout(slide, 0.3,  stat_y, "89",     "Cisco Apps\nDeployed", CISCO_BLUE)
add_stat_callout(slide, 2.85, stat_y, "2,973",  "Cloud Stacks\nwith Cisco Apps", PHASE_GREEN)
add_stat_callout(slide, 5.4,  stat_y, "73%",    "of 121 Apps\nStill Running", WARN_AMBER)
add_stat_callout(slide, 7.95, stat_y, f"~{ARCHIVED_STILL_RUNNING}",   "Archived Apps\nStill Deployed", DEBT_RED)

# Total versions across all deployed apps
total_versions = sum(v for _, v, _ in CLOUD_DATA)
add_stat_callout(slide, 10.5, stat_y, f"{total_versions}",  "Total Versions\nin Production", PURPLE)

# Arrow connectors
for ax in [2.45, 5.0, 7.55, 10.1]:
    add_text_box(slide, ax, stat_y + 0.1, 0.45, 0.45, "•", font_size=20,
                 bold=True, color=RGBColor(0xDD, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)

# ── Left side: Top 20 Bar Chart ──
chart_title_y = 2.05
add_text_box(slide, 0.3, chart_title_y, 4.5, 0.3,
             "TOP 20 MOST DEPLOYED APPS", font_size=12, bold=True, color=DARK_TEXT)
add_text_box(slide, 5.0, chart_title_y, 3.5, 0.3,
             "(by Cloud stack count)", font_size=10, color=MED_GRAY)

# Bar chart area
bar_x_label = 0.3       # label start
bar_x_start = 3.3       # bar start
bar_max_w   = 4.0        # max bar width (for 1,427)
bar_x_count = 7.5        # count label position
bar_y_start = 2.4
bar_h       = 0.185
bar_gap     = 0.04
max_stacks  = CLOUD_DATA[0][2]  # 1,427

top20 = CLOUD_DATA[:20]
for i, (name, versions, stacks) in enumerate(top20):
    y = bar_y_start + i * (bar_h + bar_gap)
    display = DISPLAY_NAMES.get(name, name)

    # Color by tier
    if stacks >= 500:
        bar_color = BAR_BLUE_1
    elif stacks >= 200:
        bar_color = BAR_BLUE_2
    elif stacks >= 100:
        bar_color = BAR_BLUE_3
    else:
        bar_color = BAR_BLUE_4

    # Label (right-aligned)
    add_text_box(slide, bar_x_label, y - 0.02, 2.9, bar_h + 0.04,
                 display, font_size=8, color=DARK_TEXT, alignment=PP_ALIGN.RIGHT)

    # Bar
    bar_w = max(0.08, (stacks / max_stacks) * bar_max_w)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(bar_x_start), Inches(y), Inches(bar_w), Inches(bar_h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bar_color
    shape.line.fill.background()

    # Stack count
    add_text_box(slide, bar_x_start + bar_w + 0.05, y - 0.02, 0.8, bar_h + 0.04,
                 f"{stacks:,}", font_size=8, bold=True, color=bar_color)

# ── Remaining 69 note ──
remaining_y = bar_y_start + 20 * (bar_h + bar_gap) + 0.05
add_text_box(slide, 0.3, remaining_y, 7.5, 0.25,
             f"+ 69 more apps deployed across 1–89 stacks each",
             font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ── Right side: Deployment Tiers ──
tier_x = 8.6
tier_w = 4.4
tier_y_start = 2.05

add_text_box(slide, tier_x, tier_y_start, tier_w, 0.3,
             "DEPLOYMENT TIERS", font_size=12, bold=True, color=DARK_TEXT)

tier_data = [
    (f"TIER 1 — 500+ Stacks", f"{len(tier1)} apps", BAR_BLUE_1,
     "ASA, IOS, ISE, Meraki, Security Cloud, ESA"),
    (f"TIER 2 — 100–499 Stacks", f"{len(tier2)} apps", BAR_BLUE_2,
     "Duo, eStreamer, Talos, Umbrella, Firepower, WSA, ThousandEyes & more"),
    (f"TIER 3 — 10–99 Stacks", f"{len(tier3)} apps", BAR_BLUE_3,
     "AppDynamics, DNA Center, SD-WAN, ACI, Intersight, NVM, Webex & more"),
    (f"TIER 4 — 1–9 Stacks", f"{len(tier4)} apps", BAR_BLUE_4,
     "Prime Infra, Spaces, OpenDNS, Threat Grid, Bug Search & more"),
]

tier_box_h = 0.88
tier_box_gap = 0.1
for i, (header, count, color, examples) in enumerate(tier_data):
    y = tier_y_start + 0.35 + i * (tier_box_h + tier_box_gap)

    # Background box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(tier_x), Inches(y), Inches(tier_w), Inches(tier_box_h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(0.75)

    # Color accent bar on left
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(tier_x), Inches(y), Inches(0.06), Inches(tier_box_h)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()

    # Header + count
    add_text_box(slide, tier_x + 0.15, y + 0.04, 2.8, 0.25,
                 header, font_size=10, bold=True, color=color)
    add_text_box(slide, tier_x + 3.0, y + 0.04, 1.2, 0.25,
                 count, font_size=10, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.RIGHT)

    # Examples
    add_text_box(slide, tier_x + 0.15, y + 0.33, tier_w - 0.3, 0.5,
                 examples, font_size=8, color=MED_GRAY)

# ── Key Takeaway ──
takeaway_y = 6.05
add_rounded_rect(slide, 0.2, takeaway_y, 12.9, 0.3, PURPLE,
                 "KEY TAKEAWAY", font_size=10, font_color=WHITE)

takeaway_lines = [
    [("73% of all Cisco Splunkbase apps ", 11, True, CISCO_BLUE),
     ("are still actively deployed in Splunk Cloud — ", 11, False, DARK_TEXT),
     ("the archiving didn't stop usage, it stopped updates", 11, True, DEBT_RED)],
    [(f"~{ARCHIVED_STILL_RUNNING} archived apps ", 11, True, ARCHIVE_ORANGE),
     ("are still running on production Cloud stacks with ", 11, False, DARK_TEXT),
     ("no path to patches or upgrades", 11, True, DEBT_RED),
     (" — ", 11, False, DARK_TEXT),
     (f"2,973 stacks ({TOTAL_STACKS:,})", 11, True, PHASE_GREEN),
     (" depend on Cisco integrations", 11, False, DARK_TEXT)],
]
add_rich_text_box(slide, 0.5, takeaway_y + 0.32, 12.4, 0.7, takeaway_lines)

# ── Footer ──
add_rounded_rect(slide, 0.2, 7.05, 12.9, 0.32, RGBColor(0x37, 0x47, 0x4F),
                 "", font_size=10, font_color=WHITE)
footer_lines = [
    [("Splunk Cloud data only", 9, True, WARN_AMBER),
     (" — Enterprise deployment numbers are not included. ", 9, False, RGBColor(0xB0, 0xBE, 0xC5)),
     ("Actual Cisco footprint across all Splunk deployments is likely significantly higher.", 9, True, WHITE)],
]
add_rich_text_box(slide, 0.5, 7.07, 12.4, 0.28, footer_lines, alignment=PP_ALIGN.CENTER)

# ─── Save ───
prs.save(OUTPUT_PATH)
size_kb = os.path.getsize(OUTPUT_PATH) // 1024
print(f"\nSaved: {OUTPUT_PATH}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {size_kb:,} KB")
