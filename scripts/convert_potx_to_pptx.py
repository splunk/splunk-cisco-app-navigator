#!/usr/bin/env python3
"""
Convert .potx template to .pptx by changing internal content type,
then inspect slide layouts.
"""
import shutil
import zipfile
import os
import tempfile

TEMPLATE = '/Users/akhamis/repo/cisco_control_center_app/docs/Cisco_PowerPoint_Template_LIGHT.potx'
OUT_PPTX = '/tmp/cisco_template_converted.pptx'

# A .potx is a .pptx zip with a different content type in [Content_Types].xml
# We need to change the content type from template to presentation
shutil.copy2(TEMPLATE, OUT_PPTX)

# Modify the content type inside the zip
tmpdir = tempfile.mkdtemp()
with zipfile.ZipFile(OUT_PPTX, 'r') as zin:
    zin.extractall(tmpdir)

# Fix content types
ct_path = os.path.join(tmpdir, '[Content_Types].xml')
with open(ct_path, 'r', encoding='utf-8') as f:
    content = f.read()

# Replace template content type with presentation content type
content = content.replace(
    'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
)

with open(ct_path, 'w', encoding='utf-8') as f:
    f.write(content)

# Repackage
with zipfile.ZipFile(OUT_PPTX, 'w', zipfile.ZIP_DEFLATED) as zout:
    for root, dirs, files in os.walk(tmpdir):
        for fname in files:
            full_path = os.path.join(root, fname)
            arcname = os.path.relpath(full_path, tmpdir)
            zout.write(full_path, arcname)

shutil.rmtree(tmpdir)
print(f"Converted template saved to: {OUT_PPTX}")
print(f"Size: {os.path.getsize(OUT_PPTX):,} bytes")

# Now inspect it
from pptx import Presentation
prs = Presentation(OUT_PPTX)

print(f"\nSlide width: {prs.slide_width / 914400:.2f} inches")
print(f"Slide height: {prs.slide_height / 914400:.2f} inches")
print(f"Number of slide layouts: {len(prs.slide_layouts)}")
print()

for i, layout in enumerate(prs.slide_layouts):
    print(f"{'='*60}")
    print(f"Layout [{i}]: '{layout.name}'")
    print(f"  Placeholders: {len(layout.placeholders)}")
    for ph in layout.placeholders:
        print(f"    idx={ph.placeholder_format.idx}, type={ph.placeholder_format.type}, "
              f"name='{ph.name}', "
              f"pos=({ph.left/914400:.2f}\", {ph.top/914400:.2f}\"), "
              f"size=({ph.width/914400:.2f}\" x {ph.height/914400:.2f}\")")
    shape_count = len(layout.shapes) - len(layout.placeholders)
    if shape_count > 0:
        print(f"  Additional shapes: {shape_count}")
    print()
