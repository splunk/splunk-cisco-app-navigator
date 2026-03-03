#!/usr/bin/env python3
"""
Generate a PowerPoint presentation from the CiscoSecurityCloud Audit Report.
Uses the Cisco brand teal (#049fd9) as accent color.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- Brand colors --
CISCO_TEAL = RGBColor(0x04, 0x9F, 0xD9)
CISCO_DARK = RGBColor(0x1B, 0x1C, 0x1D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
RED = RGBColor(0xCC, 0x00, 0x00)
AMBER = RGBColor(0xE6, 0x8A, 0x00)
GREEN = RGBColor(0x00, 0x99, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ─── Helpers ────────────────────────────────────────────────────────

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=CISCO_DARK, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, bullets, left, top, width, height,
                     font_size=16, color=CISCO_DARK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf


def section_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, CISCO_TEAL)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
                 title, font_size=40, bold=True, color=WHITE,
                 alignment=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
                     subtitle, font_size=22, color=WHITE, alignment=PP_ALIGN.LEFT)
    # bottom bar
    add_rect(slide, Inches(0), Inches(7.0), W, Inches(0.5), CISCO_DARK)
    add_text_box(slide, Inches(1), Inches(7.05), Inches(5), Inches(0.4),
                 "CiscoSecurityCloud v3.6.2 — Audit Report",
                 font_size=11, color=WHITE)
    return slide


def content_slide(title, bullets=None, body_text=None, two_col=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, WHITE)
    # header bar
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.0), CISCO_TEAL)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.7),
                 title, font_size=28, bold=True, color=WHITE)
    # footer
    add_rect(slide, Inches(0), Inches(7.0), W, Inches(0.5), LIGHT_GRAY)
    add_text_box(slide, Inches(0.6), Inches(7.05), Inches(5), Inches(0.4),
                 "CiscoSecurityCloud v3.6.2 — Audit Report",
                 font_size=10, color=MED_GRAY)

    if bullets:
        add_bullet_slide(slide, bullets,
                         Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5),
                         font_size=17)
    if body_text:
        add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5),
                     body_text, font_size=16)
    if two_col:
        left_bullets, right_bullets = two_col
        add_bullet_slide(slide, left_bullets,
                         Inches(0.6), Inches(1.3), Inches(5.5), Inches(5.5),
                         font_size=15)
        add_bullet_slide(slide, right_bullets,
                         Inches(6.6), Inches(1.3), Inches(5.5), Inches(5.5),
                         font_size=15)
    return slide


def kpi_slide(title, kpis):
    """kpis: list of (label, value, color)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.0), CISCO_TEAL)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.7),
                 title, font_size=28, bold=True, color=WHITE)
    add_rect(slide, Inches(0), Inches(7.0), W, Inches(0.5), LIGHT_GRAY)
    add_text_box(slide, Inches(0.6), Inches(7.05), Inches(5), Inches(0.4),
                 "CiscoSecurityCloud v3.6.2 — Audit Report",
                 font_size=10, color=MED_GRAY)

    cols = len(kpis)
    card_w = Inches(2.4)
    gap = Inches(0.3)
    total_w = cols * card_w + (cols - 1) * gap
    start_x = int((W - total_w) / 2)

    for i, (label, value, clr) in enumerate(kpis):
        x = start_x + i * int(card_w + gap)
        y = Inches(2.0)
        card = add_rect(slide, x, y, card_w, Inches(3.0), LIGHT_GRAY, MED_GRAY)
        card.line.width = Pt(0.5)
        add_text_box(slide, x, Inches(2.4), card_w, Inches(1.5),
                     value, font_size=48, bold=True, color=clr,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(3.9), card_w, Inches(0.8),
                     label, font_size=14, color=CISCO_DARK,
                     alignment=PP_ALIGN.CENTER)
    return slide


def table_slide(title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.0), CISCO_TEAL)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.7),
                 title, font_size=28, bold=True, color=WHITE)
    add_rect(slide, Inches(0), Inches(7.0), W, Inches(0.5), LIGHT_GRAY)
    add_text_box(slide, Inches(0.6), Inches(7.05), Inches(5), Inches(0.4),
                 "CiscoSecurityCloud v3.6.2 — Audit Report",
                 font_size=10, color=MED_GRAY)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_w = Inches(12) if not col_widths else sum(col_widths)
    tbl_left = Inches(0.6)
    tbl_top = Inches(1.3)
    row_height = Inches(0.45)
    tbl_h = row_height * n_rows

    table = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_w, tbl_h).table

    if col_widths:
        for ci, cw in enumerate(col_widths):
            table.columns[ci].width = cw

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = CISCO_DARK

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = CISCO_DARK
                p.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT_GRAY

    return slide


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CISCO_DARK)
add_rect(slide, Inches(0), Inches(2.0), W, Inches(0.06), CISCO_TEAL)
add_text_box(slide, Inches(1), Inches(2.4), Inches(11), Inches(1.5),
             "CiscoSecurityCloud App (v3.6.2)",
             font_size=44, bold=True, color=WHITE)
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.7),
             "Comprehensive Audit Report",
             font_size=30, color=CISCO_TEAL)
add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(1),
             "February 16, 2026\nCisco Enterprise Networking — Splunk Integration Team",
             font_size=16, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — Executive Summary / Agenda
# ═══════════════════════════════════════════════════════════════════
content_slide("Executive Summary", bullets=[
    "Audit of the CiscoSecurityCloud Splunk app covering 16 Cisco security products",
    "63 sourcetype stanzas, 506 field aliases, 221 EVAL expressions, 250 transforms",
    "Built on Splunk UCC framework with modular inputs, React-based UI, CIM normalization",
    "",
    "Key findings organized by impact:",
    "  \u2022  Critical — CIM field mapping errors affecting every ES customer (dest = SrcIP)",
    "  \u2022  Critical — FTD/ASA sourcetype routing misclassifies events; confirmed by field engineers",
    "  \u2022  Critical — 23 sourcetypes with redundant EXTRACT-on-JSON anti-pattern",
    "  \u2022  High — 0% compliance on ANNOTATE_PUNCT / LEARN_SOURCETYPE across all 63 stanzas",
    "  \u2022  High — ~35 JSON sourcetypes default to KV_MODE = auto (wasted regex extraction)",
    "",
    "The 3 most impactful fixes: (1) correct dest in FTD syslog, (2) fix FTD/ASA routing, (3) remove EXTRACT from JSON sourcetypes"
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — App at a Glance (KPIs)
# ═══════════════════════════════════════════════════════════════════
kpi_slide("App at a Glance", [
    ("Cisco Products", "16", CISCO_TEAL),
    ("Sourcetype Stanzas", "63", CISCO_TEAL),
    ("Field Aliases", "506", CISCO_TEAL),
    ("EVAL Extractions", "221", CISCO_TEAL),
    ("Transforms", "250", CISCO_TEAL),
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — Product-to-Sourcetype Map (page 1)
# ═══════════════════════════════════════════════════════════════════
table_slide("Product-to-Sourcetype Map (1/2)",
    ["#", "Product", "Key Sourcetypes", "Count", "Collection"],
    [
        ["1", "AI Defense", "cisco:ai:defense, cisco:ai:defense:notable", "2", "HEC"],
        ["2", "Duo", "cisco:duo:authentication, :authentication_v2", "2", "REST API"],
        ["3", "ETD", "cisco:etd", "1", "REST API"],
        ["4", "EVM (Beta)", "cisco:evm:process_activity + 4 more", "5", "Agent/JSON"],
        ["5", "CII", "cisco:cii", "1", "REST + S3"],
        ["6", "Multicloud Defense", "gateway + 18 child sourcetypes", "19", "HEC"],
        ["7", "NVM", "cisco:nvm:flowdata:v2 + 5 more", "6", "Agent/JSON"],
        ["8", "Secure Endpoint", "cisco:se", "1", "REST API"],
    ],
    col_widths=[Inches(0.5), Inches(2.8), Inches(5.5), Inches(0.9), Inches(1.8)]
)

table_slide("Product-to-Sourcetype Map (2/2)",
    ["#", "Product", "Key Sourcetypes", "Count", "Collection"],
    [
        ["9", "Secure Firewall (FTD/ASA)", "cisco:ftd:syslog, cisco:sfw:estreamer, cisco:asa", "3", "Syslog+REST"],
        ["10", "Secure Workload", "cisco:secure:workload", "1", "Syslog"],
        ["11", "Isovalent", "cisco:isovalent + 5 child sourcetypes", "6", "HEC"],
        ["12", "Isovalent Edge (Alpha)", "cisco:isovalent:edge_processor + 5 children", "6", "HEC"],
        ["13", "SMA", "cisco:sma:submissions", "1", "REST API"],
        ["14", "SNA", "cisco:sna, :notable, :risk", "3", "REST API"],
        ["15", "Vulnerability Intel", "cisco:cvi", "1", "REST API"],
        ["16", "XDR", "cisco:xdr:incidents, :incidents-summary, :notable", "3", "REST API"],
    ],
    col_widths=[Inches(0.5), Inches(2.8), Inches(5.5), Inches(0.9), Inches(1.8)]
)

# ═══════════════════════════════════════════════════════════════════
# SECTION — Critical Findings
# ═══════════════════════════════════════════════════════════════════
section_slide("Critical Findings",
              "Search-Time Anti-Patterns \u2022 CIM Mapping Errors \u2022 FTD/ASA Routing")

# ═══════════════════════════════════════════════════════════════════
# SLIDE — Search-Time EXTRACT Anti-Patterns
# ═══════════════════════════════════════════════════════════════════
content_slide("Finding #1: EXTRACT-on-JSON Anti-Pattern", bullets=[
    "23 sourceted stanzas set KV_MODE = json AND define EXTRACT- regex patterns",
    "When KV_MODE = json, Splunk auto-extracts every JSON key — EXTRACT is redundant",
    "Regex patterns use .*? (lazy match) through nested JSON = excessive CPU backtracking",
    "",
    "Affected: All 18 cisco:multicloud:defense:* + cisco:ai:defense + cisco:isovalent",
    "",
    "2 EXTRACT patterns have broken syntax (silently fail):",
    '  \u2022  session_summary: \"type\":\"(SESSION_SUMMARY?P<type>[^\"]+)\" — wrong named group syntax',
    '  \u2022  auditlogs: \"clientVersion:\"... — missing opening quote',
    "",
    "Fix: Remove all EXTRACT- from JSON sourcetypes. Use FIELDALIAS to rename",
    "auto-extracted dotted paths (e.g., sessionSummaryInfo.ingressConnection.srcIp → src_ip)"
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE — CIM Field Mapping Errors
# ═══════════════════════════════════════════════════════════════════
content_slide("Finding #2: CIM Field Mapping Errors", bullets=[
    "dest mapped to Source IP in cisco:ftd:syslog (CRITICAL)",
    "  EVAL-dest = coalesce(SrcIP, dest_ip)  ← SrcIP is almost always populated",
    "  Every FTD customer has dest pointing to the WRONG IP in Enterprise Security",
    "  Confirmed by Paul Pelletier: \"src and dest aren't even getting extracted\"",
    "  Fix: EVAL-dest = coalesce(DstIP, DestIP, ResponderIP, dest_ip)",
    "",
    "dest mapped to Application Name in cisco:duo:authentication_v2",
    "  FIELDALIAS-dest = application.name AS dest",
    "  CIM dest should be target system (hostname/IP), not \"Slack\" or \"VPN\"",
    "",
    "Isovalent src/dest mapped to process names / pod names instead of IPs",
    "  Bhavin Patel: \"Can we map node_name to dest? dest should reflect the machine\""
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE — FTD/ASA Routing
# ═══════════════════════════════════════════════════════════════════
content_slide("Finding #3: FTD/ASA Sourcetype Routing Problem", bullets=[
    "FTD events with DeviceUUID and key=value fields get misrouted to cisco:asa",
    "Reported by Mikael Bjerkeland — confirmed across multiple customer environments",
    "He has distributed a customer override app as a workaround",
    "",
    "Root cause: FTD v7+ emits ASA-style syslogs (%ASA-7-302020 format)",
    "  \u2022  set_ftd_sourcetype regex only matches %FTD or specific message IDs",
    "  \u2022  set_asa_sourcetype regex is overly broad: %ASA matches everything",
    "  \u2022  FTD events fail the FTD check → match %ASA → routed to cisco:asa",
    "",
    "Impact: FTD events with rich JSON payloads land in ASA sourcetype",
    "  No field extractions match → data is effectively unusable",
    "  Affects every customer using FTD v7+ with Advanced Logging",
    "",
    "Fix: Add DeviceUUID detection to set_ftd_sourcetype; tighten set_asa_sourcetype"
])

# ═══════════════════════════════════════════════════════════════════
# SECTION — Community-Reported Issues
# ═══════════════════════════════════════════════════════════════════
section_slide("Community-Reported Issues",
              "Real-world feedback from Splunk practitioners & Cisco security engineers")

content_slide("Community-Reported Issues from the Field", bullets=[
    "ASA regex \\w does not match dashes in usernames (SBGSPLNKSA-1423)",
    "  Example: %ASA-config-7-111009: User 'admin-user' — extraction fails",
    "",
    "ASA message_id / log_level extraction broken for infixed formats",
    "  %ASA-config-7-111009, %ASA-session-6-302020 — fix merged by Sviatoslav",
    "",
    "EVENT_BREAKER_ENABLE has no effect on modular inputs (Bhavin Patel, Tomas Moser)",
    "  Only honored on Universal Forwarders for raw TCP/UDP — remove from all stanzas",
    "",
    "Isovalent parsing broke between v3.5.0 → v3.5.1 (KV_MODE changed to none)",
    "",
    "FTD v10 JSON-in-Syslog format needs robust two-phase parsing strategy",
    "",
    "DeviceUUID does not map to a friendly hostname — analysts can't identify devices",
    "",
    "Splunk Cloud: sc_admin cannot configure app (grayed out); syslog setup misleading"
])

# ═══════════════════════════════════════════════════════════════════
# SECTION — Magic Six
# ═══════════════════════════════════════════════════════════════════
section_slide("Magic Six Optimization",
              "Where it actually matters — and where it doesn't")

content_slide("Magic Six — Routing vs. Terminal Stanzas", bullets=[
    "Not every stanza needs all six settings (per Josh Wilson's feedback)",
    "",
    "ROUTING stanzas (Magic Six less critical) — events pass through quickly:",
    "  \u2022  [gateway] → 18 Multicloud Defense children",
    "  \u2022  [cisco:isovalent] → alert, processExec, processConnect, etc.",
    "  \u2022  [cisco:ftd:syslog] ↔ [cisco:asa] bidirectional routing",
    "",
    "TERMINAL stanzas (Magic Six critical) — data is parsed & indexed here:",
    "  \u2022  cisco:asa — Best configured (all 6 set). Still missing ANNOTATE_PUNCT",
    "  \u2022  cisco:ftd:syslog — Missing TIME_FORMAT (datetime.xml heuristic runs)",
    "  \u2022  38 JSON sourcetypes — Missing SHOULD_LINEMERGE",
    "  \u2022  55 sourcetypes — Missing TIME_FORMAT",
    "",
    "Universal gaps across ALL 63 stanzas:",
    "  \u2022  ANNOTATE_PUNCT = false — 0% compliance (every event generates punct::)",
    "  \u2022  LEARN_SOURCETYPE = false — 0% compliance (classification engine runs every event)"
])

content_slide("The Golden Template", bullets=[
    "Every new sourcetype stanza should start from this template:",
    "",
    "  [cisco:product:sourcetype]",
    "  SHOULD_LINEMERGE = false",
    "  LINE_BREAKER = ([\\r\\n]+)",
    "  TIME_FORMAT = %Y-%m-%dT%H:%M:%S.%3QZ",
    "  TIME_PREFIX = \"timestamp\"\\s*:\\s*\"",
    "  MAX_TIMESTAMP_LOOKAHEAD = 32",
    "  TRUNCATE = 999999",
    "  ANNOTATE_PUNCT = false",
    "  LEARN_SOURCETYPE = false",
    "  KV_MODE = json",
    "",
    "For syslog: adjust TIME_FORMAT/TIME_PREFIX, set KV_MODE = none",
    "Any deviation from this template requires documented justification"
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE — Beyond Magic Six
# ═══════════════════════════════════════════════════════════════════
content_slide("Beyond Magic Six — Additional Optimization", bullets=[
    "KV_MODE defaults to auto on ~35 JSON sourcetypes",
    "  Splunk runs regex-based key=value extraction at search time — pure waste for JSON",
    "  Fix: Set KV_MODE = json explicitly",
    "",
    "CHARSET defaults to ASCII on all but 1 sourcetype",
    "  API sourcetypes (Duo, XDR, ETD) can return non-ASCII — hex-escaped today",
    "  Fix: Set CHARSET = UTF-8 on all API-based sourcetypes",
    "",
    "TRUNCATE = 0 (unlimited) on 28 sourcetypes",
    "  Risk: unexpectedly large JSON blob → single event → memory pressure",
    "  Fix: Replace with 999999 for API JSON, 10000 for syslog",
    "",
    "fqdn_resoleve sourcetype typo — permanent data contract issue",
    "  Fix now before more customers deploy",
    "",
    "pulldown_type = 1 on 7 stanzas — unnecessary UI hint, remove it"
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE — CIM Normalization
# ═══════════════════════════════════════════════════════════════════
content_slide("CIM Normalization Completeness", bullets=[
    "What they got right:",
    "  \u2022  Strong tagging across 15+ CIM data models",
    "  \u2022  All 82 eventtypes use sourcetype= (index-time efficient)",
    "  \u2022  ASA and eStreamer have the most complete CIM mapping",
    "",
    "Key gaps:",
    "  \u2022  cisco:ftd:syslog — Missing bytes, bytes_in, bytes_out, protocol, vendor_product",
    "  \u2022  cisco:sma:submissions — Missing action, category, file_name, signature, severity",
    "  \u2022  cisco:duo:authentication (v1) — Missing app, src, dest, signature",
    "  \u2022  cisco:sna — Missing severity for Alerts data model",
    "",
    "CIM over-mapping to multiple data models (reported by Josh Wilson):",
    '  \"Do not attempt to only look at CIM from the angle of detection —',
    '   CIM is also used for dashboards, correlation, compliance, and ad-hoc analysis\"'
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE — Potential Issues That Could Break Environments
# ═══════════════════════════════════════════════════════════════════
table_slide("Issues That Could Break a Customer's Environment",
    ["Severity", "Issue", "Impact"],
    [
        ["CRITICAL", "dest = coalesce(SrcIP, dest_ip) in FTD syslog", "Every FTD customer has wrong dest in ES"],
        ["CRITICAL", "FTD/ASA routing transform too broad", "FTD events misrouted; data unusable"],
        ["CRITICAL", "23 stanzas with redundant EXTRACT-on-JSON", "Wasted CPU on every search"],
        ["CRITICAL", "2 broken EXTRACT regex patterns", "type and clientVersion never extracted"],
        ["CRITICAL", "fqdn_resoleve sourcetype typo", "Data in misspelled sourcetype"],
        ["HIGH", "0% ANNOTATE_PUNCT / LEARN_SOURCETYPE", "Every event pays unnecessary CPU tax"],
        ["HIGH", "KV_MODE = auto on ~35 JSON sourcetypes", "Wasted regex extraction per search"],
        ["HIGH", "TRUNCATE = 0 on 28 sourcetypes", "Memory pressure risk"],
        ["HIGH", "Missing TIME_FORMAT on 55 sourcetypes", "~60% CPU overhead per Magic Six research"],
        ["HIGH", "sc_admin cannot configure in Splunk Cloud", "Grayed out Configure button"],
    ],
    col_widths=[Inches(1.2), Inches(5.6), Inches(5.2)]
)

# ═══════════════════════════════════════════════════════════════════
# SLIDE — Knowledge Object Loading
# ═══════════════════════════════════════════════════════════════════
content_slide('Knowledge Object Loading — The "All Products Active" Problem', bullets=[
    "Every Splunk restart loads ALL knowledge objects for ALL 16 products",
    "  Even if the customer only uses one product",
    "",
    "A Duo-only customer loads:",
    "  506 field aliases (needs 14)  •  221 EVALs (needs 2)  •  250 transforms (needs ~3)",
    "  82 eventtypes (needs 2)  •  71 tag stanzas (needs 1)",
    "  = ~97% overhead of unused knowledge objects",
    "",
    "The 118 ASA cisco_asa_message_id_* transforms are the heaviest cost:",
    "  Splunk tests each REGEX sequentially until one matches — O(n) where n=118",
    "  At 50K+ ASA EPS, this is measurable CPU cost",
    "",
    "Solution: Per-product activation model (like the DCE pattern)",
    "  Customers load only the knowledge objects they need"
])

# ═══════════════════════════════════════════════════════════════════
# SECTION — Splunk Cloud SSAI
# ═══════════════════════════════════════════════════════════════════
section_slide("Splunk Cloud SSAI Deployment",
              "Anatomy of what was actually deployed to SHC and IDX tiers")

content_slide("SSAI Tier Partitioning — What We Found", bullets=[
    "Source app: 2,696 files  →  SHC package: 2,511 files  →  IDX package: 4 files",
    "",
    "Headline: SSAI partitions at FILE level, not SETTING level",
    "  props.conf on indexer = byte-for-byte identical to the original (126 KB)",
    "  Indexer loads 506 FIELDALIAS + 221 EVAL + 135 REPORT — none used at index time",
    "",
    "Indexer gets only 4 files: app.conf (438B) + props.conf (126KB) + transforms.conf (67KB) + local.meta",
    "  No bin/, lib/, lookups/, appserver/, eventtypes.conf, tags.conf",
    "",
    "SHC gets everything: 13 conf files, 487 Python scripts, ~2,000 lib files",
    "  botocore is 60% of the SHC package (~1,500 files for CII AWS S3 input)",
    "  Every SHC member loads these even if CII is not used",
    "",
    "SSAI overrides @placement for restmap.conf and server.conf (Cloud-specific intelligence)"
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE — App Design Assessment
# ═══════════════════════════════════════════════════════════════════
content_slide("App Design Assessment", bullets=[
    "Build Framework: Splunk UCC (splunktaucclib 8.1.0) — standard, recommended",
    "",
    "Architecture Strengths:",
    "  \u2022  Clean per-product Python modules under bin/CiscoSecurityCloud/",
    "  \u2022  Shared utilities (config, secrets, exceptions) avoid duplication",
    "  \u2022  Credentials properly encrypted via Splunk storage passwords",
    "",
    "Concerns:",
    "  \u2022  Naming inconsistency — mix of sbg_* and cisco_* input prefixes",
    '       \"SBG\" (Security Business Group) is an internal team name in customer-facing inputs',
    "  \u2022  Monolithic conf — all 16 products in single props.conf (2,340 lines)",
    "  \u2022  Version-pinned boto3 (1.28.49) — security patches are frequent",
    "  \u2022  is_configured = false in app.conf — banner may persist after setup"
])

# ═══════════════════════════════════════════════════════════════════
# SECTION — Recommendations
# ═══════════════════════════════════════════════════════════════════
section_slide("Recommendations", "Prioritized by impact and effort")

content_slide("Immediate Fixes (Days, Not Weeks)", bullets=[
    "\u2705  Fix dest mapping in cisco:ftd:syslog",
    "     EVAL-dest = coalesce(DstIP, DestIP, ResponderIP, dest_ip)",
    "",
    "\u2705  Add ANNOTATE_PUNCT = false and LEARN_SOURCETYPE = false to all 63 stanzas",
    "     Zero behavioral risk, immediate CPU savings",
    "",
    "\u2705  Remove all EXTRACT- patterns from 23 JSON sourcetypes",
    "     Replace with FIELDALIAS to auto-extracted dotted JSON paths",
    "",
    "\u2705  Fix 2 broken regex patterns in Multicloud Defense",
    "",
    "\u2705  Fix fqdn_resoleve typo in props.conf and transforms.conf",
    "",
    "\u2705  Set KV_MODE = json on ~35 JSON sourcetypes defaulting to auto"
])

content_slide("Short-Term Fixes (Weeks)", bullets=[
    "Fix FTD/ASA routing transform — add DeviceUUID detection",
    "",
    "Add SHOULD_LINEMERGE = false to ~38 terminal JSON sourcetypes",
    "",
    "Add explicit TIME_FORMAT to SNA, FTD, eStreamer",
    "",
    "Complete FTD syslog CIM mapping — bytes, protocol, vendor_product",
    "",
    "Add CHARSET = UTF-8 to all API-based sourcetypes",
    "",
    "Replace TRUNCATE = 0 with appropriate limits (999999 API, 10000 syslog)",
    "",
    "Fix Isovalent dest mapping — map to node_name instead of pod.name"
])

content_slide("Medium-Term Improvements (Months)", bullets=[
    "Per-product activation model (DCE pattern)",
    "  Duo-only customer loads 14 aliases instead of 506",
    "",
    "Audit CIM cross-model tagging",
    "  Each eventtype → single most appropriate CIM data model",
    "  Cross-model mapping = exception with documented justification",
    "",
    "Review Cisco_Security custom data model",
    "  Determine if it duplicates CIM acceleration (double cost)",
    "",
    "Update setup UI for Splunk Cloud",
    "  Remove direct TCP/UDP config; add SC4S/HF guidance",
    "  Address sc_admin permissions (SBGSPLNKSA-1446)",
    "",
    "Establish Magic Six as an engineering standard",
    "  Every new sourcetype starts from golden template",
    "  Deviations require documented justification"
])

# ═══════════════════════════════════════════════════════════════════
# SLIDE — Overall Assessment Summary
# ═══════════════════════════════════════════════════════════════════

GRADE_COLORS = {"A": GREEN, "A-": GREEN, "B+": GREEN, "B": GREEN, "B-": GREEN,
                "C+": AMBER, "C": AMBER, "C-": AMBER,
                "D+": RED, "D": RED, "D-": RED, "F": RED}

table_slide("Overall Assessment Summary",
    ["Area", "Rating", "Key Finding"],
    [
        ["Search-Time Extraction Quality", "D", "23 stanzas w/ redundant EXTRACT-on-JSON; 2 broken regex"],
        ["CIM Field Mapping Accuracy", "C-", "dest=SrcIP in FTD syslog; Isovalent dest = pod name"],
        ["FTD/ASA Routing", "D", "FTD events misrouted to cisco:asa; confirmed multi-customer"],
        ["Magic Six (where it matters)", "C", "ASA well-configured; 38 JSON stanzas missing SHOULD_LINEMERGE"],
        ["Universal Optimization", "F", "0% compliance: ANNOTATE_PUNCT & LEARN_SOURCETYPE"],
        ["CIM Normalization Completeness", "B-", "Strong ASA/eStreamer; significant FTD, SMA, Duo v1 gaps"],
        ["Knowledge Object Loading", "D", "506 aliases + 221 EVALs loaded for ALL 16 products"],
        ["Splunk Cloud Readiness", "C+", "SSAI partitions correctly; sc_admin perms; syslog UI misleading"],
        ["Architecture / Design", "B", "UCC framework; proper creds; naming inconsistency, monolithic conf"],
    ],
    col_widths=[Inches(3.5), Inches(0.9), Inches(7.6)]
)

# ═══════════════════════════════════════════════════════════════════
# SLIDE — Bottom Line
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CISCO_DARK)
add_rect(slide, Inches(0), Inches(2.0), W, Inches(0.06), CISCO_TEAL)
add_text_box(slide, Inches(1), Inches(2.4), Inches(11), Inches(0.8),
             "The Bottom Line",
             font_size=40, bold=True, color=WHITE)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(3.5),
             ("The app's ambition is right — a unified integration experience across "
              "16 Cisco security products. The Python backend is well-engineered.\n\n"
              "The Splunk-native aspects (props.conf, transforms.conf, CIM mapping) "
              "need significant attention.\n\n"
              "The 3 most impactful fixes:\n"
              "  1.  Correct dest in FTD syslog\n"
              "  2.  Fix the FTD/ASA routing transform\n"
              "  3.  Remove redundant EXTRACT patterns from JSON sourcetypes\n\n"
              "Community feedback from practitioners is clear — customers are building "
              "workarounds. Addressing these proactively will significantly improve "
              "the app's reputation and reduce support burden."),
             font_size=18, color=WHITE)


# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
output_path = "/Users/akhamis/repo/cisco_enterprise_networking_app/docs/CiscoSecurityCloud_Audit_Report.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
