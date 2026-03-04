#!/usr/bin/env python3
"""
Generate 13-slide PPTX: Cisco Splunkbase Coverage Story.
Part 1 of a two-part executive series.
Steve Jobs style — one idea per slide, clean, building narrative tension.

Title:     "The Path Forward" — Two-part series overview (Today + March 10)
Preview:   "The 5 Pillars of the Cure" — Teaser for March 10 demos
Agenda:    "Today's Agenda" — Part 1 roadmap (7 sections)
Slide 1:   "#1 Most Requested Demo" — $500B TAM, customer demand, stakes
Slide 2:   "Why We Are Here" — 4 forces: Renewal Risk, Data Integrity, Agentic Future, Demand
Slide 3:   "What Happened" — 121 → 62 archived → 59 remain
Slide 4:   "The Real Number" — Of those 59, only 18 are enterprise-grade
Slide 5:   "Still Running" — 89 apps, 2,973 Cloud stacks (evidence)
Slides 6a-6c: "The Coverage Gap" — Progressive build
Slide 7:   "Field Feedback" — 7 pain points from the field
Slide 8:   "Systemic Issues" — 4 root causes

Part 1: Today (problem)   |   Part 2: March 10 (solutions + 5 demos)
Each slide has speaker notes for delivery guidance.
"""

import shutil
import zipfile
import os
import re
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ─── Paths ───
_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
_REPO_ROOT  = os.path.join(_SCRIPT_DIR, '..')
TEMPLATE_POTX = os.path.join(_REPO_ROOT, 'docs', 'Cisco_PowerPoint_Template_LIGHT.potx')
TEMPLATE_PPTX = os.path.join(_REPO_ROOT, 'docs', 'cisco_template_converted_v2.pptx')
OUTPUT_PATH   = os.path.join(_REPO_ROOT, 'docs', 'Cisco_Splunkbase_Coverage.pptx')

# ─── Convert .potx → .pptx ───
def convert_potx_to_pptx(potx_path, pptx_path):
    shutil.copy2(potx_path, pptx_path)
    tmpdir = tempfile.mkdtemp()
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        zin.extractall(tmpdir)
    ct_path = os.path.join(tmpdir, '[Content_Types].xml')
    with open(ct_path, 'r', encoding='utf-8') as f:
        content = f.read()
    content = content.replace(
        'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
    )
    with open(ct_path, 'w', encoding='utf-8') as f:
        f.write(content)
    with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for root, dirs, files in os.walk(tmpdir):
            for fname in files:
                full_path = os.path.join(root, fname)
                arcname = os.path.relpath(full_path, tmpdir)
                zout.write(full_path, arcname)
    shutil.rmtree(tmpdir)

# ─── Load template or blank widescreen ───
use_template = False
if os.path.exists(TEMPLATE_POTX):
    convert_potx_to_pptx(TEMPLATE_POTX, TEMPLATE_PPTX)
    prs = Presentation(TEMPLATE_PPTX)
    use_template = True
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId is None:
            for k, v in prs.slides._sldIdLst[0].attrib.items():
                if k.endswith('}id') or k == 'r:id':
                    rId = v
                    break
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    print("Using Cisco PowerPoint Template (LIGHT)")
else:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    print("Cisco template not found — using blank widescreen presentation")

# ─── Colors ───
CISCO_BLUE     = RGBColor(0x04, 0x9F, 0xD9)
CISCO_BLUE_DK  = RGBColor(0x00, 0x7C, 0x92)
SPLUNK_GREEN   = RGBColor(0x65, 0xA6, 0x37)
DEBT_RED       = RGBColor(0xD3, 0x2F, 0x2F)
WARN_AMBER     = RGBColor(0xF9, 0xA8, 0x25)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE      = RGBColor(0xFA, 0xFA, 0xFA)
LIGHT_GRAY     = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY       = RGBColor(0x9E, 0x9E, 0x9E)
DARK_TEXT       = RGBColor(0x21, 0x21, 0x21)
SUBTITLE_GRAY  = RGBColor(0x61, 0x61, 0x61)
PURPLE         = RGBColor(0x4A, 0x14, 0x8C)
PHASE_GREEN    = RGBColor(0x2E, 0x7D, 0x32)
PHASE_BLUE     = RGBColor(0x15, 0x65, 0xC0)
ARCHIVE_ORANGE = RGBColor(0xE6, 0x51, 0x00)
DARK_BG        = RGBColor(0x37, 0x47, 0x4F)

# ─── Helpers ───

def add_slide(prs, use_template):
    if use_template:
        return prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    if font_name:
        p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text_box(slide, left, top, width, height, runs_list, alignment=PP_ALIGN.LEFT):
    """runs_list = list of paragraphs, each paragraph = list of (text, size, bold, color)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for pi, para_runs in enumerate(runs_list):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = alignment
        for text, font_size, bold, color in para_runs:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE, bold=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
    return shape


def add_big_number(slide, left, top, number, color, font_size=72):
    add_text_box(slide, left, top, 3.5, 1.2, str(number),
                 font_size=font_size, bold=True, color=color, alignment=PP_ALIGN.CENTER)


def add_stat_callout(slide, left, top, number, label, color=DEBT_RED):
    add_text_box(slide, left, top, 2.2, 0.65, number, font_size=36, bold=True,
                 color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + 0.55, 2.2, 0.4, label, font_size=11,
                 color=MED_GRAY, alignment=PP_ALIGN.CENTER)


def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def add_section_box(slide, left, top, width, height, header, header_color,
                    body_lines, body_size=10):
    """Rounded rect with header bar + body text. Supports inline **bold**."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = OFF_WHITE
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(0.75)

    hdr = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.32)
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = header_color
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = header
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True

    txBox = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.40),
        Inches(width - 0.30), Inches(height - 0.50)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        parts = re.split(r'(\*\*.*?\*\*)', line)
        if len(parts) > 1:
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    run = p.add_run()
                    run.text = part[2:-2]
                    run.font.size = Pt(body_size)
                    run.font.bold = True
                    run.font.color.rgb = DARK_TEXT
                elif part:
                    run = p.add_run()
                    run.text = part
                    run.font.size = Pt(body_size)
                    run.font.color.rgb = DARK_TEXT
        else:
            p.text = line
            p.font.size = Pt(body_size)
            p.font.color.rgb = DARK_TEXT
    return shape


def add_icon_card(slide, left, top, width, height, icon, title, body,
                  accent_color=CISCO_BLUE, icon_size=28, title_size=13, body_size=10):
    """Card with large emoji icon, bold title, and body text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = OFF_WHITE
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(0.75)

    # Icon
    add_text_box(slide, left + 0.12, top + 0.08, width - 0.24, 0.45, icon,
                 font_size=icon_size, color=accent_color, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, left + 0.12, top + 0.50, width - 0.24, 0.35, title,
                 font_size=title_size, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    # Body
    add_text_box(slide, left + 0.12, top + 0.82, width - 0.24, height - 0.92, body,
                 font_size=body_size, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)
    return shape


# ══════════════════════════════════════════════════════
# TITLE SLIDE: "The Path Forward"
#   Two-part series overview. Sets the frame for both
#   Thursday (problem) and Tuesday March 10 (solution).
# ══════════════════════════════════════════════════════
slide_title = add_slide(prs, use_template)

# ── Dark title bar ──
title_bar = slide_title.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(13.333), Inches(2.0)
)
title_bar.fill.solid()
title_bar.fill.fore_color.rgb = DARK_BG
title_bar.line.fill.background()

add_text_box(slide_title, 0.6, 0.30, 12.1, 0.8,
             "The Path Forward",
             font_size=38, bold=True, color=WHITE)
add_text_box(slide_title, 0.6, 1.05, 12.1, 0.35,
             "Elevating Cisco to a First-Class Citizen in the Splunk Platform",
             font_size=16, color=RGBColor(0xB0, 0xBE, 0xC5))
add_text_box(slide_title, 0.6, 1.50, 12.1, 0.35,
             "A Two-Part Executive Series   •   Amir Khamis & Alan Somerville",
             font_size=12, color=RGBColor(0x78, 0x90, 0x9C))

# ── Part 1 box (left) — THIS THURSDAY ──
p1_shape = slide_title.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.6), Inches(2.5), Inches(5.8), Inches(4.0)
)
p1_shape.fill.solid()
p1_shape.fill.fore_color.rgb = OFF_WHITE
p1_shape.line.color.rgb = CISCO_BLUE
p1_shape.line.width = Pt(2)

# Part 1 header stripe
p1_hdr = slide_title.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.6), Inches(2.5), Inches(5.8), Inches(0.50)
)
p1_hdr.fill.solid()
p1_hdr.fill.fore_color.rgb = CISCO_BLUE
p1_hdr.line.fill.background()

add_rich_text_box(slide_title, 0.8, 2.53, 5.4, 0.44, [
    [("PART 1: TODAY  ", 14, True, WHITE),
     ("—  Defining the Challenge", 14, False, RGBColor(0xBB, 0xDE, 0xFB))],
], alignment=PP_ALIGN.LEFT)

p1_bullets = [
    ("🔍", "The Reality Check", "Voices from the Field"),
    ("📊", "The Data", "Audit of the 120+ App Ecosystem"),
    ("⚠️", "The Risk", "Impact on Cisco Secure Networking GTM"),
]
for i, (icon, title, sub) in enumerate(p1_bullets):
    y = 3.25 + i * 0.90
    add_text_box(slide_title, 1.0, y, 0.5, 0.4, icon, font_size=22, alignment=PP_ALIGN.CENTER)
    add_text_box(slide_title, 1.55, y + 0.02, 4.5, 0.30, title,
                 font_size=16, bold=True, color=DARK_TEXT)
    add_text_box(slide_title, 1.55, y + 0.35, 4.5, 0.25, sub,
                 font_size=12, color=SUBTITLE_GRAY)

# "TODAY" badge
today_badge = slide_title.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(4.6), Inches(6.05), Inches(1.6), Inches(0.30)
)
today_badge.fill.solid()
today_badge.fill.fore_color.rgb = CISCO_BLUE
today_badge.line.fill.background()
add_text_box(slide_title, 4.6, 6.07, 1.6, 0.26, "▶  TODAY",
             font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# ── Part 2 box (right) — MARCH 10 ──
p2_shape = slide_title.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.9), Inches(2.5), Inches(5.8), Inches(4.0)
)
p2_shape.fill.solid()
p2_shape.fill.fore_color.rgb = OFF_WHITE
p2_shape.line.color.rgb = MED_GRAY
p2_shape.line.width = Pt(1)

# Part 2 header stripe
p2_hdr = slide_title.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.9), Inches(2.5), Inches(5.8), Inches(0.50)
)
p2_hdr.fill.solid()
p2_hdr.fill.fore_color.rgb = SUBTITLE_GRAY
p2_hdr.line.fill.background()

add_rich_text_box(slide_title, 7.1, 2.53, 5.4, 0.44, [
    [("PART 2: MARCH 10  ", 14, True, WHITE),
     ("—  Delivering the Solution", 14, False, RGBColor(0xE0, 0xE0, 0xE0))],
], alignment=PP_ALIGN.LEFT)

p2_bullets = [
    ("💊", "The \"Aspirin\"", "Immediate Customer Relief"),
    ("🧬", "The \"Cure\"", "Unified Integration Framework"),
    ("🎬", "5 Technical Demos", "Live Proof-of-Concept Innovations"),
]
for i, (icon, title, sub) in enumerate(p2_bullets):
    y = 3.25 + i * 0.90
    add_text_box(slide_title, 7.3, y, 0.5, 0.4, icon, font_size=22, alignment=PP_ALIGN.CENTER)
    add_text_box(slide_title, 7.85, y + 0.02, 4.5, 0.30, title,
                 font_size=16, bold=True, color=DARK_TEXT)
    add_text_box(slide_title, 7.85, y + 0.35, 4.5, 0.25, sub,
                 font_size=12, color=SUBTITLE_GRAY)

# "MARCH 10" badge
m10_badge = slide_title.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(10.9), Inches(6.05), Inches(1.6), Inches(0.30)
)
m10_badge.fill.solid()
m10_badge.fill.fore_color.rgb = MED_GRAY
m10_badge.line.fill.background()
add_text_box(slide_title, 10.9, 6.07, 1.6, 0.26, "MARCH 10",
             font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Footer
add_rounded_rect(slide_title, 0.3, 7.0, 12.7, 0.32, DARK_BG, "", font_size=9)
ft_title = [
    [("Objective: ", 9, True, WARN_AMBER),
     ("Unanimous executive agreement on scope and severity (today)   →   ", 9, False, RGBColor(0xB0, 0xBE, 0xC5)),
     ("Mandate to proceed with roadmap and ownership model (March 10)", 9, True, WHITE)],
]
add_rich_text_box(slide_title, 0.6, 7.02, 12.2, 0.28, ft_title, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide_title, """HOW TO DELIVER THIS SLIDE:

"Today is Part 1 of a two-part executive series."

[Point to the left box]
"Today, we're defining the challenge. We'll share voices from the field, audit the entire 120+ app ecosystem, and quantify the risk to the Cisco Secure Networking GTM."

"Our goal today is simple: unanimous agreement on the scope and severity of the problem."

[Point to the right box]
"On March 10, we'll deliver the solution. We have an 'Aspirin' — immediate customer relief that's already built. And a 'Cure' — a Unified Integration Framework that elevates Cisco to a first-class citizen in Splunk, the same way AWS and Azure are treated today."

"We'll back that up with 5 live technical demos — proof that this isn't a concept, it's engineering that's already underway."

[Transition]
"Let's start with why this matters — and what the field is telling us."

KEY POINTS:
• Frame this as a TWO-PART series — today is diagnosis, March 10 is treatment
• The "first-class citizen" language is deliberate — AWS and Azure have native Splunk platform experiences; Cisco does not
• Amir & Alan are leading this because it started from field demand at Splunk Show
• Don't rush this slide — let the audience absorb the structure before diving in
""")


# ══════════════════════════════════════════════════════
# PREVIEW SLIDE: "The 5 Pillars of the Cure"
#   Teaser for March 10. Shows what's coming without
#   stealing the thunder. Builds anticipation.
# ══════════════════════════════════════════════════════
slide_preview = add_slide(prs, use_template)

add_text_box(slide_preview, 0.5, 0.3, 12.3, 0.7,
             "Preview: The 5 Pillars of the \"Cure\"  (March 10)",
             font_size=30, bold=True, color=DARK_TEXT)

add_text_box(slide_preview, 0.5, 0.95, 12.3, 0.35,
             "Five innovations that move Cisco from fragmented ecosystem to unified, self-maintaining platform integration",
             font_size=14, color=SUBTITLE_GRAY)

# ── 5 pillar cards ──
pillar_w = 2.30
pillar_h = 4.0
pillar_gap = 0.18
pillar_start_x = 0.45
pillar_y = 1.65

pillars = [
    {
        "num": "1",
        "title": "Automated Testing\n(TACO)",
        "team": "P&T Team",
        "body": ("Technical Add-on\n"
                 "Compatibility testing\n"
                 "framework — zero-day\n"
                 "compatibility with\n"
                 "Splunk Cloud releases"),
        "color": PHASE_GREEN,
    },
    {
        "num": "2",
        "title": "AI Self-Maintainable\nApps",
        "team": "FDSE",
        "body": ("Agentic AI for\n"
                 "automated bug-fixes,\n"
                 "cloud vetting, and\n"
                 "continuous maintenance\n"
                 "without human intervention"),
        "color": PURPLE,
    },
    {
        "num": "3",
        "title": "Auto-Schematization\n(LLM)",
        "team": "Oskar Patnaik",
        "body": ("Using LLMs to automate\n"
                 "CIM-compliant field\n"
                 "mapping from raw data\n"
                 "— guided experience for\n"
                 "any data source"),
        "color": CISCO_BLUE,
    },
    {
        "num": "4",
        "title": "SCAN + Unified\nFramework",
        "team": "Amir Khamis",
        "body": ("The \"Aspirin\" — a Front\n"
                 "Door Navigator for\n"
                 "immediate relief.\n"
                 "The \"Cure\" — Unified\n"
                 "Integration Framework\n"
                 "with CDM integration"),
        "color": DEBT_RED,
    },
    {
        "num": "5",
        "title": "AI-Driven\nTA Factory",
        "team": "Apexti (External AI)",
        "body": ("AI development company\n"
                 "building next-generation\n"
                 "Technology Add-ons\n"
                 "using AI-accelerated\n"
                 "engineering pipelines"),
        "color": WARN_AMBER,
    },
]

for i, p in enumerate(pillars):
    x = pillar_start_x + i * (pillar_w + pillar_gap)

    # Card background
    card = slide_preview.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(pillar_y), Inches(pillar_w), Inches(pillar_h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = OFF_WHITE
    card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    card.line.width = Pt(0.75)

    # Number circle
    num_circle = slide_preview.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + pillar_w/2 - 0.22), Inches(pillar_y + 0.15), Inches(0.44), Inches(0.44)
    )
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = p["color"]
    num_circle.line.fill.background()
    tf = num_circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = p["num"]
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Title
    add_text_box(slide_preview, x + 0.08, pillar_y + 0.68, pillar_w - 0.16, 0.55,
                 p["title"], font_size=11, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    # Team
    add_text_box(slide_preview, x + 0.08, pillar_y + 1.28, pillar_w - 0.16, 0.25,
                 p["team"], font_size=9, bold=True, color=p["color"], alignment=PP_ALIGN.CENTER)

    # Divider line
    div = slide_preview.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x + 0.3), Inches(pillar_y + 1.58), Inches(pillar_w - 0.6), Inches(0.015)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    div.line.fill.background()

    # Body
    add_text_box(slide_preview, x + 0.10, pillar_y + 1.68, pillar_w - 0.20, pillar_h - 1.85,
                 p["body"], font_size=10, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)

# ── Bottom bar ──
add_rounded_rect(slide_preview, 0.3, 6.0, 12.7, 0.40, CISCO_BLUE, "", font_size=9)
preview_bar = [
    [("Each demo: ", 12, True, WHITE),
     ("7 minutes live + 3 minutes Q&A   •   ", 12, False, RGBColor(0xBB, 0xDE, 0xFB)),
     ("\"If we don't fix the data, we can't power the AI agents of the future.\"", 12, True, WHITE)],
]
add_rich_text_box(slide_preview, 0.6, 6.04, 12.2, 0.35, preview_bar, alignment=PP_ALIGN.CENTER)

# Footer
add_rounded_rect(slide_preview, 0.3, 7.0, 12.7, 0.32, DARK_BG, "", font_size=9)
ft_preview = [
    [("March 10 — Part 2: ", 9, True, WARN_AMBER),
     ("Delivering the Solution   |   5 live demos   |   ", 9, False, RGBColor(0xB0, 0xBE, 0xC5)),
     ("From \"Search and Hope\" to \"Detect and Direct\"", 9, True, WHITE)],
]
add_rich_text_box(slide_preview, 0.6, 7.02, 12.2, 0.28, ft_preview, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide_preview, """HOW TO DELIVER THIS SLIDE:

[Quick — don't linger, this is a teaser]
"Before we dive into the problem, let me give you a 30-second preview of where we're headed on March 10."

[Walk across the 5 pillars quickly - left to right]
"Pillar 1: We have an automated testing framework called TACO that ensures zero-day compatibility with every Splunk Cloud release. The P&T team will demo this."

"Pillar 2: AI Self-Maintainable Apps — Agentic AI that can automatically fix bugs, vet for cloud compliance, and maintain add-ons without human intervention. FDSE is building this."

"Pillar 3: Auto-Schematization — Oskar Patnaik will show how LLMs can automatically create CIM-compliant field mappings from raw data. No more manual schema work."

"Pillar 4: This is mine — SCAN, the Splunk Cisco App Navigator, which is the 'Aspirin' for immediate relief, plus the Unified Cisco Integration Framework concept with Cloud Data Manager integration — that's the 'Cure.'"

"Pillar 5: AI-Driven TA Factory — an external AI company called Apexti that's building next-generation Technology Add-ons using AI-accelerated engineering."

[Transition back to today]
"Each demo is 7 minutes with 3 minutes for Q&A. But first — you need to see WHY all of this is necessary. Let's look at the data."

KEY POINTS:
• This is a TEASER — spend no more than 60 seconds on this slide
• The purpose is to show that solutions EXIST so that the problem slides don't feel hopeless
• The quote at the bottom is Jeetu-aligned language: "If we don't fix the data, we can't power the AI agents of the future"
• Don't explain the demos in detail — save that for March 10
""")


# ══════════════════════════════════════════════════════
# AGENDA SLIDE: "Today's Agenda"
#   Clear roadmap for Part 1. Lets the audience know
#   exactly what they're walking through.
# ══════════════════════════════════════════════════════
slide_agenda = add_slide(prs, use_template)

add_text_box(slide_agenda, 0.5, 0.3, 12.3, 0.7,
             "Today's Agenda — Part 1: Defining the Challenge",
             font_size=30, bold=True, color=DARK_TEXT)

add_text_box(slide_agenda, 0.5, 0.95, 12.3, 0.35,
             "The data-driven case for why Cisco integrations are the most urgent priority for Splunk engineering",
             font_size=14, color=SUBTITLE_GRAY)

# ── Agenda items ──
agenda_items = [
    ("1", "#1 Most Requested Demo",
     "The $500B TAM, customer demand at Splunk Show, and the \"Last Mile\" risk",
     CISCO_BLUE),
    ("2", "Why We Are Here",
     "Four converging forces: Renewal Risk, Data Integrity, Agentic Future, Customer Demand",
     DEBT_RED),
    ("3", "What Just Happened on Splunkbase",
     "121 apps → 62 archived → 59 remaining — the scale of the ecosystem cleanup",
     WARN_AMBER),
    ("4", "The Number That Actually Matters",
     "Of 59 remaining apps, only 18 have enterprise-grade support (Cisco or Splunk backed)",
     PURPLE),
    ("5", "Still Running — Cisco Apps in Splunk Cloud",
     "89 apps deployed across 2,973 Cloud stacks — the archiving removed listings, not usage",
     PHASE_GREEN),
    ("6", "The Cisco Coverage Gap",
     "47 active products, 17 on the GTM roadmap, only 28% enterprise-grade coverage today",
     CISCO_BLUE_DK),
    ("7", "Voices from the Field + Systemic Issues",
     "7 pain points customers report and the 4 root causes that keep this problem recurring",
     DEBT_RED),
]

item_start_y = 1.55
item_h = 0.72
item_gap = 0.08

for i, (num, title, desc, color) in enumerate(agenda_items):
    y = item_start_y + i * (item_h + item_gap)

    # Row background
    row = slide_agenda.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(y), Inches(12.3), Inches(item_h)
    )
    row.fill.solid()
    row.fill.fore_color.rgb = OFF_WHITE
    row.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    row.line.width = Pt(0.5)

    # Number pill
    pill = slide_agenda.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.75), Inches(y + 0.12), Inches(0.48), Inches(0.48)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    tf = pill.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = num
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Title
    add_text_box(slide_agenda, 1.45, y + 0.05, 10.8, 0.30, title,
                 font_size=16, bold=True, color=DARK_TEXT)
    # Description
    add_text_box(slide_agenda, 1.45, y + 0.36, 10.8, 0.30, desc,
                 font_size=11, color=SUBTITLE_GRAY)

# Footer
add_rounded_rect(slide_agenda, 0.3, 7.0, 12.7, 0.32, DARK_BG, "", font_size=9)
ft_agenda = [
    [("Goal: ", 9, True, WARN_AMBER),
     ("Unanimous executive agreement on the scope and severity of Cisco integration challenges", 9, True, WHITE)],
]
add_rich_text_box(slide_agenda, 0.6, 7.02, 12.2, 0.28, ft_agenda, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide_agenda, """HOW TO DELIVER THIS SLIDE:

[Quick overview — 15 seconds max]
"Here's what we'll cover today. Seven sections, each building on the last."

"We'll start with the market context — why Cisco is the #1 most requested integration. Then we'll lay out the four forces converging to make this urgent."

"After that, we'll walk through the data: what happened on Splunkbase, the enterprise-grade reality, and the Cloud deployment evidence."

"We'll close with the full coverage gap analysis and what we're hearing directly from the field."

[Transition]
"Let's start with the demand signal."

KEY POINTS:
• Don't read each item — the audience can read
• This slide exists so people know the structure and can pace themselves
• If time is tight, this is the first slide to skip
""")


# ══════════════════════════════════════════════════════
# OPENING SLIDE A: "#1 Most Requested Demo"
#   Why we are here — market demand, TAM, stakes.
#   Steve Jobs: Start with the customer. What do they want?
# ══════════════════════════════════════════════════════
slide_open_a = add_slide(prs, use_template)

# Title
add_text_box(slide_open_a, 0.5, 0.3, 12.3, 0.7,
             "#1 Most Requested Demo at Splunk Show",
             font_size=32, bold=True, color=DARK_TEXT)

# Subtitle
add_text_box(slide_open_a, 0.5, 0.95, 12.3, 0.35,
             "Cisco Networking is the integration customers ask for more than any other — and they can't find the official path",
             font_size=14, color=SUBTITLE_GRAY)

# ── Hero stat: $500B ──
add_text_box(slide_open_a, 0.5, 1.8, 12.3, 1.0,
             "$500B",
             font_size=96, bold=True, color=CISCO_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide_open_a, 0.5, 3.0, 12.3, 0.35,
             "Cisco's Total Addressable Market",
             font_size=18, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
add_text_box(slide_open_a, 1.5, 3.4, 10.3, 0.35,
             "Splunk integrations are the 'Last Mile' — the link between that infrastructure spend and security outcomes",
             font_size=13, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)

# ── TAM breakdown: 4 pillars ──
pill_y = 4.1
pill_h = 1.15
pill_w = 2.8
gap = 0.35
start_x = 0.8

tam_data = [
    ("~$210B", "Networking", "Campus, Branch,\nData Center, Core", CISCO_BLUE),
    ("~$120B", "Security", "Firewall, SSE,\nIdentity, Endpoint", DEBT_RED),
    ("~$100B", "Observability", "AppDynamics,\nThousandEyes, Splunk", PURPLE),
    ("~$50B+", "Collaboration", "Webex,\nContact Center", PHASE_GREEN),
]

for i, (amount, label, detail, color) in enumerate(tam_data):
    x = start_x + i * (pill_w + gap)
    shape = slide_open_a.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(pill_y), Inches(pill_w), Inches(pill_h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Amount
    add_text_box(slide_open_a, x + 0.1, pill_y + 0.08, pill_w - 0.2, 0.45, amount,
                 font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide_open_a, x + 0.1, pill_y + 0.48, pill_w - 0.2, 0.25, label,
                 font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Detail
    add_text_box(slide_open_a, x + 0.1, pill_y + 0.72, pill_w - 0.2, 0.40, detail,
                 font_size=9, color=RGBColor(0xE0, 0xE0, 0xE0), alignment=PP_ALIGN.CENTER)

# ── Bottom punchline bar ──
add_text_box(slide_open_a, 0.8, 5.55, 11.7, 0.50,
             "The Data Gravity is astronomical. Splunk is the only platform that can provide a Single Pane of Glass across this $500B ecosystem.",
             font_size=13, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# ── The risk callout ──
risk_shape = slide_open_a.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1.5), Inches(6.1), Inches(10.3), Inches(0.60)
)
risk_shape.fill.solid()
risk_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
risk_shape.line.color.rgb = DEBT_RED
risk_shape.line.width = Pt(1)

risk_runs = [
    [("⚠  THE RISK:  ", 11, True, DEBT_RED),
     ("If our TAs are broken or fragmented, we are devaluing Cisco's $500B position — ", 11, False, DARK_TEXT),
     ("customers can't see or secure the network they just bought.", 11, True, DEBT_RED)],
]
add_rich_text_box(slide_open_a, 1.7, 6.15, 9.9, 0.50, risk_runs, alignment=PP_ALIGN.CENTER)

# Footer
add_rounded_rect(slide_open_a, 0.3, 7.0, 12.7, 0.32, DARK_BG, "", font_size=9)
footer_a = [
    [("Source: ", 9, True, WARN_AMBER),
     ("Cisco Investor Day 2023 — Total Addressable Market projections for 2025/2026   |   ", 9, False, RGBColor(0xB0, 0xBE, 0xC5)),
     ("Splunk Show: Cisco Networking — #1 requested demo, all events", 9, True, WHITE)],
]
add_rich_text_box(slide_open_a, 0.6, 7.02, 12.2, 0.28, footer_a, alignment=PP_ALIGN.CENTER)

# Speaker notes
add_speaker_notes(slide_open_a, """HOW TO DELIVER THIS SLIDE:

[Open with conviction]
"Cisco Networking is the number one most requested demo at Splunk Show. Every event, every city — customers ask for this integration more than any other."

[Pause, let it land]
"And they can't find the official path. That's where we are today."

[Point to $500B]
"Cisco sits at the center of a $500B addressable market. This isn't a niche vendor — this is THE platform that most enterprises run their infrastructure on."

[Walk through the 4 pillars]
"$210B in networking alone — campus, branch, data center. $120B in security. $100B in observability — which now includes us, Splunk. And $50B+ in collaboration."

[Deliver the insight]
"Think about the data gravity. If the market for the hardware and services is $500B, the volume of telemetry, logs, events, and flows generated by that infrastructure is astronomical. And Splunk is the only platform positioned to be the single pane of glass across all of it."

[Hit the risk]
"But here's the risk: if our TAs are broken, fragmented, or missing — we're effectively making that $500B invisible. Customers bought the network. They bought Splunk. And the integration between the two... is the part that's failing."

KEY POINTS:
• The $500B figure comes from Cisco's 2021 and 2023 Investor updates — TAM projections for 2025/2026
• Secure Networking is the core engine of Cisco's growth strategy
• "Last Mile" framing: integrations are the critical link between infrastructure spend and outcomes
• This slide establishes WHY this matters before showing WHAT happened
""")


# ══════════════════════════════════════════════════════
# OPENING SLIDE B: "Why We Are Here"
#   Four forces converging — Renewal Risk, Data Integrity,
#   Agentic Future, Customer Demand.
#   Steve Jobs: Frame the stakes before the data.
# ══════════════════════════════════════════════════════
slide_open_b = add_slide(prs, use_template)

# Title
add_text_box(slide_open_b, 0.5, 0.3, 12.3, 0.7,
             "Why We Are Here",
             font_size=32, bold=True, color=DARK_TEXT)

# Subtitle
add_text_box(slide_open_b, 0.5, 0.95, 12.3, 0.35,
             "Four forces converging to make Cisco integrations the most urgent priority for Splunk engineering",
             font_size=14, color=SUBTITLE_GRAY)

# ── Four cards in a row ──
card_w = 2.85
card_h = 3.5
card_gap = 0.20
card_start_x = 0.55
card_y = 1.65

cards = [
    {
        "icon": "🔥",
        "title": "RENEWAL RISK",
        "body": ("Field teams report customers\n"
                 "refusing to renew because\n"
                 "\"broken\" integrations prevent\n"
                 "security outcomes.\n\n"
                 "If they can't see the network,\n"
                 "they can't justify the platform."),
        "accent": DEBT_RED,
    },
    {
        "icon": "🗑️",
        "title": "DATA INTEGRITY",
        "body": ("Garbage In, Garbage Out.\n\n"
                 "Bad data breaks our premium\n"
                 "solutions — ARI, RBA, ITSI,\n"
                 "and Enterprise Security all\n"
                 "depend on correctly parsed\n"
                 "and CIM-mapped Cisco data."),
        "accent": WARN_AMBER,
    },
    {
        "icon": "🤖",
        "title": "THE AGENTIC FUTURE",
        "body": ("The SOC of the Future runs\n"
                 "on AI Agents. Those agents\n"
                 "need high-fidelity context\n"
                 "to make decisions.\n\n"
                 "Broken TAs = blind agents\n"
                 "= wrong automated actions."),
        "accent": PURPLE,
    },
    {
        "icon": "📡",
        "title": "CUSTOMER DEMAND",
        "body": ("#1 requested integration\n"
                 "at every Splunk Show.\n\n"
                 "89 Cisco apps deployed\n"
                 "across 2,973 Cloud stacks.\n\n"
                 "Customers want this.\n"
                 "They just can't find it."),
        "accent": CISCO_BLUE,
    },
]

for i, card in enumerate(cards):
    x = card_start_x + i * (card_w + card_gap)
    add_icon_card(slide_open_b, x, card_y, card_w, card_h,
                  card["icon"], card["title"], card["body"],
                  accent_color=card["accent"],
                  icon_size=32, title_size=14, body_size=11)

# ── Connector bar beneath cards ──
bar_y = card_y + card_h + 0.30
connector = slide_open_b.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.55), Inches(bar_y), Inches(12.2), Inches(0.55)
)
connector.fill.solid()
connector.fill.fore_color.rgb = CISCO_BLUE
connector.line.fill.background()

bar_runs = [
    [("The convergence: ", 12, True, WHITE),
     ("$500B in Cisco infrastructure × broken integrations × AI dependency = ", 12, False, RGBColor(0xBB, 0xDE, 0xFB)),
     ("existential priority", 12, True, WHITE)],
]
add_rich_text_box(slide_open_b, 0.75, bar_y + 0.05, 11.8, 0.45, bar_runs, alignment=PP_ALIGN.CENTER)

# ── Bottom quote ──
add_rich_text_box(slide_open_b, 1.0, bar_y + 0.75, 11.3, 0.5, [
    [("\"High-fidelity Splunk integrations are the critical link that transforms ", 12, False, SUBTITLE_GRAY),
     ("$500B in infrastructure spend ", 12, True, CISCO_BLUE),
     ("into actionable security and operational outcomes.\"", 12, False, SUBTITLE_GRAY)],
], alignment=PP_ALIGN.CENTER)

# Footer
add_rounded_rect(slide_open_b, 0.3, 7.0, 12.7, 0.32, DARK_BG, "", font_size=9)
footer_b = [
    [("Context: ", 9, True, WARN_AMBER),
     ("Renewal risk data from field SE & sales reports   |   ", 9, False, RGBColor(0xB0, 0xBE, 0xC5)),
     ("Cloud deployment data: 89 Cisco apps on 2,973 stacks (Splunk internal)", 9, True, WHITE)],
]
add_rich_text_box(slide_open_b, 0.6, 7.02, 12.2, 0.28, footer_b, alignment=PP_ALIGN.CENTER)

# Speaker notes
add_speaker_notes(slide_open_b, """HOW TO DELIVER THIS SLIDE:

[Walk left to right across the four cards]

"There are four forces converging right now that make this the most urgent engineering priority."

[Card 1 — Renewal Risk]
"First: renewal risk. We're hearing from the field that customers are refusing to renew because the integrations they depend on are broken. If they can't see their Cisco network in Splunk, they can't justify the platform spend."

[Card 2 — Data Integrity]
"Second: data integrity. Garbage in, garbage out. Every premium solution we sell — Adaptive Response, Risk-Based Alerting, ITSI — depends on correctly parsed, CIM-mapped data. When Cisco TAs are broken, all of those solutions are compromised."

[Card 3 — The Agentic Future]
"Third: the agentic future. We're building the SOC of the Future on AI agents. Those agents need high-fidelity context to make decisions. A broken TA doesn't just mean a missing dashboard — it means an AI agent making the wrong automated action because it can't understand the data."

[Card 4 — Customer Demand]
"And fourth: customer demand. This isn't speculative. Cisco Networking is the #1 most requested demo at Splunk Show. 89 Cisco apps are deployed across nearly 3,000 Cloud stacks right now. Customers want this integration. They just can't find the official path."

[Point to the connector bar]
"These four forces don't just add up — they multiply. $500 billion in Cisco infrastructure, times broken integrations, times AI dependency, equals an existential priority."

KEY POINTS:
• This slide frames WHY before showing WHAT — do not skip ahead
• The renewable risk stories are real — reference specific field conversations if you have them
• ARI = Adaptive Response Initiative, RBA = Risk-Based Alerting, ITSI = IT Service Intelligence
• The "agentic future" connects to Splunk's strategic AI roadmap — this is language leadership uses
• End with the convergence — these aren't four separate problems, they're one compounding crisis
""")


# ══════════════════════════════════════════════════════
# SLIDE 1: "What Happened to Cisco on Splunkbase"
#   Setup the context — the scale of the cleanup
#   Steve Jobs: Start with a simple fact, let it land.
# ══════════════════════════════════════════════════════
slide1 = add_slide(prs, use_template)

# Title
add_text_box(slide1, 0.5, 0.3, 12.3, 0.7,
             "Cisco on Splunkbase — What Just Happened",
             font_size=30, bold=True, color=DARK_TEXT)

# Subtitle
add_text_box(slide1, 0.5, 0.95, 12.3, 0.35,
             "Splunk archived 51% of the Cisco app ecosystem in a single sweep",
             font_size=14, color=SUBTITLE_GRAY)

# ── Three big numbers: Before → Archived → After ──
# Centered, massive, with plenty of breathing room

# BEFORE
add_big_number(slide1, 1.0, 2.0, "121", CISCO_BLUE, font_size=80)
add_text_box(slide1, 1.0, 3.2, 3.5, 0.5, "Apps & Add-ons",
             font_size=16, bold=True, color=CISCO_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, 1.0, 3.6, 3.5, 0.3, "on Splunkbase before",
             font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow
add_text_box(slide1, 4.4, 2.4, 0.8, 0.8, "→",
             font_size=48, bold=True, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ARCHIVED
add_big_number(slide1, 5.0, 2.0, "62", DEBT_RED, font_size=80)
add_text_box(slide1, 5.0, 3.2, 3.5, 0.5, "Archived",
             font_size=16, bold=True, color=DEBT_RED, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, 5.0, 3.6, 3.5, 0.3, "removed from Splunkbase",
             font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow
add_text_box(slide1, 8.4, 2.4, 0.8, 0.8, "→",
             font_size=48, bold=True, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# REMAINING
add_big_number(slide1, 9.0, 2.0, "59", WARN_AMBER, font_size=80)
add_text_box(slide1, 9.0, 3.2, 3.5, 0.5, "Remaining",
             font_size=16, bold=True, color=WARN_AMBER, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, 9.0, 3.6, 3.5, 0.3, "still on Splunkbase today",
             font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ── "But wait..." teaser line ──
add_text_box(slide1, 0.5, 4.3, 12.3, 0.4,
             "22 apps and 40 add-ons — wiped from the catalog. But the real story isn't what was removed.",
             font_size=13, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)

# ── Two clean boxes: What was lost vs What remains ──
box_y = 4.95
box_h = 1.75

add_section_box(slide1, 0.8, box_y, 5.5, box_h,
    "What Was Archived", DEBT_RED,
    [
        "**62** apps & add-ons removed (22 apps + 40 add-ons)",
        "",
        "**10** Cisco products lost all Splunkbase coverage",
        "   (all 10 were already EOL/EOS by Cisco)",
        "",
        "**4** active products impacted: AppDynamics, Webex,",
        "   Meeting Server, CUCM — now developer-supported only",
    ], body_size=11)

add_section_box(slide1, 7.0, box_y, 5.5, box_h,
    "What Still Remains", PHASE_GREEN,
    [
        "**59** apps & add-ons (27 apps + 32 add-ons)",
        "",
        "**46** unique Cisco products still have coverage",
        "",
        "**13** SOAR connectors unaffected",
        "   (tracked separately from Splunkbase)",
        "",
        "Core integrations (ASA, ISE, Meraki) intact",
    ], body_size=11)

# ── Footer ──
add_rounded_rect(slide1, 0.3, 7.0, 12.7, 0.32, DARK_BG, "", font_size=9)
footer1 = [
    [("Splunk Cloud data: ", 9, True, WARN_AMBER),
     ("89 of these 121 apps are still actively deployed across 2,973 Cloud stacks — ", 9, False, RGBColor(0xB0, 0xBE, 0xC5)),
     ("the archiving removed listings, not running code", 9, True, WHITE)],
]
add_rich_text_box(slide1, 0.6, 7.02, 12.2, 0.28, footer1, alignment=PP_ALIGN.CENTER)

# Speaker notes
add_speaker_notes(slide1, """HOW TO DELIVER THIS SLIDE:

"Let me show you what just happened to Cisco's presence on Splunkbase."

[Click — show 121]
"We had 121 apps and add-ons. That's the largest vendor ecosystem on Splunkbase outside of Splunk itself. 49 apps, 72 add-ons — covering everything from ASA firewalls to Webex collaboration."

[Click — show 62]
"Splunk archived 62 of them. Over half. Gone from the catalog in one sweep."

[Click — show 59]
"59 remain. That sounds manageable, right? That's the story most people stop at."

[Pause]
"But here's what they're missing — and it's the most important part..."

KEY POINTS TO EMPHASIZE:
• 10 retired Cisco products lost ALL coverage — but they were already EOL, so this was housekeeping
• The 4 active products (AppDynamics, Webex, CMS, CUCM) are the real concern — they now only have developer-supported add-ons, meaning no SLA
• The footer stat is your ace card if anyone says "doesn't matter" — 89 of these apps are still running on 2,973 Cloud stacks RIGHT NOW
• This slide sets up the tension. Don't resolve it here. Let them sit with "59 remaining" assuming that's fine.
""")


# ══════════════════════════════════════════════════════
# SLIDE 2: "The Number That Matters"
#   The reveal — 59 remaining sounds okay, but only 18
#   have enterprise-grade support. This is the "aha."
#   Steve Jobs: One number that reframes everything.
# ══════════════════════════════════════════════════════
slide2 = add_slide(prs, use_template)

# Title — understated
add_text_box(slide2, 0.5, 0.3, 12.3, 0.7,
             "The Number That Actually Matters",
             font_size=30, bold=True, color=DARK_TEXT)

add_text_box(slide2, 0.5, 0.95, 12.3, 0.35,
             "Not how many apps remain — how many your customers can actually deploy",
             font_size=14, color=SUBTITLE_GRAY)

# ── The big reveal: 59 → 18 ──
# Left side: "59 Remaining" (what people think)
add_text_box(slide2, 0.8, 1.8, 5.5, 0.4,
             "Apps remaining on Splunkbase",
             font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_big_number(slide2, 1.8, 2.2, "59", RGBColor(0xCC, 0xCC, 0xCC), font_size=96)
add_text_box(slide2, 0.8, 3.6, 5.5, 0.3, "What everyone sees",
             font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Right side: "18" (the real number) — in bold, vibrant
add_text_box(slide2, 7.0, 1.8, 5.5, 0.4,
             "With Cisco or Splunk support",
             font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_big_number(slide2, 8.0, 2.2, "18", PHASE_BLUE, font_size=96)
add_text_box(slide2, 7.0, 3.6, 5.5, 0.3, "What actually matters to enterprise customers",
             font_size=12, bold=True, color=PHASE_BLUE, alignment=PP_ALIGN.CENTER)

# ── Divider line ──
div = slide2.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(6.5), Inches(2.0), Inches(0.02), Inches(2.0)
)
div.fill.solid()
div.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
div.line.fill.background()

# ── Explanation: The support breakdown ──
explain_y = 4.2
add_text_box(slide2, 0.8, explain_y, 12.0, 0.35,
             "Why enterprise customers see '18' — not '59'",
             font_size=15, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Four support-level boxes across the bottom
sup_y = 4.75
sup_h = 1.45
sup_w = 2.8
sup_gap = 0.2
sx1 = 0.6
sx2 = sx1 + sup_w + sup_gap
sx3 = sx2 + sup_w + sup_gap
sx4 = sx3 + sup_w + sup_gap

# Cisco-supported
add_section_box(slide2, sx1, sup_y, sup_w, sup_h,
    "Cisco Supported", CISCO_BLUE,
    [
        "**11** apps & add-ons",
        "",
        "Full Cisco engineering & SLA",
        "Field-tested, Cisco TAC-backed",
        "Enterprise customers: ✓ deploy",
    ], body_size=10)

# Splunk-supported
add_section_box(slide2, sx2, sup_y, sup_w, sup_h,
    "Splunk Supported", SPLUNK_GREEN,
    [
        "**7** apps & add-ons",
        "",
        "Splunk-maintained with SLA",
        "Listed, vetted, supported",
        "Enterprise customers: ✓ deploy",
    ], body_size=10)

# Developer-supported
add_section_box(slide2, sx3, sup_y, sup_w, sup_h,
    "Developer Supported", ARCHIVE_ORANGE,
    [
        "**31** apps & add-ons",
        "",
        "Community/developer maintained",
        "No SLA, no guaranteed updates",
        "Enterprise customers: ✗ risky",
    ], body_size=10)

# Not supported
add_section_box(slide2, sx4, sup_y, sup_w, sup_h,
    "Not Supported", DEBT_RED,
    [
        "**10** apps & add-ons",
        "",
        "No maintainer, no updates",
        "Effectively abandoned",
        "Enterprise customers: ✗ won't deploy",
    ], body_size=10)

# ── Punchline beneath the boxes ──
punch_y = 6.4
add_rounded_rect(slide2, 0.6, punch_y, 12.1, 0.45, PHASE_BLUE, "", font_size=12)
punch_runs = [
    [("Enterprise customers require vendor-backed support SLAs.", 13, True, WHITE),
     ("  Only ", 13, False, WHITE),
     ("18 of 59 (31%) ", 13, True, WARN_AMBER),
     ("meet that bar.", 13, False, WHITE)],
]
add_rich_text_box(slide2, 0.9, punch_y + 0.05, 11.5, 0.38, punch_runs, alignment=PP_ALIGN.CENTER)

# Footer
add_rounded_rect(slide2, 0.3, 7.05, 12.7, 0.32, DARK_BG, "", font_size=9)
footer2 = [
    [("18 = 11 Cisco + 7 Splunk supported.", 9, True, WHITE),
     ("  The other 41 (developer + unsupported) are technically available, but ", 9, False, RGBColor(0xB0, 0xBE, 0xC5)),
     ("enterprise procurement won't approve them.", 9, True, WARN_AMBER)],
]
add_rich_text_box(slide2, 0.6, 7.07, 12.2, 0.28, footer2, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide2, """HOW TO DELIVER THIS SLIDE:

[Transition from Slide 1]
"So we have 59 apps remaining. Sounds reasonable. Here's the problem."

[Click — show 59 grayed out on the left]
"59 apps. That's what everyone sees on Splunkbase."

[Click — show 18 in blue on the right]
"18. That's the number that matters."

[Pause — let it land]

"Why 18? Because enterprise customers — the ones writing the big checks — they require vendor-backed support. They need an SLA. They need to know that when something breaks at 2am, someone picks up the phone."

[Walk through the four boxes]
"11 apps are Cisco-supported. Cisco engineering, Cisco TAC, full SLA."
"7 are Splunk-supported. Splunk maintains them, Splunk backs them."
"That's 18. That's the green zone."

"Now look at the orange and red."
"31 are developer-supported — community maintained, best effort, no SLA."
"10 have no support at all. No maintainer. No updates."

[Deliver the punchline]
"When an enterprise customer evaluates Cisco data onboarding for Splunk, they don't see 59 options. They see 18. The other 41 might as well not exist."

"That's a 31% coverage rate. For the #1 infrastructure vendor in the world."

WHAT IS "18 ENTERPRISE-GRADE" AND HOW TO EXPLAIN IT:
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
The "18" number is the count of Cisco-related Splunkbase apps that carry an official support commitment from either Cisco OR Splunk.

• 11 Cisco-supported = Cisco built it, Cisco maintains it, Cisco TAC will help you
• 7 Splunk-supported = Splunk built or adopted it, Splunk Support backs it

Why this matters: Enterprise security and IT teams operate under compliance and procurement policies that REQUIRE vendor-backed support. They will NOT deploy an app where:
- No one is contractually obligated to fix vulnerabilities
- There's no SLA for bug fixes or compatibility updates
- The "maintainer" is an anonymous community developer

So when we say "18 enterprise-grade" — we mean 18 apps that pass the enterprise procurement test. The other 41 are technically on Splunkbase but functionally invisible to enterprise buyers.

This is the gap we need leadership to understand. It's not a Splunkbase listing problem — it's a COVERAGE problem.
""")


# ══════════════════════════════════════════════════════
# SLIDE 2b: "Still Running — Cisco Apps in Splunk Cloud"
#   The proof point: 89 apps still deployed on 2,973 stacks.
#   Placed here so the audience sees the URGENCY right after
#   learning only 18 are enterprise-grade.
# ══════════════════════════════════════════════════════
slide_cloud = add_slide(prs, use_template)

add_text_box(slide_cloud, 0.5, 0.3, 12.3, 0.7,
             "Still Running — Cisco Apps in Splunk Cloud",
             font_size=30, bold=True, color=DARK_TEXT)

add_text_box(slide_cloud, 0.5, 0.95, 12.3, 0.35,
             "Despite the Splunkbase archiving, 73% of Cisco apps remain actively deployed across Splunk Cloud stacks",
             font_size=14, color=SUBTITLE_GRAY)

# ── Hero stat callouts ──
stat_y = 1.05
add_stat_callout(slide_cloud, 0.3,  stat_y, "89",    "Cisco Apps\nDeployed", CISCO_BLUE)
add_stat_callout(slide_cloud, 2.85, stat_y, "2,973", "Cloud Stacks\nwith Cisco Apps", PHASE_GREEN)
add_stat_callout(slide_cloud, 5.4,  stat_y, "73%",   "of 121 Apps\nStill Running", WARN_AMBER)
add_stat_callout(slide_cloud, 7.95, stat_y, "~30",   "Archived Apps\nStill Deployed", DEBT_RED)

# Total versions across all deployed apps
CLOUD_DATA = [
    ("Splunk_TA_cisco-asa", 10, 1427), ("TA-cisco_ios", 15, 1295),
    ("Splunk_TA_cisco-ise", 6, 917), ("Splunk_TA_cisco_meraki", 13, 833),
    ("cisco_ios", 13, 741), ("CiscoSecurityCloud", 25, 622),
    ("Splunk_TA_cisco-esa", 7, 463), ("Splunk_CiscoISE", 7, 459),
    ("duo_splunkapp", 11, 428), ("TA-eStreamer", 20, 334),
    ("Splunk_TA_Talos_Intelligence", 2, 304), ("Splunk_TA_sourcefire", 3, 276),
    ("TA-cisco-cloud-security-addon", 9, 268), ("TA-meraki", 4, 256),
    ("TA-cisco-cloud-security-umbrella-addon", 7, 236), ("cisco-cloud-security", 13, 229),
    ("firepower_dashboard", 6, 207), ("TA-cisco_acs", 5, 203),
    ("Splunk_TA_cisco-wsa", 8, 186), ("TA_cisco_catalyst", 7, 161),
]
total_versions = sum(v for _, v, _ in CLOUD_DATA)  # top 20 versions

add_stat_callout(slide_cloud, 10.5, stat_y, f"{total_versions}+", "App Versions\nin Production", PURPLE)

# Dots between stats
for ax in [2.45, 5.0, 7.55, 10.1]:
    add_text_box(slide_cloud, ax, stat_y + 0.1, 0.45, 0.45, "•", font_size=20,
                 bold=True, color=RGBColor(0xDD, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)

# ── Top 15 Bar Chart (left side) ──
BAR_BLUE_1 = RGBColor(0x01, 0x57, 0x9B)
BAR_BLUE_2 = RGBColor(0x02, 0x88, 0xD1)
BAR_BLUE_3 = RGBColor(0x03, 0xA9, 0xF4)

add_text_box(slide_cloud, 0.3, 2.05, 5.0, 0.3,
             "TOP 15 MOST DEPLOYED  (by Cloud stack count)", font_size=11, bold=True, color=DARK_TEXT)

DISPLAY_NAMES = {
    "Splunk_TA_cisco-asa": "Cisco ASA (TA)",
    "TA-cisco_ios": "Cisco IOS (TA)",
    "Splunk_TA_cisco-ise": "Cisco ISE (TA)",
    "Splunk_TA_cisco_meraki": "Cisco Meraki (TA)",
    "cisco_ios": "Cisco IOS (App)",
    "CiscoSecurityCloud": "Security Cloud (SCAN)",
    "Splunk_TA_cisco-esa": "Cisco ESA (TA)",
    "Splunk_CiscoISE": "Cisco ISE (App)",
    "duo_splunkapp": "Cisco Duo (App)",
    "TA-eStreamer": "Firepower eStreamer (TA)",
    "Splunk_TA_Talos_Intelligence": "Talos Intelligence (TA)",
    "Splunk_TA_sourcefire": "Sourcefire (TA)",
    "TA-cisco-cloud-security-addon": "Umbrella Cloud Sec (TA)",
    "TA-meraki": "Meraki (Community TA)",
    "TA-cisco-cloud-security-umbrella-addon": "Umbrella DNS (TA)",
}

ALL_CLOUD_STACKS = [
    1427, 1295, 917, 833, 741, 622, 463, 459, 428, 334,
    304, 276, 268, 256, 236, 229, 207, 203, 186, 161,
    154, 149, 107, 106, 89, 78, 68, 56, 52, 50, 50, 47, 45,
    40, 40, 38, 36, 36, 36, 35, 35, 31, 30, 30, 29, 27, 26,
    24, 24, 22, 22, 18, 15, 15, 14, 14, 14, 13, 13, 13, 11,
    10, 9, 9, 7, 7, 6, 6, 6, 6, 6, 5, 5, 5, 4, 4, 3, 3,
    2, 2, 2, 2, 2, 2, 2, 1, 1, 1, 1,
]

bar_y_start = 2.4
bar_h = 0.22
bar_gap = 0.04
max_stacks = 1427
top15_data = list(zip(
    ["Splunk_TA_cisco-asa", "TA-cisco_ios", "Splunk_TA_cisco-ise",
     "Splunk_TA_cisco_meraki", "cisco_ios", "CiscoSecurityCloud",
     "Splunk_TA_cisco-esa", "Splunk_CiscoISE", "duo_splunkapp",
     "TA-eStreamer", "Splunk_TA_Talos_Intelligence", "Splunk_TA_sourcefire",
     "TA-cisco-cloud-security-addon", "TA-meraki", "TA-cisco-cloud-security-umbrella-addon"],
    [1427, 1295, 917, 833, 741, 622, 463, 459, 428, 334, 304, 276, 268, 256, 236]
))

for i, (name, stacks) in enumerate(top15_data):
    y = bar_y_start + i * (bar_h + bar_gap)
    display = DISPLAY_NAMES.get(name, name)
    bar_color = BAR_BLUE_1 if stacks >= 500 else (BAR_BLUE_2 if stacks >= 200 else BAR_BLUE_3)

    add_text_box(slide_cloud, 0.3, y - 0.02, 2.7, bar_h + 0.04,
                 display, font_size=8, color=DARK_TEXT, alignment=PP_ALIGN.RIGHT)

    bar_w = max(0.08, (stacks / max_stacks) * 3.6)
    shape = slide_cloud.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.1), Inches(y), Inches(bar_w), Inches(bar_h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bar_color
    shape.line.fill.background()

    add_text_box(slide_cloud, 3.1 + bar_w + 0.05, y - 0.02, 0.7, bar_h + 0.04,
                 f"{stacks:,}", font_size=8, bold=True, color=bar_color)

remaining_y = bar_y_start + 15 * (bar_h + bar_gap) + 0.05
add_text_box(slide_cloud, 0.3, remaining_y, 7.2, 0.25,
             "+ 74 more apps deployed across 1–229 stacks each",
             font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ── Right side: Deployment Tiers ──
tier_x = 8.0
tier_w = 4.9
tier_y_start = 2.05

add_text_box(slide_cloud, tier_x, tier_y_start, tier_w, 0.3,
             "DEPLOYMENT TIERS", font_size=12, bold=True, color=DARK_TEXT)

tier1_count = sum(1 for s in ALL_CLOUD_STACKS if s >= 500)
tier2_count = sum(1 for s in ALL_CLOUD_STACKS if 100 <= s < 500)
tier3_count = sum(1 for s in ALL_CLOUD_STACKS if 10 <= s < 100)
tier4_count = sum(1 for s in ALL_CLOUD_STACKS if s < 10)

tier_entries = [
    (f"TIER 1 — 500+ Stacks", f"{tier1_count} apps", BAR_BLUE_1,
     "ASA, IOS, ISE, Meraki, Security Cloud, ESA — the core integrations every Cisco+Splunk customer runs"),
    (f"TIER 2 — 100–499 Stacks", f"{tier2_count} apps", BAR_BLUE_2,
     "Duo, eStreamer, Talos, Umbrella, Firepower, WSA, Catalyst, ThousandEyes — broad deployment across security & networking"),
    (f"TIER 3 — 10–99 Stacks", f"{tier3_count} apps", BAR_BLUE_3,
     "AppDynamics, DNA Center, SD-WAN, ACI, Intersight, NVM, CyberVision, Webex — specialized but real customers"),
    (f"TIER 4 — 1–9 Stacks", f"  {tier4_count} apps", RGBColor(0x90, 0xCA, 0xF9),
     "Prime Infra, Spaces, OpenDNS, Threat Grid, Bug Search — long-tail adoption, still running on production stacks"),
]

tier_box_h = 0.95
tier_box_gap = 0.1
for i, (header, count, color, desc) in enumerate(tier_entries):
    y = tier_y_start + 0.35 + i * (tier_box_h + tier_box_gap)

    shape = slide_cloud.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(tier_x), Inches(y), Inches(tier_w), Inches(tier_box_h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = OFF_WHITE
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(0.75)

    accent = slide_cloud.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(tier_x), Inches(y), Inches(0.06), Inches(tier_box_h)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()

    add_text_box(slide_cloud, tier_x + 0.15, y + 0.05, 3.0, 0.25,
                 header, font_size=10, bold=True, color=color)
    add_text_box(slide_cloud, tier_x + 3.5, y + 0.05, 1.2, 0.25,
                 count, font_size=10, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide_cloud, tier_x + 0.15, y + 0.35, tier_w - 0.3, 0.55,
                 desc, font_size=8, color=MED_GRAY)

# ── Key Takeaway ──
takeaway_y = 6.05
add_rounded_rect(slide_cloud, 0.2, takeaway_y, 12.9, 0.3, PURPLE,
                 "KEY TAKEAWAY", font_size=10, font_color=WHITE)

takeaway_runs = [
    [("73% of all Cisco Splunkbase apps ", 11, True, CISCO_BLUE),
     ("are still actively deployed in Splunk Cloud — ", 11, False, DARK_TEXT),
     ("the archiving didn't stop usage, it stopped updates.", 11, True, DEBT_RED)],
    [("~30 archived apps ", 11, True, ARCHIVE_ORANGE),
     ("are running on production Cloud stacks with ", 11, False, DARK_TEXT),
     ("no path to patches or upgrades  •  ", 11, True, DEBT_RED),
     ("2,973 stacks", 11, True, PHASE_GREEN),
     (" depend on Cisco integrations", 11, False, DARK_TEXT)],
]
add_rich_text_box(slide_cloud, 0.5, takeaway_y + 0.32, 12.4, 0.55, takeaway_runs)

# Footer
add_rounded_rect(slide_cloud, 0.3, 7.05, 12.7, 0.32, DARK_BG, "", font_size=9)
footer_cloud = [
    [("Splunk Cloud data only", 9, True, WARN_AMBER),
     (" — Enterprise deployment numbers are not included. ", 9, False, RGBColor(0xB0, 0xBE, 0xC5)),
     ("Actual Cisco footprint across all Splunk deployments is likely significantly higher.", 9, True, WHITE)],
]
add_rich_text_box(slide_cloud, 0.6, 7.07, 12.2, 0.28, footer_cloud, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide_cloud, """HOW TO DELIVER THIS SLIDE:

[Transition — this is the evidence that validates the urgency]
"You just saw that only 18 of 59 remaining apps are enterprise-grade. Some people might think 'well, maybe nobody uses these apps anyway.' Let me show you the data."

[Point to the hero stats]
"89 Cisco apps are actively deployed in Splunk Cloud right now. Across 2,973 Cloud stacks."

"That's 73% of everything that was ever on Splunkbase. Still running."

[Dramatic pause]
"And approximately 30 of those 89 — are from the ARCHIVED set. Apps that were removed from Splunkbase. Still running in production. With no path to updates, no patches, no upgrades."

[Point to the bar chart]
"Look at the top 15. Cisco ASA — deployed on 1,427 Cloud stacks. IOS on 1,295. ISE on 917. Meraki on 833."

"These aren't niche integrations. These are core infrastructure data sources that thousands of SOCs and NOCs depend on."

[Point to the tiers]
"Tier 1 — 6 apps on 500+ stacks each. These are the crown jewels."
"Tier 2 — 18 apps between 100 and 499 stacks. Duo, eStreamer, Talos, Umbrella."
"Tier 3 — 38 apps on 10-99 stacks. AppDynamics, DNA Center, SD-WAN."
"Even Tier 4 — 27 apps on 1-9 stacks — that's still real customers in production."

[Deliver the punchline]
"The archiving removed LISTINGS. It didn't remove USAGE. These apps are still running. The question is: who's maintaining them? Who's patching vulnerabilities? Who's ensuring compatibility when Splunk upgrades?"

"And remember — this is CLOUD ONLY. We don't even have Enterprise numbers. The actual footprint is much larger."

NOTE: If asked about Enterprise numbers —
"We couldn't get Enterprise deployment data for this analysis. But given that Cloud represents roughly 40% of the Splunk install base, you can extrapolate that Enterprise numbers are at least this large, if not larger."
""")


# ══════════════════════════════════════════════════════
# SLIDES 3a–3c: Progressive Build — "The Coverage Gap"
#   Steve Jobs: Reveal one idea at a time.
#   Each "click" is the next slide with prior content
#   already visible. Seamless in presentation mode.
# ══════════════════════════════════════════════════════

# ── Shared layout constants ──
col_y = 1.6
col_h = 3.7
col_w = 3.8
col_gap = 0.35
cx1 = 0.6
cx2 = cx1 + col_w + col_gap
cx3 = cx2 + col_w + col_gap

def add_slide3_title_and_subtitle(slide):
    add_text_box(slide, 0.5, 0.3, 12.3, 0.7,
                 "The Cisco Coverage Gap — Where We Stand",
                 font_size=30, bold=True, color=DARK_TEXT)
    add_text_box(slide, 0.5, 0.95, 12.3, 0.35,
                 "64 Cisco products in the portfolio. Only 18 integrations pass the enterprise bar.",
                 font_size=14, color=SUBTITLE_GRAY)

def add_slide3_footer(slide):
    add_rounded_rect(slide, 0.3, 7.05, 12.7, 0.32, DARK_BG, "", font_size=9)
    footer3 = [
        [("Data source: ", 9, True, WARN_AMBER),
         ("products.conf (SCAN v1.0.4) + Splunkbase listings as of March 2026  •  ", 9, False, RGBColor(0xB0, 0xBE, 0xC5)),
         ("Cloud usage: 89 apps still deployed across 2,973 stacks", 9, True, WHITE)],
    ]
    add_rich_text_box(slide, 0.6, 7.07, 12.2, 0.28, footer3, alignment=PP_ALIGN.CENTER)

def add_col1_active(slide):
    add_section_box(slide, cx1, col_y, col_w, col_h,
        "47 Active Products", CISCO_BLUE,
        [
            "Cisco products actively in-market",
            "",
            "**Covered on Splunkbase:**",
            "   36 products have at least one app/TA",
            "",
            "**Enterprise-grade coverage:**",
            "   Only 18 with Cisco or Splunk support",
            "",
            "**Gap:**",
            "   11 active products have zero coverage",
            "   18 have coverage but no official support",
        ], body_size=11)

def add_col2_roadmap(slide):
    add_section_box(slide, cx2, col_y, col_w, col_h,
        "17 GTM Roadmap Products", ARCHIVE_ORANGE,
        [
            "Products on the go-to-market roadmap",
            "with zero Splunkbase presence today",
            "",
            "**Examples:**",
            "   Hypershield",
            "   Secure AI Factory",
            "   Security Cloud Control",
            "   pxGrid / pxGrid Cloud",
            "   Industrial Networking",
            "",
            "These are NEW products — not archiving gaps",
        ], body_size=11)

def add_col3_scorecard(slide):
    add_section_box(slide, cx3, col_y, col_w, col_h,
        "The Enterprise Scorecard", PHASE_BLUE,
        [
            "",
            "**64** Cisco products in portfolio",
            "   (47 active + 17 GTM roadmap)",
            "",
            "**18** enterprise-grade integrations",
            "",
            "═══════════════════════",
            "",
            "   Enterprise coverage rate:",
            "",
        ], body_size=11)
    add_big_number(slide, cx3 - 0.05, col_y + 3.0, "28%", DEBT_RED, font_size=42)

def add_impact_section(slide):
    impact_y = 5.55
    add_text_box(slide, 0.6, impact_y, 12.0, 0.35,
                 "WHAT THIS MEANS",
                 font_size=14, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    imp_w = 3.8
    imp_h = 0.85
    imp_y = 5.95
    ix1 = 0.6
    ix2 = ix1 + imp_w + col_gap
    ix3 = ix2 + imp_w + col_gap

    add_rounded_rect(slide, ix1, imp_y, imp_w, imp_h, DEBT_RED, "", font_size=11)
    imp1 = [
        [("Customer Risk", 12, True, WHITE)],
        [("Customers deploying unsupported apps get", 10, False, RGBColor(0xFF, 0xCC, 0xCC)),
         (" no security patches, no SLA, no TAC help", 10, True, WHITE)],
    ]
    add_rich_text_box(slide, ix1 + 0.15, imp_y + 0.05, imp_w - 0.3, imp_h - 0.1, imp1, alignment=PP_ALIGN.CENTER)

    add_rounded_rect(slide, ix2, imp_y, imp_w, imp_h, ARCHIVE_ORANGE, "", font_size=11)
    imp2 = [
        [("Competitive Risk", 12, True, WHITE)],
        [("Palo Alto, CrowdStrike, and Fortinet all", 10, False, RGBColor(0xFF, 0xE0, 0xCC)),
         (" have 100% vendor-supported TAs", 10, True, WHITE)],
    ]
    add_rich_text_box(slide, ix2 + 0.15, imp_y + 0.05, imp_w - 0.3, imp_h - 0.1, imp2, alignment=PP_ALIGN.CENTER)

    add_rounded_rect(slide, ix3, imp_y, imp_w, imp_h, PHASE_BLUE, "", font_size=11)
    imp3 = [
        [("The Opportunity", 12, True, WHITE)],
        [("SCAN + expanded Cisco-supported TAs can", 10, False, RGBColor(0xBB, 0xDE, 0xFB)),
         (" close the gap and lead the ecosystem", 10, True, WHITE)],
    ]
    add_rich_text_box(slide, ix3 + 0.15, imp_y + 0.05, imp_w - 0.3, imp_h - 0.1, imp3, alignment=PP_ALIGN.CENTER)


# ── SLIDE 3a: Just the title + Column 1 (47 Active Products) ──
slide3a = add_slide(prs, use_template)
add_slide3_title_and_subtitle(slide3a)
add_col1_active(slide3a)
add_slide3_footer(slide3a)

add_speaker_notes(slide3a, """HOW TO DELIVER THIS SLIDE:

[Transition from Slide 2]
"So... 18 enterprise-grade apps. But let's zoom out. Let's look at the FULL Cisco portfolio."

[Column 1 appears alone — centered attention]
"47 Cisco products are actively in-market today."

"36 of them have at least some presence on Splunkbase — an app or a TA."

"But when we filter for enterprise-grade — Cisco or Splunk supported — we're back to 18."

"And 11 active products? Zero coverage. No app, no add-on, nothing on Splunkbase."

[Pause — let them absorb]
"But that's just the products we're selling today..."

[Click to next slide to reveal Column 2]
""")


# ── SLIDE 3b: Column 1 + Column 2 (17 GTM Roadmap) ──
slide3b = add_slide(prs, use_template)
add_slide3_title_and_subtitle(slide3b)
add_col1_active(slide3b)
add_col2_roadmap(slide3b)
add_slide3_footer(slide3b)

add_speaker_notes(slide3b, """HOW TO DELIVER THIS SLIDE:

[Column 2 appears — the GTM roadmap]
"Then there are 17 products on the go-to-market roadmap."

"Hypershield. Secure AI Factory. Security Cloud Control. pxGrid Cloud. Industrial Networking."

"These are products Cisco is actively launching — and they have ZERO Splunkbase presence."

[Important distinction]
"These aren't archiving casualties. These are brand-new products shipping without any Splunk integration strategy."

"The field is being asked to sell these products to Splunk customers with no data onboarding story."

[Click to next slide for the scorecard reveal]
""")


# ── SLIDE 3c: All three columns + impact cards (the full reveal) ──
slide3c = add_slide(prs, use_template)
add_slide3_title_and_subtitle(slide3c)
add_col1_active(slide3c)
add_col2_roadmap(slide3c)
add_col3_scorecard(slide3c)
add_impact_section(slide3c)
add_slide3_footer(slide3c)

add_speaker_notes(slide3c, """HOW TO DELIVER THIS SLIDE:

[The Enterprise Scorecard appears — the big reveal]
"Add it all up."

"64 Cisco products in the portfolio. 47 active. 17 on the GTM roadmap."

[Pause]
"18 enterprise-grade integrations."

[Let the number land]
"28 percent. That is our enterprise coverage rate."

"For the number one infrastructure vendor in the world."

[Point to the three impact cards]
"This creates three risks."

"CUSTOMER RISK — customers deploying unsupported apps get no security patches, no SLA, no TAC help. They're on their own."

"COMPETITIVE RISK — and this one stings. Palo Alto, CrowdStrike, Fortinet? They have 100% vendor-supported TAs on Splunkbase. Every single one. We're at 28%."

"But here's the thing — THE OPPORTUNITY. SCAN gives us the platform to see exactly where the holes are. Every gap is visible. Every priority is clear. We can close this."

[Close this section]
"This isn't about what Splunk archived. It's about what Cisco needs to build."

COMPETITIVE CONTEXT (if asked):
• Palo Alto Networks: ~15 apps, ALL officially Palo Alto or Splunk supported
• CrowdStrike: ~8 apps, ALL vendor-supported
• Fortinet: ~10 apps, ALL vendor-supported
• Cisco: 59 apps, only 18 vendor-supported (28% of portfolio)
""")


# ══════════════════════════════════════════════════════
# SLIDE 4: "What We're Hearing From the Field"
#   Customer-facing pain points from the Blueprint.
#   Source: Section 7.0, credit Bryan Pluta & James Young.
#   Steve Jobs: Let the customer's pain speak.
# ══════════════════════════════════════════════════════
slide4 = add_slide(prs, use_template)

add_text_box(slide4, 0.5, 0.3, 12.3, 0.7,
             "What We're Hearing From the Field",
             font_size=30, bold=True, color=DARK_TEXT)

add_text_box(slide4, 0.5, 0.95, 12.3, 0.35,
             "Real feedback from SEs, architects, and customers — the daily pain of Cisco data onboarding",
             font_size=14, color=SUBTITLE_GRAY)

# ── Seven pain point cards: 3 + 3 + 1 layout across the slide ──
# Row 1: Three cards
card_w = 3.8
card_h = 2.0
card_gap = 0.35
row1_y = 1.55
row2_y = row1_y + card_h + 0.25
rx1 = 0.6
rx2 = rx1 + card_w + card_gap
rx3 = rx2 + card_w + card_gap

# Card 1: Misleading Coverage
add_section_box(slide4, rx1, row1_y, card_w, card_h,
    "Misleading Coverage", DEBT_RED,
    [
        "Add-ons often cover more products than",
        "the name implies. A TA labeled for one",
        "product may quietly handle three others.",
        "",
        "Customers install the wrong TA, or install",
        "duplicates they didn't need.",
    ], body_size=10)

# Card 2: Naming Chaos
add_section_box(slide4, rx2, row1_y, card_w, card_h,
    "Naming Chaos", ARCHIVE_ORANGE,
    [
        "Apps and TAs are not clearly distinguished.",
        "",
        "\"Which one do I install?\"",
        "— Every Splunk admin, every time.",
        "",
        "6 different Firepower-related add-ons.",
        "4 different Meraki entries.",
    ], body_size=10)

# Card 3: Broken Parsing
add_section_box(slide4, rx3, row1_y, card_w, card_h,
    "Broken Parsing", WARN_AMBER,
    [
        "Incomplete parsing rules result in blank",
        "fields at search time.",
        "",
        "Example: cisco:sna sourcetype has EVAL",
        "errors leaving critical fields empty.",
        "",
        "Raw logs land. Nothing gets extracted.",
    ], body_size=10)

# Row 2: Three cards
# Card 4: Missing "Super Six"
add_section_box(slide4, rx1, row2_y, card_w, card_h,
    "Missing \"Super Six\" Configs", DEBT_RED,
    [
        "TAs ship without essential props.conf",
        "settings (TIME_FORMAT, SHOULD_LINEMERGE,",
        "LINE_BREAKER, etc.)",
        "",
        "Result: high CPU usage, line-merging",
        "failures, timestamp misidentification.",
        "Admins must hand-fix every deployment.",
    ], body_size=10)

# Card 5: CIM Non-Compliance
add_section_box(slide4, rx2, row2_y, card_w, card_h,
    "CIM Non-Compliance", ARCHIVE_ORANGE,
    [
        "Non-CIM-compliant TAs break Splunk ES",
        "out-of-the-box detections.",
        "",
        "**1,900+ OOTB detections** rely on CIM.",
        "",
        "If the TA doesn't map to CIM, Cisco",
        "data is invisible to the SOC.",
    ], body_size=10)

# Card 6: Redundancy
add_section_box(slide4, rx3, row2_y, card_w, card_h,
    "Redundant Add-ons", WARN_AMBER,
    [
        "Multiple TAs exist for the same product",
        "from Cisco, Splunk, and community devs.",
        "",
        "6 Firepower/eStreamer TAs",
        "4 Meraki TAs",
        "3 ISE TAs",
        "Customers don't know which to trust.",
    ], body_size=10)

# ── Bottom section: Feature Parity (full width banner) ──
fp_y = row2_y + card_h + 0.3
add_rounded_rect(slide4, 0.6, fp_y, 12.1, 0.65, RGBColor(0xF3, 0xE5, 0xF5),
                 "", font_size=11, bold=False)
fp_accent = slide4.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.6), Inches(fp_y), Inches(0.06), Inches(0.65)
)
fp_accent.fill.solid()
fp_accent.fill.fore_color.rgb = PURPLE
fp_accent.line.fill.background()

fp_runs = [
    [("Lack of Feature Parity: ", 12, True, PURPLE),
     ("New official apps often lack the depth of older community versions. ", 11, False, DARK_TEXT),
     ("Customers are asked to downgrade functionality for the sake of 'official support.'", 11, True, DEBT_RED)],
]
add_rich_text_box(slide4, 0.85, fp_y + 0.08, 11.6, 0.55, fp_runs)

# ── Footer ──
add_rounded_rect(slide4, 0.3, 7.05, 12.7, 0.32, DARK_BG, "", font_size=9)
footer4 = [
    [("Source: ", 9, True, WARN_AMBER),
     ("Strategic Blueprint §7.0 — Key Technical & Operational Challenges  •  ", 9, False, RGBColor(0xB0, 0xBE, 0xC5)),
     ("Credit: Bryan Pluta, James Young, & field engineering coalition", 9, True, WHITE)],
]
add_rich_text_box(slide4, 0.6, 7.07, 12.2, 0.28, footer4, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide4, """HOW TO DELIVER THIS SLIDE:

[Transition from Slide 3]
"We've talked numbers. Let me show you what those numbers actually feel like on the ground."

"These aren't hypothetical. These are real issues Bryan Pluta and James Young documented from field engagements. Every card here is a conversation that's happened with a customer."

[Point to cards left-to-right, top row]

"MISLEADING COVERAGE — a TA says 'Cisco Firepower' but quietly handles eStreamer, FTD, and FMC. Customers install the wrong one. Or install three when they needed one."

"NAMING CHAOS — is it an 'app' or a 'TA'? Is 'Splunk_TA_cisco-asa' different from 'TA-cisco_acs'? There are six Firepower-related entries. Four for Meraki. Customers are guessing."

"BROKEN PARSING — they get the TA installed, data starts flowing, they go to search... and fields are blank. The cisco:sna sourcetype has EVAL errors that silently drop fields. The data is there. The parsing isn't."

[Point to bottom row]

"MISSING SUPER SIX — these are the six essential props.conf settings that every TA needs. TIME_FORMAT, SHOULD_LINEMERGE, LINE_BREAKER — the basics. Without them, you get CPU spikes, mangled events, wrong timestamps. Admins have to hand-fix every single deployment."

"CIM NON-COMPLIANCE — this one should scare you. 1,900+ out-of-the-box detections in Splunk ES rely on CIM-mapped data. If the Cisco TA doesn't map to CIM, the data is invisible to the SOC. The customer bought Enterprise Security, onboarded Cisco data, and gets zero detections."

"REDUNDANCY — same product, three different TAs from three different sources. Which one does the customer deploy? The Cisco one? The Splunk one? The community one? Nobody knows. And nobody is cleaning it up."

[Point to the feature parity banner]
"And the kicker — when we DO ship a new official app, it often has FEWER features than the community version it replaced. We're literally asking customers to downgrade for 'official support.'"

[Pause]
"This is what the field deals with every day. Every SE. Every customer engagement. This is why we need SCAN."
""")


# ══════════════════════════════════════════════════════
# SLIDE 5: "The Systemic Issues"
#   Organizational / strategic challenges (items 8-11).
#   These are the HARD conversations.
#   Steve Jobs: Name the elephant in the room.
# ══════════════════════════════════════════════════════
slide5 = add_slide(prs, use_template)

add_text_box(slide5, 0.5, 0.3, 12.3, 0.7,
             "The Deeper Problem — Why This Keeps Happening",
             font_size=30, bold=True, color=DARK_TEXT)

add_text_box(slide5, 0.5, 0.95, 12.3, 0.35,
             "Technical debt is a symptom. These are the systemic causes.",
             font_size=14, color=SUBTITLE_GRAY)

# ── Four large issue cards: 2x2 grid ──
issue_w = 5.8
issue_h = 2.25
issue_gap_x = 0.5
issue_gap_y = 0.3
iy1 = 1.55
iy2 = iy1 + issue_h + issue_gap_y
ix1 = 0.6
ix2 = ix1 + issue_w + issue_gap_x

# Issue 1: Support Loop (item 8)
add_section_box(slide5, ix1, iy1, issue_w, issue_h,
    "The Support Loop", DEBT_RED,
    [
        "When a Cisco TA breaks, who does the customer call?",
        "",
        "Splunk Support says: \"It's a Cisco TA — call Cisco TAC.\"",
        "Cisco TAC says: \"It's a Splunk platform issue — call Splunk.\"",
        "",
        "Customers bounce between two support orgs",
        "with no shared diagnostic data and no SLA.",
        "",
        "**Forcing Splunk admins to call Cisco TAC is an anti-pattern.**",
    ], body_size=11)

# Issue 2: Vendor Fragmentation (item 9)
add_section_box(slide5, ix2, iy1, issue_w, issue_h,
    "Vendor Fragmentation & Velocity Loss", ARCHIVE_ORANGE,
    [
        "We are in a \"dark period\" — development transitioning from",
        "Crest to HCL with no continuity.",
        "",
        "No single vendor or unified strategy across Cisco BUs.",
        "Constant vendor onboarding drains engineering velocity.",
        "",
        "Institutional knowledge is lost with every transition.",
        "A \"Gold Standard\" is impossible to maintain.",
        "                                                  — Andrei Sheleh",
    ], body_size=11)

# Issue 3: OCSF Distraction (item 10)
add_section_box(slide5, ix1, iy2, issue_w, issue_h,
    "OCSF as a Strategic Distraction", WARN_AMBER,
    [
        "Heavy focus on OCSF has diverted resources away from",
        "solving immediate data onboarding challenges.",
        "",
        "No clear strategy for OCSF + CIM coexistence.",
        "Value remains unproven for most customers not using",
        "AWS or the latest Cisco stacks.",
        "",
        "Meanwhile, day-to-day parsing, CIM mapping, and",
        "props.conf fixes remain unfunded.",
    ], body_size=11)

# Issue 4: Org Silos (item 11)
add_section_box(slide5, ix2, iy2, issue_w, issue_h,
    "Organizational Silos → Customer Pain", PURPLE,
    [
        "Our TA ecosystem reflects Cisco's internal BU structure",
        "rather than customer-centric outcomes.",
        "",
        "Customers must navigate our org chart just to",
        "onboard a single data source.",
        "",
        "Each BU ships its own TA, its own naming convention,",
        "its own support model. No coordination.",
        "                                       — Sarav Radhakrishnan",
    ], body_size=11)

# ── Bottom banner: The Ask ──
ask_y = iy2 + issue_h + 0.25
add_rounded_rect(slide5, 0.6, ask_y, 12.1, 0.65, PHASE_BLUE, "", font_size=12)
ask_runs = [
    [("These aren't feature requests — they are structural failures.", 13, True, WHITE),
     ("  The field needs a ", 13, False, RGBColor(0xBB, 0xDE, 0xFB)),
     ("unified vendor strategy", 13, True, WARN_AMBER),
     (", a ", 13, False, RGBColor(0xBB, 0xDE, 0xFB)),
     ("shared support protocol", 13, True, WARN_AMBER),
     (", and ", 13, False, RGBColor(0xBB, 0xDE, 0xFB)),
     ("CIM-first engineering.", 13, True, WARN_AMBER)],
]
add_rich_text_box(slide5, 0.9, ask_y + 0.08, 11.5, 0.55, ask_runs, alignment=PP_ALIGN.CENTER)

# Footer
add_rounded_rect(slide5, 0.3, 7.05, 12.7, 0.32, DARK_BG, "", font_size=9)
footer5 = [
    [("Source: ", 9, True, WARN_AMBER),
     ("Strategic Blueprint §7.0  •  ", 9, False, RGBColor(0xB0, 0xBE, 0xC5)),
     ("Insights from Bryan Pluta, James Young, Andrei Sheleh, Sarav Radhakrishnan & 20+ SMEs", 9, True, WHITE)],
]
add_rich_text_box(slide5, 0.6, 7.07, 12.2, 0.28, footer5, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide5, """HOW TO DELIVER THIS SLIDE:

[Transition from Slide 4]
"Those were the symptoms. Now let me show you why this keeps happening."

"These are the four systemic issues. If we don't address these, we'll be having this same conversation next year."

[Point to top-left - The Support Loop]
"THE SUPPORT LOOP. This one is personal for every SE in the room."

"Customer deploys a Cisco TA. Something breaks. They call Splunk Support."
"Splunk says: 'That's a Cisco-maintained TA — you need to call Cisco TAC.'"
"Customer calls Cisco TAC. TAC says: 'That's a Splunk platform issue.'"

"Ping pong. Back and forth. Nobody owns the problem. The customer sits in the middle with a broken deployment."

"Bryan Pluta said it best: 'Forcing Splunk admins to call Cisco TAC is an anti-pattern.'"

[Point to top-right - Vendor Fragmentation]
"VENDOR FRAGMENTATION. This is Andrei Sheleh's observation and it's spot on."

"We keep cycling through development vendors. Crest. Now HCL. Each time, we lose institutional knowledge. Each time, we restart from zero. There is no single vendor strategy across Cisco BUs."

"You can't build a Gold Standard when the team building it changes every 18 months."

[Point to bottom-left - OCSF]
"OCSF. Now this one is nuanced."

"OCSF may be the right long-term bet. But right now, it's diverting resources from problems customers have TODAY. We have TAs shipping without basic props.conf configs. We have CIM mapping gaps breaking ES detections. And we're spending engineering cycles on a schema that most customers can't even use yet."

"We need a clear strategy for how OCSF and CIM coexist. Until we have that, OCSF is a distraction."

[Point to bottom-right - Org Silos]
"And finally — ORGANIZATIONAL SILOS. Sarav Radhakrishnan nailed this one."

"Our TA ecosystem is a mirror of Cisco's org chart. The Security BU ships one TA. The Networking BU ships another. Collaboration ships their own. No shared naming. No shared quality bar. No shared support model."

"The customer doesn't care about our org chart. They just want to onboard Cisco data. And right now, they have to navigate our internal politics to do it."

[Point to the bottom banner]
"These aren't feature requests. These are structural failures. And fixing them requires three things:"
"1. A unified vendor strategy — one team, one standard"
"2. A shared support protocol — no more ping pong"
"3. CIM-first engineering — if it doesn't map to CIM, it doesn't ship"

[Pause]
"SCAN can't fix all of this. But SCAN gives us the visibility to know exactly where we stand. And that's the first step."
""")


# ─── Save ───
prs.save(OUTPUT_PATH)
size_kb = os.path.getsize(OUTPUT_PATH) // 1024
print(f"\nSaved: {OUTPUT_PATH}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {size_kb:,} KB")
